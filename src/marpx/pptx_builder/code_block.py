"""Code block rendering for pptx_builder."""

from __future__ import annotations

from marpx.models import CodeBlockElement

from ._helpers import _set_fill_color
from .decoration import _add_decoration_shape


def _add_code_block(slide, element: CodeBlockElement) -> None:
    """Add a code block as textbox with background."""
    from .text import (
        ALIGNMENT_MAP,
        _apply_paragraph_layout,
        _apply_text_style,
        _resolve_textbox_geometry,
        _set_text_frame_margins_zero,
    )
    from pptx.enum.text import PP_ALIGN

    if element.decoration:
        _add_decoration_shape(slide, element.box, element.decoration)

    left, top, width, height = _resolve_textbox_geometry(element)

    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    _set_text_frame_margins_zero(tf)
    txbox.line.fill.background()

    if element.decoration:
        txbox.fill.background()
    elif element.code_background:
        fill = txbox.fill
        _set_fill_color(fill, element.code_background)
    else:
        txbox.fill.background()

    # Add code text with monospace font
    for i, para in enumerate(element.paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = ALIGNMENT_MAP.get(para.alignment, PP_ALIGN.LEFT)
        _apply_paragraph_layout(p, para)

        for text_run in para.runs:
            r = p.add_run()
            r.text = text_run.text
            # Force monospace
            style = text_run.style.model_copy()
            style.font_family = "Courier New"
            _apply_text_style(r, style)
