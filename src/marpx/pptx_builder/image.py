"""Image rendering functions for pptx_builder."""

from __future__ import annotations

import io
import logging
import urllib.parse
from pathlib import Path

from pptx.util import Emu

from marpx.utils.image import resolve_image_bytes
from marpx.models import ImageElement
from marpx.utils.svg import (
    MissingDependencyError,
    rasterize_svg_to_png,
)
from marpx.utils.common import px_to_emu

from .decoration import _add_decoration_shape
from .decoration import _apply_round_rect_radius_to_geom
from ._helpers import _set_blip_alpha

logger = logging.getLogger(__name__)


def _is_svg_source(src: str) -> bool:
    """Return True when the source points to SVG content."""
    if src.startswith("data:image/svg+xml"):
        return True

    parsed = urllib.parse.urlparse(src)
    candidate = parsed.path if parsed.scheme else src
    return Path(urllib.parse.unquote(candidate)).suffix.lower() == ".svg"


def _resolve_image_placement(element: ImageElement) -> tuple[Emu, Emu, Emu, Emu]:
    """Convert CSS image box/object-fit metadata to PPTX placement."""
    left_px = element.box.x
    top_px = element.box.y
    width_px = element.box.width
    height_px = element.box.height
    if element.decoration:
        left_px += (
            element.decoration.border_left.width_px + element.decoration.padding.left_px
        )
        top_px += (
            element.decoration.border_top.width_px + element.decoration.padding.top_px
        )
        width_px -= (
            element.decoration.border_left.width_px
            + element.decoration.border_right.width_px
            + element.decoration.padding.left_px
            + element.decoration.padding.right_px
        )
        height_px -= (
            element.decoration.border_top.width_px
            + element.decoration.border_bottom.width_px
            + element.decoration.padding.top_px
            + element.decoration.padding.bottom_px
        )
        width_px = max(width_px, 1)
        height_px = max(height_px, 1)

    fit = (element.object_fit or "").lower()
    natural_w = element.image_natural_width_px or 0
    natural_h = element.image_natural_height_px or 0
    if fit not in {"contain", "scale-down"} or natural_w <= 0 or natural_h <= 0:
        return (
            Emu(px_to_emu(left_px)),
            Emu(px_to_emu(top_px)),
            Emu(px_to_emu(width_px)),
            Emu(px_to_emu(height_px)),
        )

    contain_scale = min(width_px / natural_w, height_px / natural_h)
    if fit == "scale-down":
        contain_scale = min(contain_scale, 1.0)
    draw_w = natural_w * contain_scale
    draw_h = natural_h * contain_scale

    pos = (element.object_position or "50% 50%").lower()
    x_ratio = 0.5
    y_ratio = 0.5
    if "left" in pos:
        x_ratio = 0.0
    elif "right" in pos:
        x_ratio = 1.0
    if "top" in pos:
        y_ratio = 0.0
    elif "bottom" in pos:
        y_ratio = 1.0

    left_px += max(width_px - draw_w, 0) * x_ratio
    top_px += max(height_px - draw_h, 0) * y_ratio

    return (
        Emu(px_to_emu(left_px)),
        Emu(px_to_emu(top_px)),
        Emu(px_to_emu(draw_w)),
        Emu(px_to_emu(draw_h)),
    )


def _should_round_picture_geometry(element: ImageElement) -> bool:
    """Round the picture itself only when the image is expected to fill the box."""
    if not element.decoration or element.decoration.border_radius_px <= 0:
        return False
    return (element.object_fit or "").lower() not in {"contain", "scale-down"}


def _add_image(slide, element: ImageElement) -> None:
    """Add an image to the slide."""
    if not element.image_src:
        return

    if element.decoration:
        _add_decoration_shape(slide, element.box, element.decoration)

    left, top, width, height = _resolve_image_placement(element)

    try:
        picture = None
        # SVG requires rasterization first
        if _is_svg_source(element.image_src):
            image_bytes = rasterize_svg_to_png(element.image_src or "")
            picture = slide.shapes.add_picture(
                io.BytesIO(image_bytes), left, top, width, height
            )
        # Local file:// URLs: pass path directly to avoid loading bytes into memory
        elif element.image_src.startswith("file://"):
            image_path = Path(
                urllib.parse.unquote(urllib.parse.urlparse(element.image_src).path)
            )
            if image_path.exists():
                picture = slide.shapes.add_picture(
                    str(image_path), left, top, width, height
                )
            else:
                logger.warning("Image file not found: %s", element.image_src)
        # Bare local paths (no scheme, not a data URI): pass path directly
        elif (
            not element.image_src.startswith("data:")
            and not element.image_src.startswith(("http://", "https://"))
            and not element.image_data
        ):
            image_path = Path(element.image_src)
            if image_path.exists():
                picture = slide.shapes.add_picture(
                    str(image_path), left, top, width, height
                )
            else:
                logger.warning("Image file not found: %s", element.image_src)
        else:
            # Data URIs, HTTP/HTTPS URLs, pre-loaded bytes -> delegate to resolve_image_bytes
            image_bytes = resolve_image_bytes(element.image_src, element.image_data)
            if image_bytes:
                picture = slide.shapes.add_picture(
                    io.BytesIO(image_bytes), left, top, width, height
                )
            else:
                logger.warning("Failed to resolve image: %s", element.image_src)

        if picture is not None and element.image_opacity < 1.0:
            blip = picture._element.blipFill.blip
            _set_blip_alpha(blip, element.image_opacity)

        if picture is not None and _should_round_picture_geometry(element):
            picture._element.spPr.prstGeom.set("prst", "roundRect")
            _apply_round_rect_radius_to_geom(
                picture._element.spPr.prstGeom,
                int(width),
                int(height),
                px_to_emu(element.decoration.border_radius_px),
            )
    except MissingDependencyError:
        raise
    except Exception as e:
        logger.warning("Failed to add image %s: %s", element.image_src, e)
