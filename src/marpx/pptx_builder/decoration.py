"""Decoration shape rendering for pptx_builder."""

from __future__ import annotations

import logging
from pathlib import Path

from lxml import etree
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Emu

from marpx.utils.gradient import parse_linear_gradient
from marpx.models import (
    BaseSlideElement,
    Box,
    BoxDecoration,
    BoxShadow,
    ClipPath,
    CodeBlockElement,
    RGBAColor,
    SlideElement,
)
from marpx.utils.common import px_to_emu
from .scene3d import fit_scene3d_rotations

from ._helpers import (
    _build_gradient_fill_xml,
    _remove_existing_fills,
    _set_fill_color,
    _set_inner_shadow,
    _set_line_color,
    _set_outer_shadow,
    _with_opacity,
)

logger = logging.getLogger(__name__)


def _normalize_positive_fixed_angle(angle_deg: float) -> int:
    """Convert degrees to a positive OOXML fixed-angle value."""
    normalized = angle_deg % 360.0
    normalized = (normalized + 360.0) % 360.0  # Ensure positive
    return int(round(normalized * 60000))


def _apply_scene3d(
    sp_pr,
    rotation_3d_x_deg: float = 0.0,
    rotation_3d_y_deg: float = 0.0,
    rotation_3d_z_deg: float = 0.0,
) -> None:
    """Apply PowerPoint 3D scene rotation to a shape properties node."""
    if (
        abs(rotation_3d_x_deg) <= 0.01
        and abs(rotation_3d_y_deg) <= 0.01
        and abs(rotation_3d_z_deg) <= 0.01
    ):
        return

    scene3d = sp_pr.find(qn("a:scene3d"))
    if scene3d is not None:
        sp_pr.remove(scene3d)

    scene3d = etree.SubElement(sp_pr, qn("a:scene3d"))
    camera = etree.SubElement(scene3d, qn("a:camera"))
    camera.set("prst", "orthographicFront")
    rot = etree.SubElement(camera, qn("a:rot"))
    rot.set("lat", str(_normalize_positive_fixed_angle(rotation_3d_y_deg)))
    rot.set("lon", str(_normalize_positive_fixed_angle(rotation_3d_x_deg)))
    rot.set("rev", str(_normalize_positive_fixed_angle(rotation_3d_z_deg)))
    light_rig = etree.SubElement(scene3d, qn("a:lightRig"))
    light_rig.set("rig", "threePt")
    light_rig.set("dir", "t")


def _resolve_scene3d_rotations(element: BaseSlideElement) -> tuple[float, float, float]:
    """Return fitted scene3d angles when projected corners are available."""
    has_explicit_3d_rotation = (
        abs(element.rotation_3d_x_deg) > 0.01
        or abs(element.rotation_3d_y_deg) > 0.01
        or abs(element.rotation_3d_z_deg) > 0.01
    )
    if has_explicit_3d_rotation and element.projected_corners:
        return fit_scene3d_rotations(
            element.projected_corners,
            element.box,
            fallback_x_deg=element.rotation_3d_x_deg,
            fallback_y_deg=element.rotation_3d_y_deg,
            fallback_z_deg=element.rotation_3d_z_deg,
        )
    return (
        element.rotation_3d_x_deg,
        element.rotation_3d_y_deg,
        element.rotation_3d_z_deg,
    )


def _round_rect_adjustment(
    width: int | float, height: int | float, radius: int | float
) -> float:
    """Return normalized roundRect adjustment matching CSS border-radius."""
    min_dim = min(float(width), float(height))
    if min_dim <= 0:
        return 0.0
    return max(0.0, min(float(radius) / min_dim, 0.5))


def _apply_round_rect_radius(
    shape, width: int | float, height: int | float, radius: int | float
) -> None:
    """Apply a CSS-like corner radius to a rounded rectangle auto shape."""
    if radius <= 0 or width <= 0 or height <= 0:
        return
    if getattr(shape, "adjustments", None) and len(shape.adjustments) > 0:
        shape.adjustments[0] = _round_rect_adjustment(width, height, radius)


def _replace_with_round_rect_custgeom(
    shape, w_emu: int, h_emu: int, radius_emu: int
) -> None:
    """Replace a shape's preset geometry with a custGeom rounded rectangle.

    Unlike ``prstGeom="roundRect"``, this custom geometry sets the text
    rectangle to the full shape bounds so that PowerPoint does not apply
    an automatic text-area inset based on the corner radius.
    """
    sp_pr = shape._element.spPr

    prst_geom = sp_pr.find(qn("a:prstGeom"))
    if prst_geom is not None:
        sp_pr.remove(prst_geom)

    r = min(radius_emu, w_emu // 2, h_emu // 2)

    cust_geom = etree.SubElement(sp_pr, qn("a:custGeom"))
    etree.SubElement(cust_geom, qn("a:avLst"))
    etree.SubElement(cust_geom, qn("a:gdLst"))
    etree.SubElement(cust_geom, qn("a:ahLst"))
    etree.SubElement(cust_geom, qn("a:cxnLst"))

    # Text rectangle = full shape bounds (no geometry inset)
    rect = etree.SubElement(cust_geom, qn("a:rect"))
    rect.set("l", "0")
    rect.set("t", "0")
    rect.set("r", str(w_emu))
    rect.set("b", str(h_emu))

    path_lst = etree.SubElement(cust_geom, qn("a:pathLst"))
    path = etree.SubElement(
        path_lst,
        qn("a:path"),
        attrib={"w": str(w_emu), "h": str(h_emu)},
    )

    # Draw rounded rectangle: moveTo top-left+r, then lines+arcs clockwise
    _move = etree.SubElement(path, qn("a:moveTo"))
    etree.SubElement(_move, qn("a:pt"), attrib={"x": str(r), "y": "0"})

    # Top edge → top-right corner
    _ln = etree.SubElement(path, qn("a:lnTo"))
    etree.SubElement(_ln, qn("a:pt"), attrib={"x": str(w_emu - r), "y": "0"})
    _arc = etree.SubElement(path, qn("a:arcTo"))
    _arc.set("wR", str(r))
    _arc.set("hR", str(r))
    _arc.set("stAng", "16200000")  # 270°
    _arc.set("swAng", "5400000")  # +90°

    # Right edge → bottom-right corner
    _ln = etree.SubElement(path, qn("a:lnTo"))
    etree.SubElement(_ln, qn("a:pt"), attrib={"x": str(w_emu), "y": str(h_emu - r)})
    _arc = etree.SubElement(path, qn("a:arcTo"))
    _arc.set("wR", str(r))
    _arc.set("hR", str(r))
    _arc.set("stAng", "0")
    _arc.set("swAng", "5400000")

    # Bottom edge → bottom-left corner
    _ln = etree.SubElement(path, qn("a:lnTo"))
    etree.SubElement(_ln, qn("a:pt"), attrib={"x": str(r), "y": str(h_emu)})
    _arc = etree.SubElement(path, qn("a:arcTo"))
    _arc.set("wR", str(r))
    _arc.set("hR", str(r))
    _arc.set("stAng", "5400000")  # 90°
    _arc.set("swAng", "5400000")

    # Left edge → top-left corner
    _ln = etree.SubElement(path, qn("a:lnTo"))
    etree.SubElement(_ln, qn("a:pt"), attrib={"x": "0", "y": str(r)})
    _arc = etree.SubElement(path, qn("a:arcTo"))
    _arc.set("wR", str(r))
    _arc.set("hR", str(r))
    _arc.set("stAng", "10800000")  # 180°
    _arc.set("swAng", "5400000")

    etree.SubElement(path, qn("a:close"))


def _apply_round_rect_radius_to_geom(
    prst_geom, width: int | float, height: int | float, radius: int | float
) -> None:
    """Apply a CSS-like corner radius to a preset geometry node."""
    if radius <= 0 or width <= 0 or height <= 0:
        return
    adj = int(round(_round_rect_adjustment(width, height, radius) * 100000))
    prst_geom.rewrite_guides([("adj", adj)])


def _create_outer_shadow_shapes(
    slide,
    box: Box,
    decoration: BoxDecoration,
    shape_type,
) -> None:
    """Render outer (non-inset) shadow shapes behind the main background shape."""
    for shadow in decoration.box_shadows:
        if shadow.inset or shadow.color.a <= 0:
            continue
        if _is_spread_outline_shadow(shadow):
            _add_spread_outline_shape(
                slide,
                box,
                decoration,
                shape_type,
                shadow,
            )
        else:
            _add_shadow_shape(
                slide,
                box,
                decoration,
                shape_type,
                shadow,
            )


def _create_polygon_shape(
    slide,
    box: Box,
    clip_path: ClipPath,
):
    """Create a freeform polygon shape from clip-path points.

    Points are percentage values (0-100) relative to the element bounding box.
    """
    left_emu = px_to_emu(box.x)
    top_emu = px_to_emu(box.y)
    w_emu = px_to_emu(box.width)
    h_emu = px_to_emu(box.height)

    points = clip_path.points
    if len(points) < 3:
        logger.warning(
            "Polygon clip-path has fewer than 3 points, falling back to rectangle"
        )
        return None

    # Create a rectangle shape first, then replace its geometry with custom polygon
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Emu(left_emu),
        Emu(top_emu),
        Emu(w_emu),
        Emu(h_emu),
    )

    sp_pr = shape._element.spPr

    # Remove preset geometry
    prst_geom = sp_pr.find(qn("a:prstGeom"))
    if prst_geom is not None:
        sp_pr.remove(prst_geom)

    # Build custom geometry with polygon path
    cust_geom = etree.SubElement(sp_pr, qn("a:custGeom"))
    # Add required empty child elements per OOXML spec
    etree.SubElement(cust_geom, qn("a:avLst"))
    etree.SubElement(cust_geom, qn("a:gdLst"))
    etree.SubElement(cust_geom, qn("a:ahLst"))
    etree.SubElement(cust_geom, qn("a:cxnLst"))
    rect = etree.SubElement(cust_geom, qn("a:rect"))
    rect.set("l", "0")
    rect.set("t", "0")
    rect.set("r", str(w_emu))
    rect.set("b", str(h_emu))

    path_lst = etree.SubElement(cust_geom, qn("a:pathLst"))
    path = etree.SubElement(
        path_lst,
        qn("a:path"),
        attrib={"w": str(w_emu), "h": str(h_emu)},
    )

    # Move to first point
    first = points[0]
    x0 = int(round(first.x / 100.0 * w_emu))
    y0 = int(round(first.y / 100.0 * h_emu))
    move_to = etree.SubElement(path, qn("a:moveTo"))
    etree.SubElement(move_to, qn("a:pt"), attrib={"x": str(x0), "y": str(y0)})

    # Line segments to remaining points
    for pt in points[1:]:
        px = int(round(pt.x / 100.0 * w_emu))
        py = int(round(pt.y / 100.0 * h_emu))
        ln_to = etree.SubElement(path, qn("a:lnTo"))
        etree.SubElement(ln_to, qn("a:pt"), attrib={"x": str(px), "y": str(py)})

    # Close the path
    etree.SubElement(path, qn("a:close"))

    return shape


def _create_background_shape(
    slide,
    box: Box,
    decoration: BoxDecoration,
    shape_type,
):
    """Create the main background shape and apply its fill (solid color or gradient)."""
    # Use polygon shape if clip_path is present
    if decoration.clip_path and decoration.clip_path.type == "polygon":
        bg_shape = _create_polygon_shape(slide, box, decoration.clip_path)
        if bg_shape is not None:
            _remove_theme_style(bg_shape)
            # Apply fill to polygon shape
            if decoration.background_gradient and _set_shape_gradient_fill(
                bg_shape, decoration.background_gradient
            ):
                pass  # gradient applied
            elif decoration.background_color and decoration.opacity > 0:
                fill_color = _with_opacity(
                    decoration.background_color, decoration.opacity
                )
                _set_fill_color(bg_shape.fill, fill_color)
            else:
                bg_shape.fill.background()
            return bg_shape

    left = Emu(px_to_emu(box.x))
    top = Emu(px_to_emu(box.y))
    width = Emu(px_to_emu(box.width))
    height = Emu(px_to_emu(box.height))

    bg_shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    _remove_theme_style(bg_shape)
    if decoration.border_radius_px > 0:
        w_emu = int(px_to_emu(box.width))
        h_emu = int(px_to_emu(box.height))
        r_emu = int(px_to_emu(decoration.border_radius_px))
        _replace_with_round_rect_custgeom(bg_shape, w_emu, h_emu, r_emu)

    if decoration.background_gradient and not _set_shape_gradient_fill(
        bg_shape, decoration.background_gradient
    ):
        fill = bg_shape.fill
        if decoration.background_color and decoration.opacity > 0:
            fill_color = _with_opacity(decoration.background_color, decoration.opacity)
            _set_fill_color(fill, fill_color)
        else:
            fill.background()
    else:
        fill = bg_shape.fill
        if not decoration.background_gradient:
            if decoration.background_color and decoration.opacity > 0:
                fill_color = _with_opacity(
                    decoration.background_color, decoration.opacity
                )
                _set_fill_color(fill, fill_color)
            else:
                fill.background()

    return bg_shape


def _apply_uniform_border(
    slide,
    box: Box,
    decoration: BoxDecoration,
    shape_type,
    bg_shape,
    uniform_border,
) -> None:
    """Apply a uniform border to the background shape, or clear the line if none."""
    left = Emu(px_to_emu(box.x))
    top = Emu(px_to_emu(box.y))
    width = Emu(px_to_emu(box.width))
    height = Emu(px_to_emu(box.height))

    if uniform_border:
        border_color = _with_opacity(uniform_border.color, decoration.opacity)
        if border_color.a > 0:
            if hasattr(bg_shape, "line"):
                _set_line_color(bg_shape.line, border_color)
                bg_shape.line.width = Emu(px_to_emu(uniform_border.width_px))
            else:
                border_shape = slide.shapes.add_shape(
                    shape_type, left, top, width, height
                )
                _remove_theme_style(border_shape)
                border_shape.fill.background()
                if decoration.border_radius_px > 0:
                    _apply_round_rect_radius(
                        border_shape,
                        px_to_emu(box.width),
                        px_to_emu(box.height),
                        px_to_emu(decoration.border_radius_px),
                    )
                _set_line_color(border_shape.line, border_color)
                border_shape.line.width = Emu(px_to_emu(uniform_border.width_px))
        else:
            if hasattr(bg_shape, "line"):
                bg_shape.line.fill.background()
    elif hasattr(bg_shape, "line"):
        bg_shape.line.fill.background()


def _apply_left_accent_bar(
    slide,
    box: Box,
    decoration: BoxDecoration,
    *,
    rotation_3d_x_deg: float = 0.0,
    rotation_3d_y_deg: float = 0.0,
    rotation_3d_z_deg: float = 0.0,
) -> None:
    """Render a left accent bar shape for blockquote-style left-border elements."""
    accent_border = decoration.border_left
    if accent_border.color is None:
        return
    accent_left, accent_top, accent_width, accent_height = (
        _resolve_left_accent_geometry(box, decoration)
    )
    accent_shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        accent_left,
        accent_top,
        accent_width,
        accent_height,
    )
    _remove_theme_style(accent_shape)
    accent_color = _with_opacity(accent_border.color, decoration.opacity)
    _set_fill_color(accent_shape.fill, accent_color)
    accent_shape.line.fill.background()
    _apply_scene3d(
        accent_shape._element.spPr,
        rotation_3d_x_deg=rotation_3d_x_deg,
        rotation_3d_y_deg=rotation_3d_y_deg,
        rotation_3d_z_deg=rotation_3d_z_deg,
    )


def _create_inset_shadow_shapes(
    slide,
    box: Box,
    decoration: BoxDecoration,
    shape_type,
) -> None:
    """Render inset shadow shapes layered on top of the main background shape."""
    for shadow in decoration.box_shadows:
        if not shadow.inset or shadow.color.a <= 0:
            continue
        _add_shadow_shape(
            slide,
            box,
            decoration,
            shape_type,
            shadow,
            inset=True,
        )


def _add_decoration_shape(
    slide,
    box: Box,
    decoration: BoxDecoration,
    *,
    rotation_3d_x_deg: float = 0.0,
    rotation_3d_y_deg: float = 0.0,
    rotation_3d_z_deg: float = 0.0,
):
    """Render a generic decorated box and return it as the text container."""
    shape_type = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        if decoration.border_radius_px > 0
        else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )

    _create_outer_shadow_shapes(slide, box, decoration, shape_type)

    bg_shape = _create_background_shape(slide, box, decoration, shape_type)

    uniform_border = _detect_uniform_border(decoration)
    _apply_uniform_border(slide, box, decoration, shape_type, bg_shape, uniform_border)

    _apply_scene3d(
        bg_shape._element.spPr,
        rotation_3d_x_deg=rotation_3d_x_deg,
        rotation_3d_y_deg=rotation_3d_y_deg,
        rotation_3d_z_deg=rotation_3d_z_deg,
    )

    accent_border = decoration.border_left
    if (
        accent_border.width_px > 0
        and accent_border.color is not None
        and decoration.opacity > 0
        and _is_left_accent_only(decoration)
    ):
        _apply_left_accent_bar(
            slide,
            box,
            decoration,
            rotation_3d_x_deg=rotation_3d_x_deg,
            rotation_3d_y_deg=rotation_3d_y_deg,
            rotation_3d_z_deg=rotation_3d_z_deg,
        )
    elif not uniform_border:
        _add_side_border_shapes(
            slide,
            box,
            decoration,
            rotation_3d_x_deg=rotation_3d_x_deg,
            rotation_3d_y_deg=rotation_3d_y_deg,
            rotation_3d_z_deg=rotation_3d_z_deg,
        )

    _create_inset_shadow_shapes(slide, box, decoration, shape_type)

    return bg_shape


def _set_shape_gradient_fill(shape, css_gradient: str) -> bool:
    """Apply a linear gradient fill directly to a shape spPr node."""
    parsed = parse_linear_gradient(css_gradient)
    if parsed is None:
        return False

    sp_pr = shape._element.spPr
    _remove_existing_fills(sp_pr)

    _build_gradient_fill_xml(sp_pr, parsed)
    return True


def _add_shadow_shape(
    slide,
    box: Box,
    decoration: BoxDecoration,
    shape_type,
    shadow: BoxShadow,
    *,
    inset: bool = False,
):
    """Render a shadow-only helper shape behind a decoration."""
    left = Emu(px_to_emu(box.x))
    top = Emu(px_to_emu(box.y))
    width = Emu(px_to_emu(box.width))
    height = Emu(px_to_emu(box.height))
    shadow_shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    _remove_theme_style(shadow_shape)
    _set_fill_color(shadow_shape.fill, RGBAColor(r=255, g=255, b=255, a=0.0))
    shadow_shape.line.fill.background()
    if decoration.border_radius_px > 0:
        _apply_round_rect_radius(
            shadow_shape,
            px_to_emu(box.width),
            px_to_emu(box.height),
            px_to_emu(decoration.border_radius_px),
        )
    if inset:
        _set_inner_shadow(shadow_shape._element.spPr, shadow, box.width, box.height)
    else:
        _set_outer_shadow(shadow_shape._element.spPr, shadow, box.width, box.height)


def _is_spread_outline_shadow(shadow: BoxShadow) -> bool:
    """Return True for spread-only outer shadows that should render as an outline."""
    return (
        not shadow.inset
        and shadow.spread_px > 0
        and shadow.blur_radius_px <= 0
        and shadow.offset_x_px == 0
        and shadow.offset_y_px == 0
    )


def _add_spread_outline_shape(
    slide,
    box: Box,
    decoration: BoxDecoration,
    shape_type,
    shadow: BoxShadow,
):
    """Render a spread-only outer shadow as an expanded outline shape."""
    spread = shadow.spread_px
    outline_box = Box(
        x=box.x - spread,
        y=box.y - spread,
        width=box.width + (spread * 2),
        height=box.height + (spread * 2),
    )
    left = Emu(px_to_emu(outline_box.x))
    top = Emu(px_to_emu(outline_box.y))
    width = Emu(px_to_emu(outline_box.width))
    height = Emu(px_to_emu(outline_box.height))

    outline_shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    _remove_theme_style(outline_shape)
    outline_shape.fill.background()
    if decoration.border_radius_px > 0:
        _apply_round_rect_radius(
            outline_shape,
            px_to_emu(outline_box.width),
            px_to_emu(outline_box.height),
            px_to_emu(decoration.border_radius_px + spread),
        )
    _set_line_color(outline_shape.line, shadow.color)
    outline_shape.line.width = Emu(px_to_emu(spread * 2))


def _remove_theme_style(shape) -> None:
    """Remove theme style refs that can introduce unwanted effects like shadows."""
    style = shape._element.find(qn("p:style"))
    if style is not None:
        shape._element.remove(style)


def _detect_uniform_border(decoration: BoxDecoration):
    """Return a representative border side when all visible sides are equal."""
    sides = [
        decoration.border_top,
        decoration.border_right,
        decoration.border_bottom,
        decoration.border_left,
    ]
    visible = [
        side
        for side in sides
        if side.width_px > 0 and side.style and side.style != "none" and side.color
    ]
    if not visible:
        return None
    first = visible[0]
    if (
        all(
            side.width_px == first.width_px
            and side.style == first.style
            and side.color == first.color
            for side in visible
        )
        and len(visible) == 4
    ):
        return first
    return None


def _is_left_accent_only(decoration: BoxDecoration) -> bool:
    """Return True when only the left border should be rendered as an accent bar."""
    other_sides = [
        decoration.border_top,
        decoration.border_right,
        decoration.border_bottom,
    ]
    return all(
        side.width_px == 0 or side.style == "none" or side.color is None
        for side in other_sides
    )


def _add_side_border_shapes(
    slide,
    box: Box,
    decoration: BoxDecoration,
    *,
    rotation_3d_x_deg: float = 0.0,
    rotation_3d_y_deg: float = 0.0,
    rotation_3d_z_deg: float = 0.0,
) -> None:
    """Render visible non-uniform borders as separate thin rectangle shapes."""
    sides = [
        (
            decoration.border_top,
            box.x,
            box.y,
            box.width,
            decoration.border_top.width_px,
        ),
        (
            decoration.border_right,
            box.x + box.width - decoration.border_right.width_px,
            box.y,
            decoration.border_right.width_px,
            box.height,
        ),
        (
            decoration.border_bottom,
            box.x,
            box.y + box.height - decoration.border_bottom.width_px,
            box.width,
            decoration.border_bottom.width_px,
        ),
        (
            decoration.border_left,
            box.x,
            box.y,
            decoration.border_left.width_px,
            box.height,
        ),
    ]

    for border, x, y, width, height in sides:
        if (
            border.width_px <= 0
            or border.style == "none"
            or border.color is None
            or width <= 0
            or height <= 0
        ):
            continue
        border_shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Emu(px_to_emu(x)),
            Emu(px_to_emu(y)),
            Emu(px_to_emu(width)),
            Emu(px_to_emu(height)),
        )
        _remove_theme_style(border_shape)
        _set_fill_color(
            border_shape.fill,
            _with_opacity(border.color, decoration.opacity),
        )
        border_shape.line.fill.background()
        _apply_scene3d(
            border_shape._element.spPr,
            rotation_3d_x_deg=rotation_3d_x_deg,
            rotation_3d_y_deg=rotation_3d_y_deg,
            rotation_3d_z_deg=rotation_3d_z_deg,
        )


def _resolve_left_accent_geometry(
    box: Box,
    decoration: BoxDecoration,
) -> tuple[Emu, Emu, Emu, Emu]:
    """Match the visible straight segment of a left border on rounded boxes."""
    radius_px = max(decoration.border_radius_px, 0.0)
    vertical_inset_px = min(radius_px, box.height / 2)
    accent_height_px = max(box.height - (vertical_inset_px * 2), 1.0)

    return (
        Emu(px_to_emu(box.x)),
        Emu(px_to_emu(box.y + vertical_inset_px)),
        Emu(px_to_emu(decoration.border_left.width_px)),
        Emu(px_to_emu(accent_height_px)),
    )


def _add_code_block(slide, element: CodeBlockElement) -> None:
    """Add a code block as textbox with background."""
    from .text import (
        _apply_paragraph_layout,
        _apply_text_style,
        _resolve_textbox_geometry,
        _set_text_frame_margins_zero,
        ALIGNMENT_MAP,
    )
    from pptx.enum.text import PP_ALIGN

    if element.decoration:
        _add_decoration_shape(slide, element.box, element.decoration)

    left, top, width, height = _resolve_textbox_geometry(element)

    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    _set_text_frame_margins_zero(tf)
    txbox.line.fill.background()

    if element.decoration:
        txbox.fill.background()
    elif element.code_background:
        fill = txbox.fill
        _set_fill_color(fill, element.code_background)
    else:
        txbox.fill.background()

    # Add code text with monospace font
    for i, para in enumerate(element.paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = ALIGNMENT_MAP.get(para.alignment, PP_ALIGN.LEFT)
        _apply_paragraph_layout(p, para)

        for text_run in para.runs:
            r = p.add_run()
            r.text = text_run.text
            # Force monospace
            style = text_run.style.model_copy()
            style.font_family = "Courier New"
            _apply_text_style(r, style)


def _add_fallback_image(
    slide,
    element: SlideElement,
    fallback_image_path: str | None = None,
) -> None:
    """Add a fallback image for unsupported content."""
    if not fallback_image_path:
        return

    left = Emu(px_to_emu(element.box.x))
    top = Emu(px_to_emu(element.box.y))
    width = Emu(px_to_emu(element.box.width))
    height = Emu(px_to_emu(element.box.height))

    image_path = Path(fallback_image_path)
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), left, top, width, height)
    else:
        logger.warning(
            "Fallback image not found: %s",
            fallback_image_path,
        )
