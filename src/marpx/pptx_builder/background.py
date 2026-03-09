"""Background image rendering for pptx_builder."""

from __future__ import annotations

import io
import logging

from PIL import Image
from pptx.util import Emu

from marpx.utils.image import resolve_image_bytes
from marpx.models import BackgroundImage
from marpx.utils.common import px_to_emu

from marpx.utils.svg import rasterize_svg_to_png

from .image import _is_svg_source

logger = logging.getLogger(__name__)


def _resolve_split_box(
    bg_image: BackgroundImage, slide_width_emu: int, slide_height_emu: int
) -> tuple[int, int, int, int]:
    """Resolve the destination box for a split background image."""
    if bg_image.box is not None:
        return (
            px_to_emu(bg_image.box.x),
            px_to_emu(bg_image.box.y),
            px_to_emu(bg_image.box.width),
            px_to_emu(bg_image.box.height),
        )

    split_ratio = bg_image.split_ratio if bg_image.split_ratio is not None else 0.5
    split_ratio = max(0.0, min(split_ratio, 1.0))

    if bg_image.split == "left":
        box_width = round(slide_width_emu * split_ratio)
        return 0, 0, box_width, slide_height_emu
    if bg_image.split == "right":
        box_width = round(slide_width_emu * split_ratio)
        return slide_width_emu - box_width, 0, box_width, slide_height_emu
    return 0, 0, slide_width_emu, slide_height_emu


def _add_background_image(
    pptx_slide,
    bg_image: BackgroundImage,
    slide_width_emu: int,
    slide_height_emu: int,
) -> None:
    """Add a background image as a picture shape behind all other content.

    The image is placed at z-index 0 (behind text/shapes) by inserting
    its XML element right after the spTree boilerplate nodes.
    """
    image_bytes = resolve_image_bytes(bg_image.url, bg_image.image_data)

    if not image_bytes:
        return

    if _is_svg_source(bg_image.url):
        image_bytes = rasterize_svg_to_png(bg_image.url)

    image_stream = io.BytesIO(image_bytes)
    image_width_px: int | None = None
    image_height_px: int | None = None
    try:
        with Image.open(io.BytesIO(image_bytes)) as image:
            image_width_px, image_height_px = image.size
    except Exception:
        logger.warning("Failed to inspect background image size: %s", bg_image.url)
        logger.debug(
            "Failed to process background image: %s", bg_image.url, exc_info=True
        )

    box_left, box_top, box_width, box_height = _resolve_split_box(
        bg_image,
        slide_width_emu,
        slide_height_emu,
    )

    left = Emu(box_left)
    top = Emu(box_top)
    width = Emu(box_width)
    height = Emu(box_height)
    crops = (0.0, 0.0, 0.0, 0.0)

    if image_width_px and image_height_px and box_width > 0 and box_height > 0:
        left, top, width, height, crops = _resolve_background_image_placement(
            bg_image,
            image_width_px,
            image_height_px,
            box_left,
            box_top,
            box_width,
            box_height,
        )

    pic = pptx_slide.shapes.add_picture(image_stream, left, top, width, height)
    crop_left, crop_right, crop_top, crop_bottom = crops
    pic.crop_left = crop_left
    pic.crop_right = crop_right
    pic.crop_top = crop_top
    pic.crop_bottom = crop_bottom

    # Move picture to back (behind all other shapes)
    sp_tree = pptx_slide.shapes._spTree
    sp_element = pic._element
    sp_tree.remove(sp_element)
    sp_tree.insert(2, sp_element)  # Insert after cNvPr and cNvSpPr


def _resolve_background_image_placement(
    bg_image: BackgroundImage,
    image_width_px: int,
    image_height_px: int,
    box_left_emu: int,
    box_top_emu: int,
    box_width_emu: int,
    box_height_emu: int,
) -> tuple[Emu, Emu, Emu, Emu, tuple[float, float, float, float]]:
    """Resolve background image placement for contain/cover and position."""
    if image_width_px <= 0 or image_height_px <= 0:
        return (
            Emu(box_left_emu),
            Emu(box_top_emu),
            Emu(box_width_emu),
            Emu(box_height_emu),
            (0.0, 0.0, 0.0, 0.0),
        )

    image_width_emu = px_to_emu(image_width_px)
    image_height_emu = px_to_emu(image_height_px)

    mode = (bg_image.size or "cover").lower()
    if mode == "contain":
        scale = min(
            box_width_emu / image_width_emu,
            box_height_emu / image_height_emu,
        )
    else:
        scale = max(
            box_width_emu / image_width_emu,
            box_height_emu / image_height_emu,
        )

    draw_width = max(round(image_width_emu * scale), 1)
    draw_height = max(round(image_height_emu * scale), 1)

    x_ratio, y_ratio = _parse_position(bg_image.position)
    if mode != "contain":
        overflow_x = max(draw_width - box_width_emu, 0)
        overflow_y = max(draw_height - box_height_emu, 0)
        crop_left = overflow_x * x_ratio / draw_width if draw_width else 0.0
        crop_right = overflow_x * (1 - x_ratio) / draw_width if draw_width else 0.0
        crop_top = overflow_y * y_ratio / draw_height if draw_height else 0.0
        crop_bottom = overflow_y * (1 - y_ratio) / draw_height if draw_height else 0.0
        return (
            Emu(box_left_emu),
            Emu(box_top_emu),
            Emu(box_width_emu),
            Emu(box_height_emu),
            (crop_left, crop_right, crop_top, crop_bottom),
        )

    left = box_left_emu + round((box_width_emu - draw_width) * x_ratio)
    top = box_top_emu + round((box_height_emu - draw_height) * y_ratio)

    return (
        Emu(left),
        Emu(top),
        Emu(draw_width),
        Emu(draw_height),
        (0.0, 0.0, 0.0, 0.0),
    )


def _parse_position(position: str | None) -> tuple[float, float]:
    """Parse CSS background-position to x/y anchor ratios."""
    if not position:
        return 0.5, 0.5

    def parse_token(token: str, axis: str) -> float | None:
        token = token.lower()
        keyword_map = {
            "left": 0.0,
            "center": 0.5,
            "right": 1.0,
            "top": 0.0,
            "bottom": 1.0,
        }
        if token in keyword_map:
            return keyword_map[token]
        if token.endswith("%"):
            try:
                return max(0.0, min(float(token[:-1]) / 100.0, 1.0))
            except ValueError:
                return None
        if axis == "x" and token in {"top", "bottom"}:
            return 0.5
        if axis == "y" and token in {"left", "right"}:
            return 0.5
        return None

    tokens = position.lower().split()
    if not tokens:
        return 0.5, 0.5
    if len(tokens) == 1:
        token = tokens[0]
        if token in {"top", "bottom"}:
            return 0.5, 0.0 if token == "top" else 1.0
        if token in {"left", "right"}:
            return 0.0 if token == "left" else 1.0, 0.5
        parsed = parse_token(token, "x")
        return (parsed if parsed is not None else 0.5), 0.5

    if tokens[0] in {"top", "bottom"} and tokens[1] in {"left", "right"}:
        x_ratio = parse_token(tokens[1], "x")
        y_ratio = parse_token(tokens[0], "y")
        return (
            x_ratio if x_ratio is not None else 0.5,
            y_ratio if y_ratio is not None else 0.5,
        )

    x_ratio = parse_token(tokens[0], "x")
    y_ratio = parse_token(tokens[1], "y")
    return (
        x_ratio if x_ratio is not None else 0.5,
        y_ratio if y_ratio is not None else 0.5,
    )
