"""Main build_pptx() function and slide background helper."""

from __future__ import annotations

import io
import logging
from pathlib import Path

from pptx import Presentation as PptxPresentation
from pptx.util import Emu

from marpx.models import (
    Background,
    CodeBlockElement,
    ImageElement,
    ListElement,
    Presentation,
    TableElement,
    TextElement,
    UnsupportedElement,
)
from marpx.extraction.capabilities import Capability
from marpx.pipeline import SlideRenderInfo
from marpx.utils.common import px_to_emu

from ._helpers import _set_fill_color
from .background import _add_background_image
from .code_block import _add_code_block
from .fallback import _add_fallback_image
from marpx.utils.gradient import render_gradient_png
from .directives import _add_footer, _add_header, _add_page_number
from .image import MissingDependencyError, _add_image
from .table import _add_table
from .text import _add_textbox
from .text_grouping import (
    _add_grouped_textbox,
    _group_adjacent_text_elements,
    _is_groupable_text_element,
)

logger = logging.getLogger(__name__)


def _set_slide_background(
    pptx_slide, background: Background, slide_width_px: float, slide_height_px: float
) -> None:
    """Set solid background color on a slide."""
    if background.background_gradient:
        gradient_bytes = render_gradient_png(
            background.background_gradient,
            max(int(round(slide_width_px)), 1),
            max(int(round(slide_height_px)), 1),
        )
        if gradient_bytes:
            pptx_slide.shapes.add_picture(
                io.BytesIO(gradient_bytes),
                Emu(0),
                Emu(0),
                Emu(px_to_emu(slide_width_px)),
                Emu(px_to_emu(slide_height_px)),
            )
            return
    if background.color and background.color.a > 0:
        bg = pptx_slide.background
        fill = bg.fill
        _set_fill_color(fill, background.color)


def build_pptx(
    presentation: Presentation,
    output_path: str | Path,
    slide_render_info: dict[int, SlideRenderInfo] | None = None,
) -> Path:
    """Build a PPTX file from a Presentation model.

    Args:
        presentation: Normalized presentation data.
        output_path: Path for the output .pptx file.
        slide_render_info: Rendering decisions from the pipeline.

    Returns:
        Path to the generated PPTX file.
    """
    if slide_render_info is None:
        slide_render_info = {}

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    pptx = PptxPresentation()
    if presentation.slides:
        slide_sizes = {
            (round(slide.width_px, 4), round(slide.height_px, 4))
            for slide in presentation.slides
        }
        if len(slide_sizes) > 1:
            raise ValueError(
                "PowerPoint does not support mixed slide sizes in one deck."
            )
        first_slide = presentation.slides[0]
        pptx.slide_width = Emu(px_to_emu(first_slide.width_px))
        pptx.slide_height = Emu(px_to_emu(first_slide.height_px))

    for slide_idx, slide_data in enumerate(presentation.slides):
        s_info = slide_render_info.get(slide_idx, SlideRenderInfo())

        # Add blank slide
        layout = pptx.slide_layouts[6]  # Blank layout
        pptx_slide = pptx.slides.add_slide(layout)

        # Handle full-slide fallback
        is_full_slide_fallback = bool(s_info.is_fallback and s_info.fallback_image_path)
        if is_full_slide_fallback:
            img_path = Path(s_info.fallback_image_path)
            if img_path.exists():
                pptx_slide.shapes.add_picture(
                    str(img_path),
                    Emu(0),
                    Emu(0),
                    Emu(px_to_emu(slide_data.width_px)),
                    Emu(px_to_emu(slide_data.height_px)),
                )
        else:
            # Set background
            _set_slide_background(
                pptx_slide,
                slide_data.background,
                slide_data.width_px,
                slide_data.height_px,
            )

            # Add background images (behind content)
            if slide_data.background.images:
                slide_w_emu = px_to_emu(slide_data.width_px)
                slide_h_emu = px_to_emu(slide_data.height_px)
                for bg_image in slide_data.background.images:
                    try:
                        _add_background_image(
                            pptx_slide,
                            bg_image,
                            slide_w_emu,
                            slide_h_emu,
                        )
                    except Exception as e:
                        logger.warning("Failed to add background image: %s", e)
                        logger.debug(
                            "Failed to process background: %s", e, exc_info=True
                        )

            # Add elements sorted by z-index, grouping adjacent plain text content
            skipped_count = 0
            _element_to_idx = {
                id(el): idx for idx, el in enumerate(slide_data.elements)
            }
            sorted_elements = sorted(slide_data.elements, key=lambda e: e.z_index)
            for group in _group_adjacent_text_elements(sorted_elements):
                if len(group) > 1 and all(
                    _is_groupable_text_element(el) for el in group
                ):
                    _add_grouped_textbox(pptx_slide, group)
                    continue

                element = group[0]
                # Look up element render info for capability-driven dispatch
                el_idx = _element_to_idx.get(id(element))
                el_info = (
                    s_info.element_info.get(el_idx) if el_idx is not None else None
                )

                try:
                    # Check if this element should be rendered as fallback
                    if (
                        el_info is not None
                        and el_info.capability == Capability.SUBTREE_FALLBACK
                    ):
                        _add_fallback_image(
                            pptx_slide, element, el_info.fallback_image_path
                        )
                    elif isinstance(element, (TextElement, ListElement)):
                        _add_textbox(pptx_slide, element)
                    elif isinstance(element, ImageElement):
                        _add_image(pptx_slide, element)
                    elif isinstance(element, TableElement):
                        _add_table(pptx_slide, element)
                    elif isinstance(element, CodeBlockElement):
                        _add_code_block(pptx_slide, element)
                    elif isinstance(element, UnsupportedElement):
                        fallback_path = el_info.fallback_image_path if el_info else None
                        _add_fallback_image(pptx_slide, element, fallback_path)
                except MissingDependencyError:
                    raise
                except Exception as e:
                    logger.warning(
                        "Failed to add element %s: %s",
                        element.element_type.value,
                        e,
                    )
                    logger.debug("Failed to build element: %s", e, exc_info=True)
                    skipped_count += 1

            if skipped_count:
                logger.info(
                    "Slide %d: skipped %d element(s)",
                    slide_idx + 1,
                    skipped_count,
                )

            # Add directives (header, footer, page number)
            slide_w_emu = px_to_emu(slide_data.width_px)
            slide_h_emu = px_to_emu(slide_data.height_px)

            if slide_data.header_text:
                _add_header(pptx_slide, slide_data.header_text, slide_w_emu)

            if slide_data.footer_text:
                _add_footer(
                    pptx_slide, slide_data.footer_text, slide_w_emu, slide_h_emu
                )

            if slide_data.paginate and slide_data.page_number is not None:
                _add_page_number(
                    pptx_slide,
                    slide_data.page_number,
                    slide_data.page_total,
                    slide_w_emu,
                    slide_h_emu,
                )

        # Add speaker notes
        if slide_data.notes:
            notes_slide = pptx_slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = slide_data.notes

    pptx.save(str(output_path))
    logger.info("PPTX saved: %s", output_path)
    return output_path
