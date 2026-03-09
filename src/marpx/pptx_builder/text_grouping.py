"""Text grouping/merging logic for pptx_builder."""

from __future__ import annotations

from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.util import Emu

from marpx.models import ListElement, SlideElement, TextElement
from marpx.utils import (
    boxes_have_horizontal_overlap,
    boxes_have_mergeable_vertical_gap,
    boxes_share_column,
    px_to_emu,
    union_boxes,
)

from .text import (
    GROUPABLE_TEXT_TYPES,
    _populate_text_frame,
    _set_text_frame_margins_zero,
)


def _is_groupable_text_element(element: SlideElement) -> bool:
    """Return True when the element should be grouped into a shared textbox."""
    if not isinstance(element, (TextElement, ListElement)):
        return False
    return (
        element.element_type in GROUPABLE_TEXT_TYPES
        and element.decoration is None
        and abs(element.rotation_deg) <= 0.01
        and abs(element.rotation_3d_x_deg) <= 0.01
        and abs(element.rotation_3d_y_deg) <= 0.01
        and abs(element.rotation_3d_z_deg) <= 0.01
    )


def _text_elements_should_merge(first: SlideElement, second: SlideElement) -> bool:
    """Heuristic for merging adjacent text-capable elements into one textbox."""
    if not (_is_groupable_text_element(first) and _is_groupable_text_element(second)):
        return False

    same_column = boxes_share_column(first.box, second.box)
    similar_width = abs(first.box.width - second.box.width) <= 64
    aligned = same_column or (
        similar_width and boxes_have_horizontal_overlap(first.box, second.box)
    )
    return aligned and boxes_have_mergeable_vertical_gap(
        first.box,
        second.box,
        min_gap_px=-8,
        max_gap_px=64,
    )


def _group_adjacent_text_elements(
    elements: list[SlideElement],
) -> list[list[SlideElement]]:
    """Group adjacent text elements that should share one text frame."""
    groups: list[list[SlideElement]] = []
    current_group: list[SlideElement] = []

    for element in elements:
        if not current_group:
            current_group = [element]
            continue

        if _text_elements_should_merge(current_group[-1], element):
            current_group.append(element)
        else:
            groups.append(current_group)
            current_group = [element]

    if current_group:
        groups.append(current_group)

    return groups


def _add_grouped_textbox(slide, elements: list[TextElement | ListElement]) -> None:
    """Render multiple adjacent text-like elements into one textbox."""
    box = union_boxes([element.box for element in elements])
    txbox = slide.shapes.add_textbox(
        Emu(px_to_emu(box.x)),
        Emu(px_to_emu(box.y)),
        Emu(px_to_emu(max(box.width, 1))),
        Emu(px_to_emu(max(box.height, 1))),
    )
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    _set_text_frame_margins_zero(tf)
    _populate_text_frame(tf, elements)
