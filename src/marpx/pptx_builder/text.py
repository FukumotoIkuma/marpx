"""Text rendering functions for pptx_builder."""

from __future__ import annotations

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from marpx.models import (
    BaseSlideElement,
    ElementType,
    ListElement,
    SlideElement,
    TextElement,
    TextRun,
    TextStyle,
)
from marpx.gradient_utils import parse_linear_gradient
from marpx.utils import (
    blend_alpha,
    px_to_emu,
    px_to_pt,
)

from ._helpers import _build_gradient_fill_xml, _remove_existing_fills, _to_rgb
from .decoration import _add_decoration_shape
from .decoration import _resolve_scene3d_rotations

# Fill tags relevant to run properties (no pattFill).
_RUN_FILL_TAGS: frozenset[str] = frozenset(
    {qn("a:solidFill"), qn("a:gradFill"), qn("a:noFill"), qn("a:blipFill")}
)

ALIGNMENT_MAP: dict[str, PP_ALIGN] = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
    "start": PP_ALIGN.LEFT,
    "end": PP_ALIGN.RIGHT,
}

VERTICAL_ALIGNMENT_MAP = {
    "top": MSO_VERTICAL_ANCHOR.TOP,
    "middle": MSO_VERTICAL_ANCHOR.MIDDLE,
    "bottom": MSO_VERTICAL_ANCHOR.BOTTOM,
}

GROUPABLE_TEXT_TYPES: tuple[ElementType, ...] = (
    ElementType.HEADING,
    ElementType.PARAGRAPH,
    ElementType.BLOCKQUOTE,
    ElementType.UNORDERED_LIST,
    ElementType.ORDERED_LIST,
)


def _apply_text_style(run, style: TextStyle) -> None:
    """Apply TextStyle to a python-pptx run."""
    run.font.name = style.font_family
    run.font.size = Pt(px_to_pt(style.font_size_px))
    run.font.bold = style.bold
    run.font.italic = style.italic
    run.font.underline = style.underline
    r_pr = run._r.get_or_add_rPr()
    if style.strike:
        r_pr.set("strike", "sngStrike")
    elif "strike" in r_pr.attrib:
        del r_pr.attrib["strike"]
    if style.color.a < 1.0:
        run.font.color.rgb = RGBColor(style.color.r, style.color.g, style.color.b)
        solid_fill = r_pr.find(qn("a:solidFill"))
        if solid_fill is None:
            solid_fill = etree.SubElement(r_pr, qn("a:solidFill"))
        srgb = solid_fill.find(qn("a:srgbClr"))
        if srgb is None:
            srgb = etree.SubElement(solid_fill, qn("a:srgbClr"))
        srgb.set("val", f"{style.color.r:02X}{style.color.g:02X}{style.color.b:02X}")
        for child in list(srgb):
            if child.tag == qn("a:alpha"):
                srgb.remove(child)
        alpha = etree.SubElement(srgb, qn("a:alpha"))
        alpha.set("val", str(int(style.color.a * 100000)))
    else:
        run.font.color.rgb = _to_rgb(style.color)
    if style.text_gradient:
        _set_run_gradient_fill(r_pr, style.text_gradient)
    if style.background_color and style.background_color.a > 0:
        r_pr = run._r.get_or_add_rPr()
        existing = r_pr.find(qn("a:highlight"))
        if existing is not None:
            r_pr.remove(existing)
        highlight = etree.Element(qn("a:highlight"))
        bg = blend_alpha(style.background_color)
        color = etree.SubElement(highlight, qn("a:srgbClr"))
        color.set("val", f"{bg.r:02X}{bg.g:02X}{bg.b:02X}")
        insert_before = None
        for child in r_pr:
            if child.tag in {
                qn("a:latin"),
                qn("a:ea"),
                qn("a:cs"),
                qn("a:sym"),
                qn("a:hlinkClick"),
                qn("a:hlinkMouseOver"),
            }:
                insert_before = child
                break
        if insert_before is None:
            r_pr.append(highlight)
        else:
            insert_before.addprevious(highlight)


def _set_run_gradient_fill(r_pr, css_gradient: str) -> None:
    """Apply a linear gradient fill directly to a run properties node."""
    parsed = parse_linear_gradient(css_gradient)
    if parsed is None:
        return

    _remove_existing_fills(r_pr, fill_tags=_RUN_FILL_TAGS)

    # Build a detached gradFill element (parent_node=None) so we can
    # insert it at the correct position within r_pr afterwards.
    grad_fill = _build_gradient_fill_xml(
        None,
        parsed,
        extra_attrs={"flip": "none"},
        extra_children=[etree.Element(qn("a:tileRect"))],
    )

    insert_before = None
    for child in r_pr:
        if child.tag in {
            qn("a:latin"),
            qn("a:ea"),
            qn("a:cs"),
            qn("a:sym"),
            qn("a:hlinkClick"),
            qn("a:hlinkMouseOver"),
        }:
            insert_before = child
            break
    if insert_before is None:
        r_pr.append(grad_fill)
    else:
        insert_before.addprevious(grad_fill)


def _add_paragraph_runs(pptx_para, runs: list[TextRun]) -> None:
    """Add text runs to a python-pptx paragraph."""
    for i, text_run in enumerate(runs):
        if i == 0:
            r = pptx_para.runs[0] if pptx_para.runs else pptx_para.add_run()
        else:
            r = pptx_para.add_run()
        r.text = text_run.text
        _apply_text_style(r, text_run.style)
        if text_run.link_url:
            r.hyperlink.address = text_run.link_url


def _ordered_numbering_type(style_type: str | None) -> str:
    """Map CSS list-style-type to OOXML numbering type."""
    normalized = (style_type or "").lower()
    return {
        "decimal": "arabicPeriod",
        "decimal-leading-zero": "arabicPeriod",
        "lower-alpha": "alphaLcPeriod",
        "lower-latin": "alphaLcPeriod",
        "upper-alpha": "alphaUcPeriod",
        "upper-latin": "alphaUcPeriod",
        "lower-roman": "romanLcPeriod",
        "upper-roman": "romanUcPeriod",
    }.get(normalized, "arabicPeriod")


def _unordered_bullet_char(style_type: str | None) -> str:
    """Map CSS list-style-type to a bullet character."""
    normalized = (style_type or "").lower()
    return {
        "circle": "\u25e6",
        "square": "\u25aa",
    }.get(normalized, "\u2022")


def _configure_list_paragraph(
    pptx_para,
    level: int,
    ordered: bool,
    list_style_type: str | None = None,
    order_number: int | None = None,
) -> None:
    """Apply bullet or numbering settings to a paragraph."""
    pptx_para.level = level
    pPr = pptx_para._p.get_or_add_pPr()
    normalized_style = (list_style_type or "").lower()
    if normalized_style == "none":
        etree.SubElement(pPr, qn("a:buNone"))
        return
    if ordered:
        buAutoNum = etree.SubElement(pPr, qn("a:buAutoNum"))
        buAutoNum.set("type", _ordered_numbering_type(list_style_type))
        if order_number and order_number > 1:
            buAutoNum.set("startAt", str(order_number))
    else:
        buChar = etree.SubElement(pPr, qn("a:buChar"))
        buChar.set("char", _unordered_bullet_char(list_style_type))


def _iter_text_payloads(element: TextElement | ListElement):
    """Yield paragraph payloads for a groupable text element."""
    if isinstance(element, TextElement):
        for para in element.paragraphs:
            yield {
                "paragraph": para,
                "is_list": False,
                "ordered": False,
                "level": 0,
            }
        return

    if isinstance(element, ListElement):
        for item in element.list_items:
            yield {
                "paragraph": None,
                "runs": item.runs,
                "is_list": True,
                "ordered": element.element_type == ElementType.ORDERED_LIST,
                "level": item.level,
                "style_type": item.list_style_type,
                "order_number": item.order_number,
                "alignment": item.alignment,
                "line_height_px": item.line_height_px,
                "space_before_px": item.space_before_px,
                "space_after_px": item.space_after_px,
            }


def _append_payload_to_text_frame(
    text_frame,
    payload,
    index: int,
    *,
    suppress_line_spacing: bool = False,
) -> None:
    """Append a paragraph payload to a text frame."""
    if index == 0:
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()

    if payload["is_list"]:
        _configure_list_paragraph(
            p,
            payload["level"],
            payload["ordered"],
            payload.get("style_type"),
            payload.get("order_number"),
        )
        p.alignment = ALIGNMENT_MAP.get(payload.get("alignment", "left"), PP_ALIGN.LEFT)
        _apply_spacing(
            p,
            None if suppress_line_spacing else payload.get("line_height_px"),
            payload.get("space_before_px", 0.0),
            payload.get("space_after_px", 0.0),
        )
        _add_paragraph_runs(p, payload["runs"])
        return

    para = payload["paragraph"]
    if para.list_level is not None:
        _configure_list_paragraph(
            p,
            para.list_level,
            para.list_ordered,
            para.list_style_type,
            para.order_number,
        )
    p.alignment = ALIGNMENT_MAP.get(para.alignment, PP_ALIGN.LEFT)
    _apply_paragraph_layout(p, para, suppress_line_spacing=suppress_line_spacing)
    _add_paragraph_runs(p, para.runs)


def _populate_text_frame(text_frame, elements: list[TextElement | ListElement]) -> None:
    """Populate a text frame from one or more text-bearing elements."""
    payloads: list[tuple[SlideElement, dict]] = []
    for element in elements:
        for payload in _iter_text_payloads(element):
            payloads.append((element, payload))

    suppress_line_spacing = False
    if len(payloads) == 1:
        element, payload = payloads[0]
        if (
            element.vertical_align == "middle"
            and not payload["is_list"]
            and len(element.paragraphs) == 1
        ):
            para = payload["paragraph"]
            if para is not None and all("\n" not in run.text for run in para.runs):
                suppress_line_spacing = True

    payload_index = 0
    for _element, payload in payloads:
        _append_payload_to_text_frame(
            text_frame,
            payload,
            payload_index,
            suppress_line_spacing=suppress_line_spacing,
        )
        payload_index += 1


def _resolve_textbox_geometry(
    element: TextElement | ListElement,
) -> tuple[Emu, Emu, Emu, Emu]:
    """Return the element's content box when available, else its outer box."""
    if element.content_box is not None:
        source_box = element.content_box
    elif element.decoration is not None:
        source_box = _derive_content_box_from_decoration(element)
    else:
        source_box = element.box
    left_px = source_box.x
    top_px = source_box.y
    width_px = max(source_box.width, 1)
    height_px = max(source_box.height, 1)
    return (
        Emu(px_to_emu(left_px)),
        Emu(px_to_emu(top_px)),
        Emu(px_to_emu(width_px)),
        Emu(px_to_emu(height_px)),
    )


def _derive_content_box_from_decoration(element: BaseSlideElement):
    """Derive a content box from an outer box plus extracted border/padding."""
    decoration = element.decoration
    assert decoration is not None
    left_inset = decoration.border_left.width_px + decoration.padding.left_px
    top_inset = decoration.border_top.width_px + decoration.padding.top_px
    right_inset = decoration.border_right.width_px + decoration.padding.right_px
    bottom_inset = decoration.border_bottom.width_px + decoration.padding.bottom_px
    return element.box.model_copy(
        update={
            "x": element.box.x + left_inset,
            "y": element.box.y + top_inset,
            "width": max(element.box.width - left_inset - right_inset, 1),
            "height": max(element.box.height - top_inset - bottom_inset, 1),
        }
    )


def _set_text_frame_margins_zero(text_frame) -> None:
    """Use extracted box geometry as-is and disable PowerPoint default insets."""
    zero = Emu(0)
    text_frame.margin_top = zero
    text_frame.margin_right = zero
    text_frame.margin_bottom = zero
    text_frame.margin_left = zero


def _set_text_frame_margins_from_element(text_frame, element: BaseSlideElement) -> None:
    """Apply content-box insets as text-frame margins on the element's outer box."""
    if element.content_box is not None:
        content_box = element.content_box
    elif element.decoration is not None:
        content_box = _derive_content_box_from_decoration(element)
    else:
        _set_text_frame_margins_zero(text_frame)
        return

    left = max(content_box.x - element.box.x, 0)
    top = max(content_box.y - element.box.y, 0)
    right = max(
        (element.box.x + element.box.width) - (content_box.x + content_box.width), 0
    )
    bottom = max(
        (element.box.y + element.box.height) - (content_box.y + content_box.height), 0
    )
    text_frame.margin_left = Emu(px_to_emu(left))
    text_frame.margin_top = Emu(px_to_emu(top))
    text_frame.margin_right = Emu(px_to_emu(right))
    text_frame.margin_bottom = Emu(px_to_emu(bottom))


def _apply_paragraph_layout(
    pptx_para, para, *, suppress_line_spacing: bool = False
) -> None:
    """Apply paragraph-level layout such as line-height and spacing."""
    _apply_spacing(
        pptx_para,
        None if suppress_line_spacing else para.line_height_px,
        para.space_before_px,
        para.space_after_px,
    )


def _apply_spacing(
    pptx_para,
    line_height_px: float | None,
    space_before_px: float,
    space_after_px: float,
) -> None:
    """Apply line-height and paragraph spacing values."""
    if line_height_px and line_height_px > 0:
        pptx_para.line_spacing = Pt(px_to_pt(line_height_px))
    if space_before_px > 0:
        pptx_para.space_before = Pt(px_to_pt(space_before_px))
    if space_after_px > 0:
        pptx_para.space_after = Pt(px_to_pt(space_after_px))


def _add_textbox(slide, element: TextElement | ListElement) -> None:
    """Add a textbox or decorated text container."""
    if element.decoration:
        use_shape_text_frame = (
            abs(element.rotation_3d_x_deg) > 0.01
            or abs(element.rotation_3d_y_deg) > 0.01
            or abs(element.rotation_3d_z_deg) > 0.01
        )
        scene3d_x_deg, scene3d_y_deg, scene3d_z_deg = _resolve_scene3d_rotations(
            element
        )
        bg_shape = _add_decoration_shape(
            slide,
            element.box,
            element.decoration,
            rotation_3d_x_deg=scene3d_x_deg,
            rotation_3d_y_deg=scene3d_y_deg,
            rotation_3d_z_deg=scene3d_z_deg,
        )
        has_content = (isinstance(element, TextElement) and element.paragraphs) or (
            isinstance(element, ListElement) and element.list_items
        )
        if not has_content:
            return
        if use_shape_text_frame:
            text_container = bg_shape
        else:
            left, top, width, height = _resolve_textbox_geometry(element)
            text_container = slide.shapes.add_textbox(left, top, width, height)
            text_container.fill.background()
            text_container.line.fill.background()
    else:
        left, top, width, height = _resolve_textbox_geometry(element)
        text_container = slide.shapes.add_textbox(left, top, width, height)

    tf = text_container.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = VERTICAL_ALIGNMENT_MAP.get(
        element.vertical_align,
        MSO_VERTICAL_ANCHOR.TOP,
    )
    if element.decoration and (
        abs(element.rotation_3d_x_deg) > 0.01
        or abs(element.rotation_3d_y_deg) > 0.01
        or abs(element.rotation_3d_z_deg) > 0.01
    ):
        _set_text_frame_margins_from_element(tf, element)
    else:
        _set_text_frame_margins_zero(tf)
    _populate_text_frame(tf, [element])
