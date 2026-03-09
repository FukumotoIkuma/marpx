"""Shadow rendering for decoration shapes."""

from __future__ import annotations

from pptx.util import Emu

from marpx.models import Box, BoxDecoration, BoxShadow, RGBAColor
from marpx.utils.common import px_to_emu

from .._helpers import (
    _set_fill_color,
    _set_inner_shadow,
    _set_line_color,
    _set_outer_shadow,
)
from .shapes import _apply_round_rect_radius, _remove_theme_style

# Alpha thresholds for filtering visually imperceptible shadows.
# Spread-only outlines (concentric rings etc.) render as visible coloured
# lines, so they use a very low threshold.  Blur-based shadows create
# transparent helper shapes that clutter the z-order, so they need a
# higher bar.  Inset shadows sit ON TOP of content and are the most
# disruptive, so they use the strictest threshold.
_SPREAD_OUTLINE_ALPHA_THRESHOLD = 0.01
_OUTER_SHADOW_ALPHA_THRESHOLD = 0.10
_INSET_SHADOW_ALPHA_THRESHOLD = 0.15

# Maximum number of blur-based shadow shapes per element.  CSS often uses
# 10+ layered box-shadows for depth, but each becomes a separate transparent
# shape in PPTX.  We keep only the most prominent ones.
_MAX_OUTER_BLUR_SHADOWS = 1
_MAX_INSET_SHADOWS = 1


def _shadow_prominence(shadow: BoxShadow) -> float:
    """Score a shadow by visual weight (higher = more prominent)."""
    return shadow.color.a * max(shadow.blur_radius_px, shadow.spread_px, 1.0)


def _create_outer_shadow_shapes(
    slide,
    box: Box,
    decoration: BoxDecoration,
    shape_type,
) -> None:
    """Render outer (non-inset) shadow shapes behind the main background shape."""
    outlines = []
    blurs = []
    for shadow in decoration.box_shadows:
        if shadow.inset:
            continue
        if _is_spread_outline_shadow(shadow):
            if shadow.color.a >= _SPREAD_OUTLINE_ALPHA_THRESHOLD:
                outlines.append(shadow)
        else:
            if shadow.color.a >= _OUTER_SHADOW_ALPHA_THRESHOLD:
                blurs.append(shadow)

    # Spread-only outlines: keep all (they are visible coloured lines).
    for shadow in outlines:
        _add_spread_outline_shape(slide, box, decoration, shape_type, shadow)

    # Blur shadows: keep only the most prominent ones.
    blurs.sort(key=_shadow_prominence, reverse=True)
    for shadow in blurs[:_MAX_OUTER_BLUR_SHADOWS]:
        _add_shadow_shape(slide, box, decoration, shape_type, shadow)


def _add_shadow_shape(
    slide,
    box: Box,
    decoration: BoxDecoration,
    shape_type,
    shadow: BoxShadow,
    *,
    inset: bool = False,
):
    """Render a shadow-only helper shape behind a decoration."""
    left = Emu(px_to_emu(box.x))
    top = Emu(px_to_emu(box.y))
    width = Emu(px_to_emu(box.width))
    height = Emu(px_to_emu(box.height))
    shadow_shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    _remove_theme_style(shadow_shape)
    _set_fill_color(shadow_shape.fill, RGBAColor(r=255, g=255, b=255, a=0.0))
    shadow_shape.line.fill.background()
    if decoration.border_radius_px > 0:
        _apply_round_rect_radius(
            shadow_shape,
            px_to_emu(box.width),
            px_to_emu(box.height),
            px_to_emu(decoration.border_radius_px),
        )
    if inset:
        _set_inner_shadow(shadow_shape._element.spPr, shadow, box.width, box.height)
    else:
        _set_outer_shadow(shadow_shape._element.spPr, shadow, box.width, box.height)


def _is_spread_outline_shadow(shadow: BoxShadow) -> bool:
    """Return True for spread-only outer shadows that should render as an outline."""
    return (
        not shadow.inset
        and shadow.spread_px > 0
        and shadow.blur_radius_px <= 0
        and shadow.offset_x_px == 0
        and shadow.offset_y_px == 0
    )


def _add_spread_outline_shape(
    slide,
    box: Box,
    decoration: BoxDecoration,
    shape_type,
    shadow: BoxShadow,
):
    """Render a spread-only outer shadow as an expanded outline shape."""
    spread = shadow.spread_px
    outline_box = Box(
        x=box.x - spread,
        y=box.y - spread,
        width=box.width + (spread * 2),
        height=box.height + (spread * 2),
    )
    left = Emu(px_to_emu(outline_box.x))
    top = Emu(px_to_emu(outline_box.y))
    width = Emu(px_to_emu(outline_box.width))
    height = Emu(px_to_emu(outline_box.height))

    outline_shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    _remove_theme_style(outline_shape)
    outline_shape.fill.background()
    if decoration.border_radius_px > 0:
        _apply_round_rect_radius(
            outline_shape,
            px_to_emu(outline_box.width),
            px_to_emu(outline_box.height),
            px_to_emu(decoration.border_radius_px + spread),
        )
    _set_line_color(outline_shape.line, shadow.color)
    outline_shape.line.width = Emu(px_to_emu(spread * 2))


def _create_inset_shadow_shapes(
    slide,
    box: Box,
    decoration: BoxDecoration,
    shape_type,
) -> None:
    """Render inset shadow shapes layered on top of the main background shape."""
    candidates = [
        s
        for s in decoration.box_shadows
        if s.inset and s.color.a >= _INSET_SHADOW_ALPHA_THRESHOLD
    ]
    candidates.sort(key=_shadow_prominence, reverse=True)
    for shadow in candidates[:_MAX_INSET_SHADOWS]:
        _add_shadow_shape(
            slide,
            box,
            decoration,
            shape_type,
            shadow,
            inset=True,
        )
