"""Background shapes and geometry helpers for decoration rendering."""

from __future__ import annotations

import io
import logging

from lxml import etree
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Emu

from marpx.models import Box, BoxDecoration, ClipPath
from marpx.utils.common import emu_to_px, px_to_emu
from marpx.utils.gradient import parse_linear_gradient, render_gradient_png

from .._helpers import (
    _build_gradient_fill_xml,
    _remove_existing_fills,
    _set_fill_color,
    _with_opacity,
)

logger = logging.getLogger(__name__)


def _round_rect_adjustment(
    width: int | float, height: int | float, radius: int | float
) -> float:
    """Return normalized roundRect adjustment matching CSS border-radius."""
    min_dim = min(float(width), float(height))
    if min_dim <= 0:
        return 0.0
    return max(0.0, min(float(radius) / min_dim, 0.5))


def _apply_round_rect_radius(
    shape, width: int | float, height: int | float, radius: int | float
) -> None:
    """Apply a CSS-like corner radius to a rounded rectangle auto shape."""
    if radius <= 0 or width <= 0 or height <= 0:
        return
    if getattr(shape, "adjustments", None) and len(shape.adjustments) > 0:
        shape.adjustments[0] = _round_rect_adjustment(width, height, radius)


def _replace_with_round_rect_custgeom(
    shape, w_emu: int, h_emu: int, radius_emu: int
) -> None:
    """Replace a shape's preset geometry with a custGeom rounded rectangle.

    Unlike ``prstGeom="roundRect"``, this custom geometry sets the text
    rectangle to the full shape bounds so that PowerPoint does not apply
    an automatic text-area inset based on the corner radius.
    """
    sp_pr = shape._element.spPr

    prst_geom = sp_pr.find(qn("a:prstGeom"))
    if prst_geom is not None:
        sp_pr.remove(prst_geom)

    r = min(radius_emu, w_emu // 2, h_emu // 2)

    cust_geom = etree.SubElement(sp_pr, qn("a:custGeom"))
    etree.SubElement(cust_geom, qn("a:avLst"))
    etree.SubElement(cust_geom, qn("a:gdLst"))
    etree.SubElement(cust_geom, qn("a:ahLst"))
    etree.SubElement(cust_geom, qn("a:cxnLst"))

    # Text rectangle = full shape bounds (no geometry inset)
    rect = etree.SubElement(cust_geom, qn("a:rect"))
    rect.set("l", "0")
    rect.set("t", "0")
    rect.set("r", str(w_emu))
    rect.set("b", str(h_emu))

    path_lst = etree.SubElement(cust_geom, qn("a:pathLst"))
    path = etree.SubElement(
        path_lst,
        qn("a:path"),
        attrib={"w": str(w_emu), "h": str(h_emu)},
    )

    # Draw rounded rectangle: moveTo top-left+r, then lines+arcs clockwise
    _move = etree.SubElement(path, qn("a:moveTo"))
    etree.SubElement(_move, qn("a:pt"), attrib={"x": str(r), "y": "0"})

    # Top edge -> top-right corner
    _ln = etree.SubElement(path, qn("a:lnTo"))
    etree.SubElement(_ln, qn("a:pt"), attrib={"x": str(w_emu - r), "y": "0"})
    _arc = etree.SubElement(path, qn("a:arcTo"))
    _arc.set("wR", str(r))
    _arc.set("hR", str(r))
    _arc.set("stAng", "16200000")  # 270
    _arc.set("swAng", "5400000")  # +90

    # Right edge -> bottom-right corner
    _ln = etree.SubElement(path, qn("a:lnTo"))
    etree.SubElement(_ln, qn("a:pt"), attrib={"x": str(w_emu), "y": str(h_emu - r)})
    _arc = etree.SubElement(path, qn("a:arcTo"))
    _arc.set("wR", str(r))
    _arc.set("hR", str(r))
    _arc.set("stAng", "0")
    _arc.set("swAng", "5400000")

    # Bottom edge -> bottom-left corner
    _ln = etree.SubElement(path, qn("a:lnTo"))
    etree.SubElement(_ln, qn("a:pt"), attrib={"x": str(r), "y": str(h_emu)})
    _arc = etree.SubElement(path, qn("a:arcTo"))
    _arc.set("wR", str(r))
    _arc.set("hR", str(r))
    _arc.set("stAng", "5400000")  # 90
    _arc.set("swAng", "5400000")

    # Left edge -> top-left corner
    _ln = etree.SubElement(path, qn("a:lnTo"))
    etree.SubElement(_ln, qn("a:pt"), attrib={"x": "0", "y": str(r)})
    _arc = etree.SubElement(path, qn("a:arcTo"))
    _arc.set("wR", str(r))
    _arc.set("hR", str(r))
    _arc.set("stAng", "10800000")  # 180
    _arc.set("swAng", "5400000")

    etree.SubElement(path, qn("a:close"))


def _apply_round_rect_radius_to_geom(
    prst_geom, width: int | float, height: int | float, radius: int | float
) -> None:
    """Apply a CSS-like corner radius to a preset geometry node."""
    if radius <= 0 or width <= 0 or height <= 0:
        return
    adj = int(round(_round_rect_adjustment(width, height, radius) * 100000))
    prst_geom.rewrite_guides([("adj", adj)])


def _create_polygon_shape(
    slide,
    box: Box,
    clip_path: ClipPath,
):
    """Create a freeform polygon shape from clip-path points.

    Points are percentage values (0-100) relative to the element bounding box.
    """
    left_emu = px_to_emu(box.x)
    top_emu = px_to_emu(box.y)
    w_emu = px_to_emu(box.width)
    h_emu = px_to_emu(box.height)

    points = clip_path.points
    if len(points) < 3:
        logger.warning(
            "Polygon clip-path has fewer than 3 points, falling back to rectangle"
        )
        return None

    # Create a rectangle shape first, then replace its geometry with custom polygon
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Emu(left_emu),
        Emu(top_emu),
        Emu(w_emu),
        Emu(h_emu),
    )

    sp_pr = shape._element.spPr

    # Remove preset geometry
    prst_geom = sp_pr.find(qn("a:prstGeom"))
    if prst_geom is not None:
        sp_pr.remove(prst_geom)

    # Build custom geometry with polygon path
    cust_geom = etree.SubElement(sp_pr, qn("a:custGeom"))
    # Add required empty child elements per OOXML spec
    etree.SubElement(cust_geom, qn("a:avLst"))
    etree.SubElement(cust_geom, qn("a:gdLst"))
    etree.SubElement(cust_geom, qn("a:ahLst"))
    etree.SubElement(cust_geom, qn("a:cxnLst"))
    rect = etree.SubElement(cust_geom, qn("a:rect"))
    rect.set("l", "0")
    rect.set("t", "0")
    rect.set("r", str(w_emu))
    rect.set("b", str(h_emu))

    path_lst = etree.SubElement(cust_geom, qn("a:pathLst"))
    path = etree.SubElement(
        path_lst,
        qn("a:path"),
        attrib={"w": str(w_emu), "h": str(h_emu)},
    )

    # Move to first point
    first = points[0]
    x0 = int(round(first.x / 100.0 * w_emu))
    y0 = int(round(first.y / 100.0 * h_emu))
    move_to = etree.SubElement(path, qn("a:moveTo"))
    etree.SubElement(move_to, qn("a:pt"), attrib={"x": str(x0), "y": str(y0)})

    # Line segments to remaining points
    for pt in points[1:]:
        px = int(round(pt.x / 100.0 * w_emu))
        py = int(round(pt.y / 100.0 * h_emu))
        ln_to = etree.SubElement(path, qn("a:lnTo"))
        etree.SubElement(ln_to, qn("a:pt"), attrib={"x": str(px), "y": str(py)})

    # Close the path
    etree.SubElement(path, qn("a:close"))

    return shape


def _create_background_shape(
    slide,
    box: Box,
    decoration: BoxDecoration,
    shape_type,
):
    """Create the main background shape and apply its fill (solid color or gradient)."""
    # Use polygon shape if clip_path is present
    if decoration.clip_path and decoration.clip_path.type == "polygon":
        bg_shape = _create_polygon_shape(slide, box, decoration.clip_path)
        if bg_shape is not None:
            _remove_theme_style(bg_shape)
            # Apply fill to polygon shape
            if decoration.background_gradient and _set_shape_gradient_fill(
                bg_shape, decoration.background_gradient, slide=slide
            ):
                pass  # gradient applied
            elif decoration.background_color and decoration.opacity > 0:
                fill_color = _with_opacity(
                    decoration.background_color, decoration.opacity
                )
                _set_fill_color(bg_shape.fill, fill_color)
            else:
                bg_shape.fill.background()
            return bg_shape

    left = Emu(px_to_emu(box.x))
    top = Emu(px_to_emu(box.y))
    width = Emu(px_to_emu(box.width))
    height = Emu(px_to_emu(box.height))

    bg_shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    _remove_theme_style(bg_shape)
    if decoration.border_radius_px > 0:
        w_emu = int(px_to_emu(box.width))
        h_emu = int(px_to_emu(box.height))
        r_emu = int(px_to_emu(decoration.border_radius_px))
        _replace_with_round_rect_custgeom(bg_shape, w_emu, h_emu, r_emu)

    if decoration.background_gradient and not _set_shape_gradient_fill(
        bg_shape, decoration.background_gradient, slide=slide
    ):
        fill = bg_shape.fill
        if decoration.background_color and decoration.opacity > 0:
            fill_color = _with_opacity(decoration.background_color, decoration.opacity)
            _set_fill_color(fill, fill_color)
        else:
            fill.background()
    else:
        fill = bg_shape.fill
        if not decoration.background_gradient:
            if decoration.background_color and decoration.opacity > 0:
                fill_color = _with_opacity(
                    decoration.background_color, decoration.opacity
                )
                _set_fill_color(fill, fill_color)
            else:
                fill.background()

    return bg_shape


def _set_shape_gradient_fill(shape, css_gradient: str, slide=None) -> bool:
    """Apply a gradient fill to a shape.

    Tries native OOXML linear gradient first.  When that is not possible
    (e.g. radial-gradient), falls back to rasterising the gradient to a
    PNG and applying it as a picture (blipFill) on the shape.

    Parameters
    ----------
    shape:
        A python-pptx shape object whose ``_element.spPr`` will be modified.
    css_gradient:
        The raw CSS gradient string (e.g. ``"radial-gradient(red, blue)"``).
    slide:
        The python-pptx slide object, required for the PNG fallback path
        so that the image can be registered as a relationship.
    """
    parsed = parse_linear_gradient(css_gradient)
    if parsed is not None:
        sp_pr = shape._element.spPr
        _remove_existing_fills(sp_pr)
        _build_gradient_fill_xml(sp_pr, parsed)
        return True

    # --- PNG rasterization fallback for non-linear gradients ----------
    if slide is None:
        logger.debug(
            "Cannot apply non-linear gradient without slide reference: %s",
            css_gradient,
        )
        return False

    # Derive pixel dimensions from the shape's extent
    xfrm = shape._element.spPr.find(qn("a:xfrm"))
    if xfrm is None:
        return False
    ext = xfrm.find(qn("a:ext"))
    if ext is None:
        return False
    w_emu = int(ext.get("cx", "0"))
    h_emu = int(ext.get("cy", "0"))
    if w_emu <= 0 or h_emu <= 0:
        return False

    width_px = max(int(round(emu_to_px(w_emu))), 1)
    height_px = max(int(round(emu_to_px(h_emu))), 1)

    png_bytes = render_gradient_png(css_gradient, width_px, height_px)
    if png_bytes is None:
        return False

    _apply_picture_fill(shape, slide, png_bytes)
    return True


def _apply_picture_fill(shape, slide, png_bytes: bytes) -> None:
    """Replace a shape's fill with a picture (blipFill) from PNG bytes."""
    # Register the image with the slide's part to get a relationship ID
    slide_part = slide.part
    image_part, r_id = slide_part.get_or_add_image_part(io.BytesIO(png_bytes))

    sp_pr = shape._element.spPr
    _remove_existing_fills(sp_pr)

    blip_fill = etree.SubElement(sp_pr, qn("a:blipFill"))
    blip_fill.set("rotWithShape", "1")
    blip = etree.SubElement(blip_fill, qn("a:blip"))
    blip.set(qn("r:embed"), r_id)
    stretch = etree.SubElement(blip_fill, qn("a:stretch"))
    etree.SubElement(stretch, qn("a:fillRect"))


def _remove_theme_style(shape) -> None:
    """Remove theme style refs that can introduce unwanted effects like shadows."""
    style = shape._element.find(qn("p:style"))
    if style is not None:
        shape._element.remove(style)
