"""Border rendering for decoration shapes."""

from __future__ import annotations

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Emu

from marpx.models import Box, BoxDecoration
from marpx.utils.common import px_to_emu

from .._helpers import _set_fill_color, _set_line_color, _with_opacity
from .shapes import _apply_round_rect_radius, _remove_theme_style


def _apply_uniform_border(
    slide,
    box: Box,
    decoration: BoxDecoration,
    shape_type,
    bg_shape,
    uniform_border,
) -> None:
    """Apply a uniform border to the background shape, or clear the line if none."""
    left = Emu(px_to_emu(box.x))
    top = Emu(px_to_emu(box.y))
    width = Emu(px_to_emu(box.width))
    height = Emu(px_to_emu(box.height))

    if uniform_border:
        border_color = _with_opacity(uniform_border.color, decoration.opacity)
        if border_color.a > 0:
            if hasattr(bg_shape, "line"):
                _set_line_color(bg_shape.line, border_color)
                bg_shape.line.width = Emu(px_to_emu(uniform_border.width_px))
            else:
                border_shape = slide.shapes.add_shape(
                    shape_type, left, top, width, height
                )
                _remove_theme_style(border_shape)
                border_shape.fill.background()
                if decoration.border_radius_px > 0:
                    _apply_round_rect_radius(
                        border_shape,
                        px_to_emu(box.width),
                        px_to_emu(box.height),
                        px_to_emu(decoration.border_radius_px),
                    )
                _set_line_color(border_shape.line, border_color)
                border_shape.line.width = Emu(px_to_emu(uniform_border.width_px))
        else:
            if hasattr(bg_shape, "line"):
                bg_shape.line.fill.background()
    elif hasattr(bg_shape, "line"):
        bg_shape.line.fill.background()


def _apply_left_accent_bar(
    slide,
    box: Box,
    decoration: BoxDecoration,
    *,
    rotation_3d_x_deg: float = 0.0,
    rotation_3d_y_deg: float = 0.0,
    rotation_3d_z_deg: float = 0.0,
    perspective_px: float = 0.0,
    element_height_px: float = 0.0,
) -> None:
    """Render a left accent bar shape for blockquote-style left-border elements."""
    from .core import _apply_scene3d

    accent_border = decoration.border_left
    if accent_border.color is None:
        return
    accent_left, accent_top, accent_width, accent_height = (
        _resolve_left_accent_geometry(box, decoration)
    )
    accent_shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        accent_left,
        accent_top,
        accent_width,
        accent_height,
    )
    _remove_theme_style(accent_shape)
    accent_color = _with_opacity(accent_border.color, decoration.opacity)
    _set_fill_color(accent_shape.fill, accent_color)
    accent_shape.line.fill.background()
    _apply_scene3d(
        accent_shape._element.spPr,
        rotation_3d_x_deg=rotation_3d_x_deg,
        rotation_3d_y_deg=rotation_3d_y_deg,
        rotation_3d_z_deg=rotation_3d_z_deg,
        perspective_px=perspective_px,
        element_height_px=element_height_px,
    )


def _detect_uniform_border(decoration: BoxDecoration):
    """Return a representative border side when all visible sides are equal."""
    sides = [
        decoration.border_top,
        decoration.border_right,
        decoration.border_bottom,
        decoration.border_left,
    ]
    visible = [
        side
        for side in sides
        if side.width_px > 0 and side.style and side.style != "none" and side.color
    ]
    if not visible:
        return None
    first = visible[0]
    if (
        all(
            side.width_px == first.width_px
            and side.style == first.style
            and side.color == first.color
            for side in visible
        )
        and len(visible) == 4
    ):
        return first
    return None


def _is_left_accent_only(decoration: BoxDecoration) -> bool:
    """Return True when only the left border should be rendered as an accent bar."""
    other_sides = [
        decoration.border_top,
        decoration.border_right,
        decoration.border_bottom,
    ]
    return all(
        side.width_px == 0 or side.style == "none" or side.color is None
        for side in other_sides
    )


def _add_side_border_shapes(
    slide,
    box: Box,
    decoration: BoxDecoration,
    *,
    rotation_3d_x_deg: float = 0.0,
    rotation_3d_y_deg: float = 0.0,
    rotation_3d_z_deg: float = 0.0,
    perspective_px: float = 0.0,
    element_height_px: float = 0.0,
) -> None:
    """Render visible non-uniform borders as separate thin rectangle shapes."""
    from .core import _apply_scene3d

    sides = [
        (
            decoration.border_top,
            box.x,
            box.y,
            box.width,
            decoration.border_top.width_px,
        ),
        (
            decoration.border_right,
            box.x + box.width - decoration.border_right.width_px,
            box.y,
            decoration.border_right.width_px,
            box.height,
        ),
        (
            decoration.border_bottom,
            box.x,
            box.y + box.height - decoration.border_bottom.width_px,
            box.width,
            decoration.border_bottom.width_px,
        ),
        (
            decoration.border_left,
            box.x,
            box.y,
            decoration.border_left.width_px,
            box.height,
        ),
    ]

    for border, x, y, width, height in sides:
        if (
            border.width_px <= 0
            or border.style == "none"
            or border.color is None
            or width <= 0
            or height <= 0
        ):
            continue
        border_shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Emu(px_to_emu(x)),
            Emu(px_to_emu(y)),
            Emu(px_to_emu(width)),
            Emu(px_to_emu(height)),
        )
        _remove_theme_style(border_shape)
        _set_fill_color(
            border_shape.fill,
            _with_opacity(border.color, decoration.opacity),
        )
        border_shape.line.fill.background()
        _apply_scene3d(
            border_shape._element.spPr,
            rotation_3d_x_deg=rotation_3d_x_deg,
            rotation_3d_y_deg=rotation_3d_y_deg,
            rotation_3d_z_deg=rotation_3d_z_deg,
            perspective_px=perspective_px,
            element_height_px=element_height_px,
        )


def _resolve_left_accent_geometry(
    box: Box,
    decoration: BoxDecoration,
) -> tuple[Emu, Emu, Emu, Emu]:
    """Match the visible straight segment of a left border on rounded boxes."""
    radius_px = max(decoration.border_radius_px, 0.0)
    vertical_inset_px = min(radius_px, box.height / 2)
    accent_height_px = max(box.height - (vertical_inset_px * 2), 1.0)

    return (
        Emu(px_to_emu(box.x)),
        Emu(px_to_emu(box.y + vertical_inset_px)),
        Emu(px_to_emu(decoration.border_left.width_px)),
        Emu(px_to_emu(accent_height_px)),
    )
