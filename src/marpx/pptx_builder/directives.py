"""Directives rendering (header, footer, page number) for pptx_builder."""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Pt

from marpx.utils.common import px_to_emu


def _add_header(pptx_slide, header_text: str, slide_width_emu: int) -> None:
    """Add header text as a textbox at the top of the slide."""
    margin = Emu(px_to_emu(20))
    height = Emu(px_to_emu(30))
    width = Emu(slide_width_emu) - margin * 2

    txbox = pptx_slide.shapes.add_textbox(margin, margin, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = header_text
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(128, 128, 128)


def _add_footer(
    pptx_slide, footer_text: str, slide_width_emu: int, slide_height_emu: int
) -> None:
    """Add footer text at the bottom of the slide."""
    margin = Emu(px_to_emu(20))
    height = Emu(px_to_emu(25))
    width = Emu(slide_width_emu // 2)
    top = Emu(slide_height_emu) - height - margin

    txbox = pptx_slide.shapes.add_textbox(margin, top, width, height)
    tf = txbox.text_frame
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = footer_text
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(128, 128, 128)


def _add_page_number(
    pptx_slide,
    page_number: int,
    page_total: int | None,
    slide_width_emu: int,
    slide_height_emu: int,
) -> None:
    """Add page number at the bottom-right of the slide."""
    margin = Emu(px_to_emu(20))
    height = Emu(px_to_emu(25))
    width = Emu(px_to_emu(100))
    left = Emu(slide_width_emu) - width - margin
    top = Emu(slide_height_emu) - height - margin

    txbox = pptx_slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    if page_total:
        run.text = f"{page_number} / {page_total}"
    else:
        run.text = str(page_number)
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(128, 128, 128)
