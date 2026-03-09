"""Tests for radial-gradient and PNG fallback in decoration shapes."""

from __future__ import annotations


from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Emu

from marpx.models import Box, BoxDecoration, BoxPadding, BoxShadow, RGBAColor
from marpx.pptx_builder.decoration.shapes import (
    _apply_picture_fill,
    _create_background_shape,
    _set_shape_gradient_fill,
)
from marpx.utils.gradient import render_gradient_png


def _make_slide():
    """Create a real python-pptx slide for testing."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    return slide


def _make_shape(slide, width_emu=914400, height_emu=914400):
    """Add a simple rectangle shape to the slide."""
    return slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Emu(0),
        Emu(0),
        Emu(width_emu),
        Emu(height_emu),
    )


def _default_border():
    return BoxShadow(
        color=RGBAColor(r=0, g=0, b=0, a=0.0),
        offset_x_px=0,
        offset_y_px=0,
        blur_radius_px=0,
        spread_px=0,
    )


def _default_decoration(**kwargs):
    defaults = {
        "background_color": RGBAColor(r=255, g=255, b=255, a=1.0),
        "background_gradient": None,
        "opacity": 1.0,
        "border_radius_px": 0,
        "border_top": BoxDecoration.__dataclass_fields__["border_top"].default
        if hasattr(BoxDecoration, "__dataclass_fields__")
        else None,
        "clip_path": None,
    }
    defaults.update(kwargs)
    return defaults


# ---------------------------------------------------------------------------
# TestSetShapeGradientFill - linear gradient (existing behavior)
# ---------------------------------------------------------------------------
class TestSetShapeGradientFillLinear:
    """Existing linear gradient path should still work."""

    def test_linear_gradient_applies_grad_fill(self) -> None:
        slide = _make_slide()
        shape = _make_shape(slide)
        result = _set_shape_gradient_fill(
            shape, "linear-gradient(red, blue)", slide=slide
        )
        assert result is True
        sp_pr = shape._element.spPr
        grad_fill = sp_pr.find(qn("a:gradFill"))
        assert grad_fill is not None

    def test_linear_gradient_without_slide_still_works(self) -> None:
        slide = _make_slide()
        shape = _make_shape(slide)
        result = _set_shape_gradient_fill(
            shape, "linear-gradient(red, blue)", slide=None
        )
        assert result is True

    def test_invalid_gradient_returns_false(self) -> None:
        slide = _make_slide()
        shape = _make_shape(slide)
        result = _set_shape_gradient_fill(shape, "not-a-gradient", slide=slide)
        assert result is False


# ---------------------------------------------------------------------------
# TestSetShapeGradientFill - radial gradient (PNG fallback)
# ---------------------------------------------------------------------------
class TestSetShapeGradientFillRadial:
    """Radial gradient should use PNG rasterization fallback."""

    def test_radial_gradient_applies_blip_fill(self) -> None:
        slide = _make_slide()
        shape = _make_shape(slide)
        result = _set_shape_gradient_fill(
            shape, "radial-gradient(red, blue)", slide=slide
        )
        assert result is True
        sp_pr = shape._element.spPr
        blip_fill = sp_pr.find(qn("a:blipFill"))
        assert blip_fill is not None
        blip = blip_fill.find(qn("a:blip"))
        assert blip is not None
        assert qn("r:embed") in blip.attrib

    def test_radial_gradient_without_slide_returns_false(self) -> None:
        slide = _make_slide()
        shape = _make_shape(slide)
        result = _set_shape_gradient_fill(
            shape, "radial-gradient(red, blue)", slide=None
        )
        assert result is False

    def test_radial_gradient_has_stretch_fill_rect(self) -> None:
        slide = _make_slide()
        shape = _make_shape(slide)
        _set_shape_gradient_fill(shape, "radial-gradient(red, blue)", slide=slide)
        sp_pr = shape._element.spPr
        blip_fill = sp_pr.find(qn("a:blipFill"))
        stretch = blip_fill.find(qn("a:stretch"))
        assert stretch is not None
        fill_rect = stretch.find(qn("a:fillRect"))
        assert fill_rect is not None

    def test_radial_gradient_removes_previous_fills(self) -> None:
        slide = _make_slide()
        shape = _make_shape(slide)
        # Apply solid fill first
        shape.fill.solid()
        # Now apply radial gradient
        _set_shape_gradient_fill(shape, "radial-gradient(red, blue)", slide=slide)
        sp_pr = shape._element.spPr
        solid_fills = sp_pr.findall(qn("a:solidFill"))
        assert len(solid_fills) == 0

    def test_radial_gradient_with_center(self) -> None:
        slide = _make_slide()
        shape = _make_shape(slide)
        result = _set_shape_gradient_fill(
            shape,
            "radial-gradient(circle at 25% 75%, red, blue)",
            slide=slide,
        )
        assert result is True

    def test_radial_gradient_three_stops(self) -> None:
        slide = _make_slide()
        shape = _make_shape(slide)
        result = _set_shape_gradient_fill(
            shape,
            "radial-gradient(red, green, blue)",
            slide=slide,
        )
        assert result is True


# ---------------------------------------------------------------------------
# TestApplyPictureFill
# ---------------------------------------------------------------------------
class TestApplyPictureFill:
    """Tests for _apply_picture_fill helper."""

    def test_creates_blip_fill_xml(self) -> None:
        slide = _make_slide()
        shape = _make_shape(slide)
        png_bytes = render_gradient_png("radial-gradient(red, blue)", 100, 100)
        assert png_bytes is not None
        _apply_picture_fill(shape, slide, png_bytes, 914400, 914400)
        sp_pr = shape._element.spPr
        blip_fill = sp_pr.find(qn("a:blipFill"))
        assert blip_fill is not None
        assert blip_fill.get("rotWithShape") == "1"


# ---------------------------------------------------------------------------
# TestCreateBackgroundShapeWithRadialGradient
# ---------------------------------------------------------------------------
class TestCreateBackgroundShapeWithRadialGradient:
    """Integration: _create_background_shape handles radial-gradient."""

    def _make_decoration(self, gradient: str) -> BoxDecoration:
        from marpx.models import BorderSide

        border = BorderSide(width_px=0, style="none", color=None)
        padding = BoxPadding(top_px=0, right_px=0, bottom_px=0, left_px=0)
        return BoxDecoration(
            background_color=RGBAColor(r=255, g=255, b=255, a=1.0),
            background_gradient=gradient,
            opacity=1.0,
            border_radius_px=0,
            border_top=border,
            border_right=border,
            border_bottom=border,
            border_left=border,
            padding=padding,
            box_shadows=[],
            clip_path=None,
        )

    def test_radial_gradient_decoration_shape(self) -> None:
        slide = _make_slide()
        box = Box(x=0, y=0, width=100, height=100)
        decoration = self._make_decoration("radial-gradient(red, blue)")
        shape = _create_background_shape(
            slide,
            box,
            decoration,
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        )
        assert shape is not None
        sp_pr = shape._element.spPr
        blip_fill = sp_pr.find(qn("a:blipFill"))
        assert blip_fill is not None

    def test_linear_gradient_decoration_shape_still_uses_grad_fill(self) -> None:
        slide = _make_slide()
        box = Box(x=0, y=0, width=100, height=100)
        decoration = self._make_decoration("linear-gradient(red, blue)")
        shape = _create_background_shape(
            slide,
            box,
            decoration,
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        )
        assert shape is not None
        sp_pr = shape._element.spPr
        grad_fill = sp_pr.find(qn("a:gradFill"))
        assert grad_fill is not None
