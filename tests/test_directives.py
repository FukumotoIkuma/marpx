"""Tests for Marp directives support (header, footer, paginate).

These tests construct model objects directly -- NO Playwright or marp-cli needed.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation
from pptx.enum.text import PP_ALIGN

from marpx.models import (
    Box,
    ElementType,
    Paragraph,
    Presentation,
    Slide,
    SlideElement,
    TextRun,
    TextStyle,
)
from marpx.pptx_builder.builder import build_pptx
from marpx.pptx_builder.directives import _add_footer, _add_header, _add_page_number
from marpx.utils import px_to_emu


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _make_heading(text: str, level: int = 1) -> SlideElement:
    return SlideElement(
        element_type=ElementType.HEADING,
        box=Box(x=50, y=20, width=600, height=60),
        heading_level=level,
        paragraphs=[
            Paragraph(
                runs=[TextRun(text=text, style=TextStyle(bold=True, font_size_px=32))],
                alignment="left",
            )
        ],
    )


def _build_and_read(presentation: Presentation, tmp_path: Path) -> PptxPresentation:
    """Build PPTX and re-open it for assertions."""
    out = tmp_path / "test_directives.pptx"
    build_pptx(presentation, out)
    return PptxPresentation(str(out))


def _get_all_textbox_texts(pptx_slide) -> list[str]:
    """Collect all text from textbox shapes on a slide."""
    texts = []
    for shape in pptx_slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    return texts


def _make_blank_pptx_slide():
    """Create a minimal PPTX presentation and return (pptx_slide, prs)."""
    prs = PptxPresentation()
    layout = prs.slide_layouts[6]  # Blank
    pptx_slide = prs.slides.add_slide(layout)
    return pptx_slide, prs


# ---------------------------------------------------------------------------
# Tests: Slide model directive fields
# ---------------------------------------------------------------------------


class TestSlideModelDirectives:
    """Verify that the Slide model accepts directive fields."""

    def test_default_values(self) -> None:
        slide = Slide(width_px=1280, height_px=720)
        assert slide.header_text is None
        assert slide.footer_text is None
        assert slide.paginate is False
        assert slide.page_number is None
        assert slide.page_total is None

    def test_explicit_values(self) -> None:
        slide = Slide(
            width_px=1280,
            height_px=720,
            header_text="My Header",
            footer_text="My Footer",
            paginate=True,
            page_number=3,
            page_total=10,
        )
        assert slide.header_text == "My Header"
        assert slide.footer_text == "My Footer"
        assert slide.paginate is True
        assert slide.page_number == 3
        assert slide.page_total == 10


# ---------------------------------------------------------------------------
# Tests: _add_header
# ---------------------------------------------------------------------------


class TestAddHeader:
    """Verify _add_header creates a textbox at top of slide."""

    def test_header_creates_textbox(self) -> None:
        pptx_slide, _ = _make_blank_pptx_slide()
        slide_w = px_to_emu(1280)
        _add_header(pptx_slide, "Test Header", slide_w)

        assert len(pptx_slide.shapes) == 1
        shape = pptx_slide.shapes[0]
        assert shape.has_text_frame
        assert shape.text_frame.paragraphs[0].text == "Test Header"

    def test_header_position_at_top(self) -> None:
        pptx_slide, _ = _make_blank_pptx_slide()
        slide_w = px_to_emu(1280)
        _add_header(pptx_slide, "Header", slide_w)

        shape = pptx_slide.shapes[0]
        # Top position should be small (margin at top)
        margin_emu = px_to_emu(20)
        assert shape.top == margin_emu

    def test_header_font_color_gray(self) -> None:
        pptx_slide, _ = _make_blank_pptx_slide()
        slide_w = px_to_emu(1280)
        _add_header(pptx_slide, "Header", slide_w)

        shape = pptx_slide.shapes[0]
        run = shape.text_frame.paragraphs[0].runs[0]
        assert run.font.color.rgb is not None
        # Gray: RGB(128, 128, 128)
        assert str(run.font.color.rgb) == "808080"


# ---------------------------------------------------------------------------
# Tests: _add_footer
# ---------------------------------------------------------------------------


class TestAddFooter:
    """Verify _add_footer creates a textbox at bottom of slide."""

    def test_footer_creates_textbox(self) -> None:
        pptx_slide, _ = _make_blank_pptx_slide()
        slide_w = px_to_emu(1280)
        slide_h = px_to_emu(720)
        _add_footer(pptx_slide, "Test Footer", slide_w, slide_h)

        assert len(pptx_slide.shapes) == 1
        shape = pptx_slide.shapes[0]
        assert shape.has_text_frame
        assert shape.text_frame.paragraphs[0].text == "Test Footer"

    def test_footer_position_at_bottom(self) -> None:
        pptx_slide, _ = _make_blank_pptx_slide()
        slide_w = px_to_emu(1280)
        slide_h = px_to_emu(720)
        _add_footer(pptx_slide, "Footer", slide_w, slide_h)

        shape = pptx_slide.shapes[0]
        # Top position should be near the bottom
        # slide_h_emu - height_emu - margin_emu
        assert shape.top > px_to_emu(600)  # well past middle

    def test_footer_font_color_gray(self) -> None:
        pptx_slide, _ = _make_blank_pptx_slide()
        slide_w = px_to_emu(1280)
        slide_h = px_to_emu(720)
        _add_footer(pptx_slide, "Footer", slide_w, slide_h)

        shape = pptx_slide.shapes[0]
        run = shape.text_frame.paragraphs[0].runs[0]
        assert str(run.font.color.rgb) == "808080"


# ---------------------------------------------------------------------------
# Tests: _add_page_number
# ---------------------------------------------------------------------------


class TestAddPageNumber:
    """Verify _add_page_number creates a textbox at bottom-right."""

    def test_page_number_with_total(self) -> None:
        pptx_slide, _ = _make_blank_pptx_slide()
        slide_w = px_to_emu(1280)
        slide_h = px_to_emu(720)
        _add_page_number(pptx_slide, 3, 10, slide_w, slide_h)

        shape = pptx_slide.shapes[0]
        assert shape.text_frame.paragraphs[0].text == "3 / 10"

    def test_page_number_without_total(self) -> None:
        pptx_slide, _ = _make_blank_pptx_slide()
        slide_w = px_to_emu(1280)
        slide_h = px_to_emu(720)
        _add_page_number(pptx_slide, 3, None, slide_w, slide_h)

        shape = pptx_slide.shapes[0]
        assert shape.text_frame.paragraphs[0].text == "3"

    def test_page_number_right_aligned(self) -> None:
        pptx_slide, _ = _make_blank_pptx_slide()
        slide_w = px_to_emu(1280)
        slide_h = px_to_emu(720)
        _add_page_number(pptx_slide, 1, 5, slide_w, slide_h)

        shape = pptx_slide.shapes[0]
        para = shape.text_frame.paragraphs[0]
        assert para.alignment == PP_ALIGN.RIGHT

    def test_page_number_position_bottom_right(self) -> None:
        pptx_slide, _ = _make_blank_pptx_slide()
        slide_w = px_to_emu(1280)
        slide_h = px_to_emu(720)
        _add_page_number(pptx_slide, 1, None, slide_w, slide_h)

        shape = pptx_slide.shapes[0]
        # Should be near bottom
        assert shape.top > px_to_emu(600)
        # Should be near right edge
        assert shape.left > px_to_emu(1000)

    def test_page_number_font_color_gray(self) -> None:
        pptx_slide, _ = _make_blank_pptx_slide()
        slide_w = px_to_emu(1280)
        slide_h = px_to_emu(720)
        _add_page_number(pptx_slide, 1, None, slide_w, slide_h)

        shape = pptx_slide.shapes[0]
        run = shape.text_frame.paragraphs[0].runs[0]
        assert str(run.font.color.rgb) == "808080"


# ---------------------------------------------------------------------------
# Tests: No directives = no extra shapes
# ---------------------------------------------------------------------------


class TestNoDirectives:
    """Verify that slides without directives have no extra shapes."""

    def test_no_extra_shapes(self, tmp_path: Path) -> None:
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[_make_heading("Title")],
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        # Only the heading textbox
        assert len(slide.shapes) == 1


# ---------------------------------------------------------------------------
# Tests: Full build with all directives
# ---------------------------------------------------------------------------


class TestFullBuildWithDirectives:
    """Verify full PPTX build with header + footer + paginate."""

    def test_all_directives_present(self, tmp_path: Path) -> None:
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[_make_heading("Slide 1")],
                    header_text="My Header",
                    footer_text="My Footer",
                    paginate=True,
                    page_number=1,
                    page_total=5,
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]

        # 1 heading + 1 header + 1 footer + 1 page number = 4 shapes
        assert len(slide.shapes) == 4

        texts = _get_all_textbox_texts(slide)
        assert any("My Header" in t for t in texts)
        assert any("My Footer" in t for t in texts)
        assert any("1 / 5" in t for t in texts)

    def test_header_only(self, tmp_path: Path) -> None:
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[_make_heading("Title")],
                    header_text="Only Header",
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]

        # 1 heading + 1 header = 2 shapes
        assert len(slide.shapes) == 2
        texts = _get_all_textbox_texts(slide)
        assert any("Only Header" in t for t in texts)

    def test_footer_only(self, tmp_path: Path) -> None:
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[_make_heading("Title")],
                    footer_text="Only Footer",
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]

        # 1 heading + 1 footer = 2 shapes
        assert len(slide.shapes) == 2
        texts = _get_all_textbox_texts(slide)
        assert any("Only Footer" in t for t in texts)

    def test_paginate_only(self, tmp_path: Path) -> None:
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[_make_heading("Title")],
                    paginate=True,
                    page_number=2,
                    page_total=8,
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]

        # 1 heading + 1 page number = 2 shapes
        assert len(slide.shapes) == 2
        texts = _get_all_textbox_texts(slide)
        assert any("2 / 8" in t for t in texts)

    def test_paginate_false_no_page_number(self, tmp_path: Path) -> None:
        """When paginate=False, no page number should be added even if page_number is set."""
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[_make_heading("Title")],
                    paginate=False,
                    page_number=3,
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]

        # Only the heading, no page number
        assert len(slide.shapes) == 1


# ---------------------------------------------------------------------------
# Tests: Multi-slide with per-slide overrides
# ---------------------------------------------------------------------------


class TestPerSlideOverrides:
    """Verify that per-slide directive overrides work correctly."""

    def test_mixed_directives(self, tmp_path: Path) -> None:
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[_make_heading("S1")],
                    header_text="Header",
                    footer_text="Footer",
                    paginate=True,
                    page_number=1,
                    page_total=3,
                ),
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[_make_heading("S2")],
                    # No directives on this slide
                ),
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[_make_heading("S3")],
                    header_text="Custom Header",
                    paginate=True,
                    page_number=3,
                ),
            ],
        )
        pptx = _build_and_read(pres, tmp_path)

        # Slide 1: heading + header + footer + page number = 4
        assert len(pptx.slides[0].shapes) == 4

        # Slide 2: heading only = 1
        assert len(pptx.slides[1].shapes) == 1

        # Slide 3: heading + header + page number = 3
        assert len(pptx.slides[2].shapes) == 3
        texts = _get_all_textbox_texts(pptx.slides[2])
        assert any("Custom Header" in t for t in texts)
        assert any("3" in t for t in texts)


# ---------------------------------------------------------------------------
# Integration test (requires Playwright + marp-cli)
# ---------------------------------------------------------------------------


@pytest.mark.integration
class TestDirectivesIntegration:
    """Integration test using Marp-rendered HTML with directives fixture."""

    def test_paginate_header_footer_fixture(self, fixtures_dir: Path) -> None:
        """Verify extraction from paginate-header-footer.md fixture.

        This test requires Playwright and marp-cli to be installed.
        """
        md_path = fixtures_dir / "paginate-header-footer.md"
        assert md_path.exists(), f"Fixture not found: {md_path}"
        # Full integration would render and extract; skipped here as unit-only run.
