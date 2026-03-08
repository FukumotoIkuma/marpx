"""End-to-end tests for the marpx conversion pipeline.

All tests require Playwright and marp-cli, so they are marked as integration tests.
"""

from __future__ import annotations

from pathlib import Path

import pytest

from pptx import Presentation as PptxPresentation

from marpx.converter import ConversionError, convert

FIXTURES_DIR = Path(__file__).parent / "fixtures"


@pytest.fixture(scope="module")
def simple_pptx(session_output_dir: Path) -> Path:
    output = session_output_dir / "simple.e2e.pptx"
    convert(FIXTURES_DIR / "simple.md", output)
    return output


@pytest.fixture(scope="module")
def table_pptx(session_output_dir: Path) -> Path:
    output = session_output_dir / "table.e2e.pptx"
    convert(FIXTURES_DIR / "table.md", output)
    return output


@pytest.fixture(scope="module")
def complex_pptx(session_output_dir: Path) -> Path:
    output = session_output_dir / "complex.e2e.pptx"
    convert(FIXTURES_DIR / "complex-unsupported.md", output)
    return output


@pytest.fixture(scope="module")
def rendered_layout_boxes_pptx(session_output_dir: Path) -> Path:
    output = session_output_dir / "rendered_layout_boxes.e2e.pptx"
    convert(FIXTURES_DIR / "rendered-layout-boxes.md", output)
    return output


@pytest.fixture(scope="module")
def quote_box_list_pptx(session_output_dir: Path) -> Path:
    output = session_output_dir / "quote_box_list.e2e.pptx"
    convert(FIXTURES_DIR / "quote-box-list.md", output)
    return output


@pytest.fixture(scope="module")
def decorated_raw_text_pptx(session_output_dir: Path) -> Path:
    output = session_output_dir / "decorated_raw_text.e2e.pptx"
    convert(FIXTURES_DIR / "decorated-raw-text.md", output)
    return output


@pytest.fixture(scope="module")
def rendered_layout_images_pptx(session_output_dir: Path) -> Path:
    output = session_output_dir / "rendered_layout_images.e2e.pptx"
    convert(FIXTURES_DIR / "rendered-layout-images.md", output)
    return output


@pytest.fixture(scope="module")
def multi_paragraph_pptx(session_output_dir: Path) -> Path:
    output = session_output_dir / "multi_paragraph.e2e.pptx"
    convert(FIXTURES_DIR / "multi-paragraph.md", output)
    return output


@pytest.mark.integration
class TestSimpleConversion:
    """End-to-end tests with simple.md."""

    def test_produces_pptx(self, simple_pptx: Path) -> None:
        assert simple_pptx.exists()
        assert simple_pptx.stat().st_size > 0

    def test_correct_slide_count(self, simple_pptx: Path) -> None:
        pptx = PptxPresentation(str(simple_pptx))
        assert len(pptx.slides) == 2

    def test_has_shapes(self, simple_pptx: Path) -> None:
        pptx = PptxPresentation(str(simple_pptx))
        # First slide should have shapes (heading, paragraph, list)
        assert len(pptx.slides[0].shapes) >= 1


@pytest.mark.integration
class TestTableConversion:
    """End-to-end tests with table.md."""

    def test_produces_pptx(self, table_pptx: Path) -> None:
        assert table_pptx.exists()

    def test_has_table(self, table_pptx: Path) -> None:
        pptx = PptxPresentation(str(table_pptx))
        # Look for at least one table shape across all slides
        has_table = False
        for slide in pptx.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    has_table = True
                    break
        assert has_table


@pytest.mark.integration
class TestComplexUnsupported:
    """End-to-end tests with complex-unsupported.md."""

    def test_does_not_crash(self, complex_pptx: Path) -> None:
        assert complex_pptx.exists()
        assert complex_pptx.stat().st_size > 0

    def test_slide_count(self, complex_pptx: Path) -> None:
        pptx = PptxPresentation(str(complex_pptx))
        assert len(pptx.slides) == 3


@pytest.mark.integration
class TestRenderedLayoutConversion:
    """End-to-end tests for rendered-layout capture."""

    def test_decorated_blocks_render_as_multiple_shapes(
        self, rendered_layout_boxes_pptx: Path
    ) -> None:
        pptx = PptxPresentation(str(rendered_layout_boxes_pptx))
        slide = pptx.slides[0]
        assert len(slide.shapes) >= 5

    def test_checklist_marker_is_present(
        self, rendered_layout_boxes_pptx: Path
    ) -> None:
        pptx = PptxPresentation(str(rendered_layout_boxes_pptx))
        texts = [
            para.text
            for shape in pptx.slides[0].shapes
            if shape.has_text_frame
            for para in shape.text_frame.paragraphs
        ]
        assert any("☐ First item" in text for text in texts)

    def test_quote_box_with_list_stays_decorated(
        self, rendered_layout_boxes_pptx: Path
    ) -> None:
        pptx = PptxPresentation(str(rendered_layout_boxes_pptx))
        text_shapes = [
            shape
            for shape in pptx.slides[0].shapes
            if shape.has_text_frame and len(shape.text_frame.paragraphs) >= 4
        ]
        assert text_shapes
        assert any(shape.shape_type != 17 for shape in pptx.slides[0].shapes)

    def test_quote_box_markdown_list_renders_in_decorated_shape(
        self, quote_box_list_pptx: Path
    ) -> None:
        pptx = PptxPresentation(str(quote_box_list_pptx))
        text_shapes = [
            shape
            for shape in pptx.slides[0].shapes
            if shape.has_text_frame and "sample text" in shape.text
        ]
        assert len(text_shapes) == 1
        quote_shape = text_shapes[0]
        assert any(shape.shape_type != 17 for shape in pptx.slides[0].shapes)
        paragraphs = quote_shape.text_frame.paragraphs
        assert len(paragraphs) == 4
        assert paragraphs[0].runs[0].font.bold is True
        assert paragraphs[1].level == 0
        assert paragraphs[-1].runs[0].font.bold is True

    def test_decorated_raw_text_container_renders_all_text(
        self, decorated_raw_text_pptx: Path
    ) -> None:
        pptx = PptxPresentation(str(decorated_raw_text_pptx))
        text_shapes = [
            shape
            for shape in pptx.slides[0].shapes
            if shape.has_text_frame and "neutral placeholder prose" in shape.text
        ]
        assert len(text_shapes) == 1
        quote_shape = text_shapes[0]
        assert any(shape.shape_type != 17 for shape in pptx.slides[0].shapes)
        assert len(quote_shape.text_frame.paragraphs) == 2
        assert "Example" in quote_shape.text_frame.paragraphs[1].text

    def test_object_fit_images_are_not_stretched(
        self, rendered_layout_images_pptx: Path
    ) -> None:
        pptx = PptxPresentation(str(rendered_layout_images_pptx))
        picture_shapes = [
            shape for shape in pptx.slides[0].shapes if getattr(shape, "image", None)
        ]
        assert len(picture_shapes) == 2
        assert all(shape.width == shape.height for shape in picture_shapes)


@pytest.mark.integration
class TestParagraphGroupingConversion:
    """End-to-end tests for keeping multiple paragraphs in one textbox."""

    def test_multi_paragraph_stays_in_one_text_shape(
        self, multi_paragraph_pptx: Path
    ) -> None:
        pptx = PptxPresentation(str(multi_paragraph_pptx))
        text_shapes = [
            shape
            for shape in pptx.slides[0].shapes
            if shape.has_text_frame and "plain placeholder sentence" in shape.text
        ]
        assert len(text_shapes) == 1
        assert len(text_shapes[0].text_frame.paragraphs) >= 2


@pytest.mark.integration
class TestErrorHandling:
    """Edge-case error handling."""

    def test_nonexistent_input(self, tmp_output_dir: Path) -> None:
        with pytest.raises(ConversionError, match="not found"):
            convert("/nonexistent/file.md", tmp_output_dir / "out.pptx")
