"""Tests for marpx.pptx_builder module.

These tests construct model objects directly -- NO Playwright or marp-cli needed.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from marpx.models import (
    Background,
    BorderSide,
    Box,
    BoxDecoration,
    BoxShadow,
    BoxPadding,
    ElementType,
    ListItem,
    Paragraph,
    Presentation,
    Point,
    RGBAColor,
    Slide,
    SlideElement,
    TableCell,
    TableRow,
    TextRun,
    TextStyle,
    UnsupportedInfo,
)
from marpx.pptx_builder.builder import build_pptx
from marpx.pptx_builder.image import MissingDependencyError, _resolve_image_placement
from marpx.pptx_builder.scene3d import fit_scene3d_rotations
from marpx.utils import px_to_emu


# ---------------------------------------------------------------------------
# Helpers to build test data
# ---------------------------------------------------------------------------


def _make_heading(text: str, level: int = 1) -> SlideElement:
    return SlideElement(
        element_type=ElementType.HEADING,
        box=Box(x=50, y=20, width=600, height=60),
        heading_level=level,
        paragraphs=[
            Paragraph(
                runs=[TextRun(text=text, style=TextStyle(bold=True, font_size_px=32))],
                alignment="left",
            )
        ],
    )


def _make_paragraph(text: str) -> SlideElement:
    return SlideElement(
        element_type=ElementType.PARAGRAPH,
        box=Box(x=50, y=100, width=600, height=40),
        paragraphs=[Paragraph(runs=[TextRun(text=text)], alignment="left")],
    )


def _make_list(items: list[tuple[str, int]], ordered: bool = False) -> SlideElement:
    etype = ElementType.ORDERED_LIST if ordered else ElementType.UNORDERED_LIST
    return SlideElement(
        element_type=etype,
        box=Box(x=50, y=160, width=600, height=200),
        list_items=[
            ListItem(
                runs=[TextRun(text=text)],
                level=level,
                order_number=(i + 1) if ordered else None,
            )
            for i, (text, level) in enumerate(items)
        ],
    )


def _make_table(rows_data: list[list[str]], has_header: bool = True) -> SlideElement:
    table_rows = []
    for r_idx, row in enumerate(rows_data):
        cells = [
            TableCell(
                paragraphs=[Paragraph(runs=[TextRun(text=cell_text)])],
                is_header=(r_idx == 0 and has_header),
            )
            for cell_text in row
        ]
        table_rows.append(TableRow(cells=cells))
    return SlideElement(
        element_type=ElementType.TABLE,
        box=Box(x=50, y=100, width=600, height=300),
        table_rows=table_rows,
    )


def _make_code_block(
    code: str,
    language: str = "python",
    decoration: BoxDecoration | None = None,
) -> SlideElement:
    box = Box(x=50, y=100, width=600, height=150)
    return SlideElement(
        element_type=ElementType.CODE_BLOCK,
        box=box,
        content_box=_content_box_from_decoration(box, decoration)
        if decoration is not None
        else None,
        paragraphs=[
            Paragraph(
                runs=[TextRun(text=code, style=TextStyle(font_family="Courier New"))]
            )
        ],
        code_language=language,
        code_background=RGBAColor(r=40, g=42, b=54),
        decoration=decoration,
    )


def _content_box_from_decoration(box: Box, decoration: BoxDecoration) -> Box:
    left_inset = decoration.border_left.width_px + decoration.padding.left_px
    top_inset = decoration.border_top.width_px + decoration.padding.top_px
    right_inset = decoration.border_right.width_px + decoration.padding.right_px
    bottom_inset = decoration.border_bottom.width_px + decoration.padding.bottom_px
    return Box(
        x=box.x + left_inset,
        y=box.y + top_inset,
        width=max(box.width - left_inset - right_inset, 1),
        height=max(box.height - top_inset - bottom_inset, 1),
    )


def _make_decorated_block(text: str) -> SlideElement:
    box = Box(x=50, y=100, width=400, height=140)
    decoration = BoxDecoration(
        background_color=RGBAColor(r=248, g=244, b=255),
        border_left=BorderSide(
            width_px=5,
            style="solid",
            color=RGBAColor(r=74, g=58, b=138),
        ),
        padding=BoxPadding(top_px=12, right_px=16, bottom_px=12, left_px=16),
        border_radius_px=4,
    )
    return SlideElement(
        element_type=ElementType.DECORATED_BLOCK,
        box=box,
        content_box=_content_box_from_decoration(box, decoration),
        paragraphs=[Paragraph(runs=[TextRun(text=text)])],
        decoration=decoration,
    )


def _make_decorated_list_block() -> SlideElement:
    box = Box(x=50, y=100, width=400, height=180)
    decoration = BoxDecoration(
        background_color=RGBAColor(r=248, g=244, b=255),
        border_left=BorderSide(
            width_px=5,
            style="solid",
            color=RGBAColor(r=74, g=58, b=138),
        ),
        padding=BoxPadding(top_px=12, right_px=16, bottom_px=12, left_px=16),
        border_radius_px=4,
    )
    return SlideElement(
        element_type=ElementType.DECORATED_BLOCK,
        box=box,
        content_box=_content_box_from_decoration(box, decoration),
        paragraphs=[
            Paragraph(runs=[TextRun(text="Lead")]),
            Paragraph(runs=[TextRun(text="First")], list_level=0),
            Paragraph(runs=[TextRun(text="Second")], list_level=1),
        ],
        decoration=decoration,
    )


def _build_and_read(presentation: Presentation, tmp_path: Path) -> PptxPresentation:
    """Build PPTX and re-open it for assertions."""
    out = tmp_path / "test_output.pptx"
    build_pptx(presentation, out)
    return PptxPresentation(str(out))


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------


class TestBuildSlideCount:
    """Verify correct number of slides."""

    def test_single_slide(self, tmp_path: Path) -> None:
        pres = Presentation(
            slides=[
                Slide(width_px=1280, height_px=720, elements=[_make_heading("Slide 1")])
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        assert len(pptx.slides) == 1

    def test_multiple_slides(self, tmp_path: Path) -> None:
        pres = Presentation(
            slides=[
                Slide(width_px=1280, height_px=720, elements=[_make_heading("S1")]),
                Slide(width_px=1280, height_px=720, elements=[_make_heading("S2")]),
                Slide(width_px=1280, height_px=720, elements=[_make_heading("S3")]),
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        assert len(pptx.slides) == 3


class TestSlideDimensions:
    """Verify slide width and height."""

    def test_default_dimensions(self, tmp_path: Path) -> None:
        pres = Presentation(
            slides=[Slide(width_px=1280, height_px=720, elements=[])],
        )
        pptx = _build_and_read(pres, tmp_path)
        assert pptx.slide_width == px_to_emu(1280)
        assert pptx.slide_height == px_to_emu(720)

    def test_mixed_slide_sizes_raise(self, tmp_path: Path) -> None:
        pres = Presentation(
            slides=[
                Slide(width_px=1280, height_px=720, elements=[]),
                Slide(width_px=1024, height_px=768, elements=[]),
            ],
        )
        out = tmp_path / "mixed_sizes.pptx"
        with pytest.raises(ValueError, match="mixed slide sizes"):
            build_pptx(pres, out)


class TestTextboxContent:
    """Verify text content in textboxes."""

    def test_heading_text(self, tmp_path: Path) -> None:
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[_make_heading("My Title")],
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    texts.append(para.text)
        assert any("My Title" in t for t in texts)

    def test_paragraph_text(self, tmp_path: Path) -> None:
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[_make_paragraph("Hello World")],
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    texts.append(para.text)
        assert any("Hello World" in t for t in texts)

    def test_strikethrough_run_is_emitted(self, tmp_path: Path) -> None:
        element = SlideElement(
            element_type=ElementType.PARAGRAPH,
            box=Box(x=50, y=100, width=600, height=40),
            paragraphs=[
                Paragraph(
                    runs=[
                        TextRun(
                            text="struck",
                            style=TextStyle(strike=True),
                        )
                    ]
                )
            ],
        )
        pres = Presentation(
            slides=[Slide(width_px=1280, height_px=720, elements=[element])]
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        textbox = next(shape for shape in slide.shapes if shape.has_text_frame)
        run = textbox.text_frame.paragraphs[0].runs[0]
        assert 'strike="sngStrike"' in run._r.xml

    def test_gradient_text_run_emits_grad_fill(self, tmp_path: Path) -> None:
        element = SlideElement(
            element_type=ElementType.PARAGRAPH,
            box=Box(x=50, y=100, width=600, height=80),
            paragraphs=[
                Paragraph(
                    runs=[
                        TextRun(
                            text="2.4M",
                            style=TextStyle(
                                bold=True,
                                font_size_px=56,
                                text_gradient="linear-gradient(135deg, #60a5fa, #c084fc)",
                            ),
                        )
                    ]
                )
            ],
        )
        pres = Presentation(
            slides=[Slide(width_px=1280, height_px=720, elements=[element])]
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        textbox = next(shape for shape in slide.shapes if shape.has_text_frame)
        run = textbox.text_frame.paragraphs[0].runs[0]

        assert "a:gradFill" in run._r.xml
        assert 'val="60A5FA"' in run._r.xml
        assert 'val="C084FC"' in run._r.xml
        assert "<a:tileRect/>" in run._r.xml
        assert run._r.xml.index("a:gradFill") < run._r.xml.index("a:latin")

    def test_middle_vertical_align_sets_center_anchor(self, tmp_path: Path) -> None:
        element = SlideElement(
            element_type=ElementType.DECORATED_BLOCK,
            box=Box(x=50, y=100, width=32, height=32),
            content_box=Box(x=50, y=100, width=32, height=32),
            vertical_align="middle",
            decoration=BoxDecoration(
                background_color=RGBAColor(r=59, g=130, b=246),
                border_radius_px=16,
            ),
            paragraphs=[
                Paragraph(
                    runs=[
                        TextRun(
                            text="1",
                            style=TextStyle(
                                bold=True,
                                color=RGBAColor(r=255, g=255, b=255),
                            ),
                        )
                    ],
                    alignment="center",
                )
            ],
        )
        pres = Presentation(
            slides=[Slide(width_px=1280, height_px=720, elements=[element])]
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        textbox = next(
            shape
            for shape in slide.shapes
            if shape.has_text_frame and shape.text == "1"
        )
        assert textbox.text_frame.vertical_anchor == MSO_VERTICAL_ANCHOR.MIDDLE
        assert "<a:lnSpc>" not in textbox.text_frame.paragraphs[0]._p.xml

    def test_inline_code_run_sets_highlight_color(self, tmp_path: Path) -> None:
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[
                        SlideElement(
                            element_type=ElementType.PARAGRAPH,
                            box=Box(x=50, y=100, width=600, height=40),
                            paragraphs=[
                                Paragraph(
                                    runs=[
                                        TextRun(text="This uses "),
                                        TextRun(
                                            text="inline code",
                                            style=TextStyle(
                                                font_family="Courier New",
                                                background_color=RGBAColor(
                                                    r=248, g=250, b=252
                                                ),
                                            ),
                                        ),
                                    ]
                                )
                            ],
                        )
                    ],
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        shape = next(shape for shape in slide.shapes if shape.has_text_frame)
        xml = "".join(p._p.xml for p in shape.text_frame.paragraphs)

        assert "a:highlight" in xml
        assert 'val="F8FAFC"' in xml
        assert xml.index("a:highlight") < xml.index('a:latin typeface="Courier New"')

    def test_transparent_placeholder_run_emits_alpha_zero(self, tmp_path: Path) -> None:
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[
                        SlideElement(
                            element_type=ElementType.PARAGRAPH,
                            box=Box(x=50, y=100, width=600, height=40),
                            paragraphs=[
                                Paragraph(
                                    runs=[
                                        TextRun(text="Before "),
                                        TextRun(
                                            text="inline code",
                                            style=TextStyle(
                                                font_family="Courier New",
                                                color=RGBAColor(r=0, g=0, b=0, a=0.0),
                                            ),
                                        ),
                                        TextRun(text=" after"),
                                    ]
                                )
                            ],
                        )
                    ],
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        shape = next(shape for shape in slide.shapes if shape.has_text_frame)
        xml = "".join(p._p.xml for p in shape.text_frame.paragraphs)

        assert 'a:alpha val="0"' in xml


class TestShapeCount:
    """Verify number of shapes matches element count."""

    def test_adjacent_text_elements_share_one_shape(self, tmp_path: Path) -> None:
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[
                        _make_heading("Title"),
                        _make_paragraph("Body text"),
                    ],
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        assert len(slide.shapes) == 1

    def test_heading_paragraph_and_list_share_one_shape(self, tmp_path: Path) -> None:
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[
                        _make_heading("Title"),
                        _make_paragraph("Body text"),
                        _make_list([("A", 0), ("B", 0)]),
                    ],
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        assert len(slide.shapes) == 1
        text_frame = slide.shapes[0].text_frame
        assert len(text_frame.paragraphs) >= 4

    def test_shape_only_decorated_block_does_not_add_empty_textbox(
        self, tmp_path: Path
    ) -> None:
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[
                        SlideElement(
                            element_type=ElementType.DECORATED_BLOCK,
                            box=Box(x=50, y=100, width=300, height=120),
                            paragraphs=[],
                            decoration=BoxDecoration(
                                background_color=RGBAColor(r=238, g=244, b=255),
                                border_radius_px=16,
                            ),
                        ),
                        SlideElement(
                            element_type=ElementType.PARAGRAPH,
                            box=Box(x=80, y=140, width=200, height=40),
                            paragraphs=[Paragraph(runs=[TextRun(text="Body")])],
                        ),
                    ],
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        assert len(slide.shapes) == 2

    def test_shape_only_decorated_block_uses_css_radius_adjustment(
        self, tmp_path: Path
    ) -> None:
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[
                        SlideElement(
                            element_type=ElementType.DECORATED_BLOCK,
                            box=Box(x=50, y=100, width=200, height=100),
                            paragraphs=[],
                            decoration=BoxDecoration(
                                background_color=RGBAColor(r=238, g=244, b=255),
                                border_radius_px=16,
                            ),
                        )
                    ],
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        shape = next(shape for shape in slide.shapes if shape.text == "")
        assert 'prst="roundRect"' in shape._element.xml
        assert 'a:gd name="adj" fmla="val 16000"' in shape._element.xml

    def test_full_width_heading_does_not_merge_with_narrow_column_text(
        self, tmp_path: Path
    ) -> None:
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[
                        SlideElement(
                            element_type=ElementType.HEADING,
                            box=Box(x=72, y=100, width=1136, height=44),
                            heading_level=2,
                            paragraphs=[Paragraph(runs=[TextRun(text="Split Layout")])],
                        ),
                        SlideElement(
                            element_type=ElementType.HEADING,
                            box=Box(x=72, y=185, width=535, height=30),
                            heading_level=3,
                            paragraphs=[Paragraph(runs=[TextRun(text="Text Region")])],
                        ),
                        SlideElement(
                            element_type=ElementType.UNORDERED_LIST,
                            box=Box(x=72, y=231, width=535, height=160),
                            list_items=[
                                ListItem(runs=[TextRun(text="Left item")], level=0),
                                ListItem(
                                    runs=[TextRun(text="Another left item")], level=0
                                ),
                            ],
                        ),
                    ],
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        text_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
        assert len(text_shapes) == 2
        assert any(shape.text == "Split Layout" for shape in text_shapes)
        assert any("Text Region" in shape.text for shape in text_shapes)
        grouped_shape = next(
            shape for shape in text_shapes if "Text Region" in shape.text
        )
        assert grouped_shape.text_frame.margin_left == 0
        assert grouped_shape.text_frame.margin_right == 0
        assert grouped_shape.text_frame.margin_top == 0
        assert grouped_shape.text_frame.margin_bottom == 0

    def test_nested_list_preserves_marker_styles_and_spacing(
        self, tmp_path: Path
    ) -> None:
        list_element = SlideElement(
            element_type=ElementType.ORDERED_LIST,
            box=Box(x=50, y=160, width=600, height=260),
            list_items=[
                ListItem(
                    runs=[TextRun(text="Top")],
                    level=0,
                    order_number=1,
                    list_style_type="decimal",
                    line_height_px=43.5,
                    space_before_px=0,
                    space_after_px=0,
                ),
                ListItem(
                    runs=[TextRun(text="Nested i")],
                    level=1,
                    order_number=1,
                    list_style_type="lower-roman",
                    line_height_px=43.5,
                    space_before_px=7.25,
                    space_after_px=0,
                ),
            ],
        )
        bullet_element = SlideElement(
            element_type=ElementType.UNORDERED_LIST,
            box=Box(x=680, y=160, width=300, height=220),
            list_items=[
                ListItem(runs=[TextRun(text="Disc")], level=0, list_style_type="disc"),
                ListItem(
                    runs=[TextRun(text="Circle")], level=1, list_style_type="circle"
                ),
                ListItem(
                    runs=[TextRun(text="Square")], level=2, list_style_type="square"
                ),
            ],
        )
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[list_element, bullet_element],
                )
            ]
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        ordered_shape = next(
            shape
            for shape in slide.shapes
            if shape.has_text_frame and "Nested i" in shape.text
        )
        bullet_shape = next(
            shape
            for shape in slide.shapes
            if shape.has_text_frame and "Square" in shape.text
        )

        ordered_xml = "".join(p._p.xml for p in ordered_shape.text_frame.paragraphs)
        bullet_xml = "".join(p._p.xml for p in bullet_shape.text_frame.paragraphs)

        assert 'type="romanLcPeriod"' in ordered_xml
        assert 'val="3262"' in ordered_xml
        assert 'val="543"' in ordered_xml
        assert 'char="•"' in bullet_xml
        assert 'char="◦"' in bullet_xml
        assert 'char="▪"' in bullet_xml

    def test_checklist_items_disable_ppt_bullet_when_list_style_is_none(
        self, tmp_path: Path
    ) -> None:
        checklist = SlideElement(
            element_type=ElementType.UNORDERED_LIST,
            box=Box(x=50, y=160, width=600, height=180),
            list_items=[
                ListItem(
                    runs=[TextRun(text="☐ another checklist item")],
                    level=0,
                    list_style_type="none",
                )
            ],
        )
        pres = Presentation(
            slides=[Slide(width_px=1280, height_px=720, elements=[checklist])]
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        shape = next(
            shape
            for shape in slide.shapes
            if shape.has_text_frame and "☐ another checklist item" in shape.text
        )
        xml = "".join(p._p.xml for p in shape.text_frame.paragraphs)

        assert "a:buNone" in xml
        assert "a:buChar" not in xml
        assert "a:buAutoNum" not in xml

    def test_decorated_block_adds_background_shape_with_text(
        self, tmp_path: Path
    ) -> None:
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[_make_decorated_block("Quoted text")],
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        assert len(slide.shapes) == 3
        texts = [
            para.text
            for shape in slide.shapes
            if shape.has_text_frame
            for para in shape.text_frame.paragraphs
        ]
        assert any("Quoted text" in text for text in texts)

    def test_decorated_block_places_textbox_at_content_box(
        self, tmp_path: Path
    ) -> None:
        element = _make_decorated_block("Quoted text")
        pres = Presentation(
            slides=[Slide(width_px=1280, height_px=720, elements=[element])],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        text_shapes = [
            shape
            for shape in slide.shapes
            if shape.has_text_frame and "Quoted text" in shape.text
        ]
        assert len(text_shapes) == 1
        shape = text_shapes[0]
        assert shape.left == px_to_emu(71)
        assert shape.top == px_to_emu(112)
        assert shape.width == px_to_emu(363)
        assert shape.height == px_to_emu(116)
        assert shape.text_frame.margin_left == 0
        assert shape.text_frame.margin_top == 0

    def test_decorated_block_with_3d_rotation_uses_scene3d_on_shape(
        self, tmp_path: Path
    ) -> None:
        element = _make_decorated_block("AI Engine v3")
        element.rotation_3d_x_deg = 20
        element.rotation_3d_y_deg = 40
        element.rotation_3d_z_deg = 60
        element.vertical_align = "middle"
        pres = Presentation(
            slides=[Slide(width_px=1280, height_px=720, elements=[element])],
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        text_shapes = [
            shape
            for shape in slide.shapes
            if shape.has_text_frame and "AI Engine v3" in shape.text
        ]
        assert len(text_shapes) == 1
        shape = text_shapes[0]
        xml = shape._element.xml
        assert "a:scene3d" in xml
        assert 'prst="orthographicFront"' in xml
        assert 'lat="2400000"' in xml
        assert 'lon="1200000"' in xml
        assert 'rev="3600000"' in xml
        assert shape.text_frame.margin_left == px_to_emu(21)
        assert shape.text_frame.margin_top == px_to_emu(12)

    def test_fit_scene3d_rotations_from_projected_quad(self) -> None:
        box = Box(x=0, y=0, width=282, height=202)
        projected = [
            Point(x=5.0, y=4.0),
            Point(x=282.0, y=0.0),
            Point(x=282.0, y=206.0),
            Point(x=0.0, y=201.0),
        ]
        x_deg, y_deg, z_deg = fit_scene3d_rotations(
            projected,
            box,
            fallback_x_deg=4.0,
            fallback_y_deg=8.0,
            fallback_z_deg=0.0,
        )
        assert (x_deg, y_deg, z_deg) != (4.0, 8.0, 0.0)
        assert max(abs(x_deg), abs(y_deg), abs(z_deg)) >= 8.0

    def test_projected_corners_without_3d_rotation_do_not_emit_scene3d(
        self, tmp_path: Path
    ) -> None:
        element = _make_decorated_block("Professional")
        element.projected_corners = [
            Point(x=element.box.x, y=element.box.y),
            Point(x=element.box.x + element.box.width, y=element.box.y),
            Point(
                x=element.box.x + element.box.width,
                y=element.box.y + element.box.height,
            ),
            Point(x=element.box.x, y=element.box.y + element.box.height),
        ]
        pres = Presentation(
            slides=[Slide(width_px=1280, height_px=720, elements=[element])],
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        text_shapes = [
            shape
            for shape in slide.shapes
            if shape.has_text_frame and "Professional" in shape.text
        ]
        assert len(text_shapes) == 1
        assert "a:scene3d" not in text_shapes[0]._element.xml

    def test_decorated_block_left_accent_respects_rounded_corners(
        self, tmp_path: Path
    ) -> None:
        box = Box(x=50, y=100, width=400, height=140)
        decoration = BoxDecoration(
            background_color=RGBAColor(r=248, g=244, b=255),
            border_left=BorderSide(
                width_px=6,
                style="solid",
                color=RGBAColor(r=74, g=58, b=138),
            ),
            padding=BoxPadding(top_px=12, right_px=16, bottom_px=12, left_px=16),
            border_radius_px=14,
        )
        element = SlideElement(
            element_type=ElementType.DECORATED_BLOCK,
            box=box,
            content_box=_content_box_from_decoration(box, decoration),
            paragraphs=[Paragraph(runs=[TextRun(text="Quoted text")])],
            decoration=decoration,
        )
        pres = Presentation(
            slides=[Slide(width_px=1280, height_px=720, elements=[element])]
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        accent_shape = next(
            shape for shape in slide.shapes if shape.width == px_to_emu(6)
        )

        assert accent_shape.left == px_to_emu(50)
        assert accent_shape.top == px_to_emu(114)
        assert accent_shape.height == px_to_emu(112)

    def test_decorated_block_preserves_list_levels(self, tmp_path: Path) -> None:
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[_make_decorated_list_block()],
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        shape = next(
            shape
            for shape in slide.shapes
            if shape.has_text_frame and "Lead" in shape.text
        )
        levels = [para.level for para in shape.text_frame.paragraphs]
        assert levels == [0, 0, 1]

    def test_decorated_block_preserves_list_styles_and_order_numbers(
        self, tmp_path: Path
    ) -> None:
        box = Box(x=50, y=100, width=420, height=220)
        decoration = BoxDecoration(
            background_color=RGBAColor(r=248, g=244, b=255),
            border_left=BorderSide(
                width_px=5,
                style="solid",
                color=RGBAColor(r=74, g=58, b=138),
            ),
            padding=BoxPadding(top_px=12, right_px=16, bottom_px=12, left_px=16),
            border_radius_px=4,
        )
        element = SlideElement(
            element_type=ElementType.DECORATED_BLOCK,
            box=box,
            content_box=_content_box_from_decoration(box, decoration),
            paragraphs=[
                Paragraph(runs=[TextRun(text="Lead")]),
                Paragraph(
                    runs=[TextRun(text="Top")],
                    list_level=0,
                    list_ordered=True,
                    list_style_type="decimal",
                    order_number=7,
                    line_height_px=43.5,
                ),
                Paragraph(
                    runs=[TextRun(text="Nested")],
                    list_level=1,
                    list_ordered=True,
                    list_style_type="lower-roman",
                    order_number=2,
                    line_height_px=43.5,
                    space_before_px=7.25,
                ),
            ],
            decoration=decoration,
        )
        pres = Presentation(
            slides=[Slide(width_px=1280, height_px=720, elements=[element])]
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        shape = next(
            shape
            for shape in slide.shapes
            if shape.has_text_frame and "Nested" in shape.text
        )
        xml = "".join(p._p.xml for p in shape.text_frame.paragraphs)

        assert 'type="arabicPeriod"' in xml
        assert 'type="romanLcPeriod"' in xml
        assert 'startAt="7"' in xml
        assert 'startAt="2"' in xml

    def test_decoration_opacity_emits_fill_and_accent_alpha(
        self, tmp_path: Path
    ) -> None:
        box = Box(x=50, y=100, width=400, height=160)
        decoration = BoxDecoration(
            background_color=RGBAColor(r=255, g=0, b=0),
            border_left=BorderSide(
                width_px=6,
                style="solid",
                color=RGBAColor(r=0, g=0, b=255),
            ),
            opacity=0.5,
        )
        element = SlideElement(
            element_type=ElementType.DECORATED_BLOCK,
            box=box,
            content_box=_content_box_from_decoration(box, decoration),
            paragraphs=[Paragraph(runs=[TextRun(text="Opaque enough")])],
            decoration=decoration,
        )
        pres = Presentation(
            slides=[Slide(width_px=1280, height_px=720, elements=[element])]
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        bg_shape = next(shape for shape in slide.shapes if shape.has_text_frame)
        accent_shape = next(
            shape for shape in slide.shapes if shape.width == px_to_emu(6)
        )

        assert 'a:srgbClr val="FF0000"' in bg_shape._element.xml
        assert 'a:alpha val="50000"' in bg_shape._element.xml
        assert 'a:srgbClr val="0000FF"' in accent_shape._element.xml
        assert 'a:alpha val="50000"' in accent_shape._element.xml

    def test_linear_gradient_decoration_renders_native_shape_background(
        self, tmp_path: Path
    ) -> None:
        box = Box(x=50, y=100, width=320, height=120)
        decoration = BoxDecoration(
            background_gradient="linear-gradient(135deg, rgba(255, 0, 0, 0.5), rgba(0, 0, 255, 0.5))",
            border_radius_px=12,
        )
        element = SlideElement(
            element_type=ElementType.DECORATED_BLOCK,
            box=box,
            content_box=_content_box_from_decoration(box, decoration),
            paragraphs=[Paragraph(runs=[TextRun(text="Gradient card")])],
            decoration=decoration,
        )
        pres = Presentation(
            slides=[Slide(width_px=1280, height_px=720, elements=[element])]
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]

        bg_shape = next(
            shape for shape in slide.shapes if shape.has_text_frame and shape.text == ""
        )
        assert "a:gradFill" in bg_shape._element.xml
        assert 'val="FF0000"' in bg_shape._element.xml
        assert 'val="0000FF"' in bg_shape._element.xml
        assert 'a:lin ang="18900000"' in bg_shape._element.xml
        assert any(
            shape.has_text_frame and "Gradient card" in shape.text
            for shape in slide.shapes
        )

    def test_horizontal_linear_gradient_decoration_maps_to_ooxml_zero_angle(
        self, tmp_path: Path
    ) -> None:
        box = Box(x=50, y=100, width=320, height=6)
        decoration = BoxDecoration(
            background_gradient="linear-gradient(90deg, #ff0000, #0000ff)",
        )
        element = SlideElement(
            element_type=ElementType.DECORATED_BLOCK,
            box=box,
            content_box=box,
            paragraphs=[],
            decoration=decoration,
        )
        pres = Presentation(
            slides=[Slide(width_px=1280, height_px=720, elements=[element])]
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        bg_shape = next(shape for shape in slide.shapes if shape.has_text_frame)

        assert 'a:lin ang="0"' in bg_shape._element.xml

    def test_box_shadow_decoration_emits_outer_shadow(self, tmp_path: Path) -> None:
        box = Box(x=50, y=100, width=320, height=120)
        decoration = BoxDecoration(
            background_color=RGBAColor(r=255, g=255, b=255),
            border_radius_px=12,
            box_shadows=[
                BoxShadow(
                    offset_x_px=0,
                    offset_y_px=8,
                    blur_radius_px=32,
                    color=RGBAColor(r=59, g=130, b=246, a=0.15),
                )
            ],
        )
        element = SlideElement(
            element_type=ElementType.DECORATED_BLOCK,
            box=box,
            content_box=_content_box_from_decoration(box, decoration),
            paragraphs=[Paragraph(runs=[TextRun(text="Shadow card")])],
            decoration=decoration,
        )
        pres = Presentation(
            slides=[Slide(width_px=1280, height_px=720, elements=[element])]
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        shadow_shape = next(
            shape for shape in slide.shapes if "a:outerShdw" in shape._element.xml
        )
        xml = shadow_shape._element.xml
        assert 'blurRad="304800"' in xml
        assert 'dist="76200"' in xml
        assert 'dir="5400000"' in xml
        assert 'val="3B82F6"' in xml
        assert 'a:alpha val="15000"' in xml

    def test_spread_only_box_shadow_renders_outline_shape(self, tmp_path: Path) -> None:
        box = Box(x=50, y=100, width=48, height=48)
        decoration = BoxDecoration(
            background_gradient="linear-gradient(135deg, #3b82f6, #8b5cf6)",
            border_radius_px=24,
            box_shadows=[
                BoxShadow(
                    offset_x_px=0,
                    offset_y_px=0,
                    blur_radius_px=0,
                    spread_px=2,
                    color=RGBAColor(r=139, g=92, b=246),
                )
            ],
        )
        element = SlideElement(
            element_type=ElementType.DECORATED_BLOCK,
            box=box,
            content_box=box,
            paragraphs=[],
            decoration=decoration,
        )
        pres = Presentation(
            slides=[Slide(width_px=1280, height_px=720, elements=[element])]
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        outline_shape = next(
            shape
            for shape in slide.shapes
            if f'w="{px_to_emu(4)}"' in shape._element.xml
            and 'val="8B5CF6"' in shape._element.xml
        )
        assert "a:outerShdw" not in outline_shape._element.xml

    def test_inset_box_shadow_emits_inner_shadow(self, tmp_path: Path) -> None:
        box = Box(x=50, y=100, width=320, height=120)
        decoration = BoxDecoration(
            background_color=RGBAColor(r=255, g=255, b=255),
            border_radius_px=12,
            box_shadows=[
                BoxShadow(
                    offset_x_px=0,
                    offset_y_px=2,
                    blur_radius_px=8,
                    spread_px=4,
                    color=RGBAColor(r=15, g=23, b=42, a=0.08),
                    inset=True,
                )
            ],
        )
        element = SlideElement(
            element_type=ElementType.DECORATED_BLOCK,
            box=box,
            content_box=_content_box_from_decoration(box, decoration),
            paragraphs=[Paragraph(runs=[TextRun(text="Inset card")])],
            decoration=decoration,
        )
        pres = Presentation(
            slides=[Slide(width_px=1280, height_px=720, elements=[element])]
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        shadow_shape = next(
            shape for shape in slide.shapes if "a:innerShdw" in shape._element.xml
        )
        xml = shadow_shape._element.xml
        assert 'blurRad="76200"' in xml
        assert 'dist="19050"' in xml
        assert 'dir="5400000"' in xml
        assert 'sx="102500"' in xml
        assert 'sy="106667"' in xml
        assert 'val="0F172A"' in xml
        assert 'a:alpha val="8000"' in xml

    def test_blockquote_uses_extracted_padding_and_left_accent(
        self, tmp_path: Path
    ) -> None:
        box = Box(x=80, y=120, width=420, height=110)
        decoration = BoxDecoration(
            border_left=BorderSide(
                width_px=4,
                style="solid",
                color=RGBAColor(r=148, g=163, b=184),
            ),
            padding=BoxPadding(left_px=16),
        )
        element = SlideElement(
            element_type=ElementType.BLOCKQUOTE,
            box=box,
            content_box=_content_box_from_decoration(box, decoration),
            paragraphs=[Paragraph(runs=[TextRun(text="Quoted text")])],
            decoration=decoration,
        )
        pres = Presentation(
            slides=[Slide(width_px=1280, height_px=720, elements=[element])]
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        text_shape = next(
            shape
            for shape in slide.shapes
            if shape.has_text_frame and "Quoted text" in shape.text
        )
        accent_shape = next(
            shape for shape in slide.shapes if shape.width == px_to_emu(4)
        )

        assert text_shape.left == px_to_emu(100)
        assert text_shape.width == px_to_emu(400)
        assert text_shape.text_frame.margin_left == 0
        assert accent_shape.left == px_to_emu(80)
        assert accent_shape.height == px_to_emu(110)

    def test_math_fallback_adds_picture_shape(self, tmp_path: Path) -> None:
        fallback_image = tmp_path / "math.png"
        fallback_image.write_bytes(
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
            b"\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx\x9cc\xf8\xcf\xc0"
            b"\x00\x00\x03\x01\x01\x00\xc9\xfe\x92\xef\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        math_element = SlideElement(
            element_type=ElementType.MATH,
            box=Box(x=50, y=50, width=100, height=40),
            unsupported_info=UnsupportedInfo(
                reason="Math expression",
                tag_name="mjx-container",
                fallback_image_path=str(fallback_image),
            ),
        )
        pres = Presentation(
            slides=[Slide(width_px=1280, height_px=720, elements=[math_element])]
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        assert len(slide.shapes) == 1
        assert slide.shapes[0].shape_type is not None


class TestListParagraphLevels:
    """Verify list items have correct paragraph levels."""

    def test_flat_list_level_0(self, tmp_path: Path) -> None:
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[
                        _make_list([("A", 0), ("B", 0), ("C", 0)]),
                    ],
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    assert para.level == 0

    def test_nested_list_levels(self, tmp_path: Path) -> None:
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[
                        _make_list([("A", 0), ("A1", 1), ("A2", 2), ("B", 0)]),
                    ],
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        for shape in slide.shapes:
            if shape.has_text_frame:
                levels = [p.level for p in shape.text_frame.paragraphs]
                assert levels == [0, 1, 2, 0]


class TestTableBuilding:
    """Verify table structure."""

    def test_table_rows_and_cols(self, tmp_path: Path) -> None:
        table_data = [
            ["H1", "H2", "H3"],
            ["A1", "A2", "A3"],
            ["B1", "B2", "B3"],
        ]
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[_make_table(table_data)],
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]

        # Find the table shape
        table_shape = None
        for shape in slide.shapes:
            if shape.has_table:
                table_shape = shape
                break
        assert table_shape is not None

        table = table_shape.table
        assert len(table.rows) == 3
        assert len(table.columns) == 3

    def test_table_cell_text(self, tmp_path: Path) -> None:
        table_data = [
            ["Name", "Value"],
            ["Alpha", "100"],
        ]
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[_make_table(table_data)],
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]

        for shape in slide.shapes:
            if shape.has_table:
                cell_text = shape.table.cell(1, 0).text
                assert "Alpha" in cell_text

    def test_table_preserves_extracted_column_widths(self, tmp_path: Path) -> None:
        element = SlideElement(
            element_type=ElementType.TABLE,
            box=Box(x=72, y=100, width=1136, height=160),
            table_rows=[
                TableRow(
                    cells=[
                        TableCell(
                            paragraphs=[Paragraph(runs=[TextRun(text="A")])],
                            width_px=100,
                            is_header=True,
                        ),
                        TableCell(
                            paragraphs=[Paragraph(runs=[TextRun(text="B")])],
                            width_px=230,
                            is_header=True,
                        ),
                        TableCell(
                            paragraphs=[Paragraph(runs=[TextRun(text="C")])],
                            width_px=110,
                            is_header=True,
                        ),
                    ]
                ),
                TableRow(
                    cells=[
                        TableCell(
                            paragraphs=[Paragraph(runs=[TextRun(text="1")])],
                            width_px=100,
                        ),
                        TableCell(
                            paragraphs=[Paragraph(runs=[TextRun(text="2")])],
                            width_px=230,
                        ),
                        TableCell(
                            paragraphs=[Paragraph(runs=[TextRun(text="3")])],
                            width_px=110,
                        ),
                    ]
                ),
            ],
        )
        pres = Presentation(
            slides=[Slide(width_px=1280, height_px=720, elements=[element])]
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        table_shape = next(shape for shape in slide.shapes if shape.has_table)
        widths = [col.width for col in table_shape.table.columns]
        assert widths[0] == px_to_emu(100)
        assert widths[1] == px_to_emu(230)
        assert widths[2] == px_to_emu(110)

    def test_table_cell_emits_gradient_fill_padding_and_border(
        self, tmp_path: Path
    ) -> None:
        element = SlideElement(
            element_type=ElementType.TABLE,
            box=Box(x=72, y=100, width=400, height=120),
            table_rows=[
                TableRow(
                    cells=[
                        TableCell(
                            paragraphs=[Paragraph(runs=[TextRun(text="Header")])],
                            is_header=True,
                            background_gradient="linear-gradient(135deg, #3b82f6, #2563eb)",
                            padding=BoxPadding(
                                top_px=14, right_px=16, bottom_px=14, left_px=16
                            ),
                            border_bottom=BorderSide(
                                width_px=1,
                                style="solid",
                                color=RGBAColor(r=226, g=232, b=240),
                            ),
                        )
                    ]
                )
            ],
        )
        pres = Presentation(
            slides=[Slide(width_px=1280, height_px=720, elements=[element])]
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        table_shape = next(shape for shape in slide.shapes if shape.has_table)
        cell = table_shape.table.cell(0, 0)

        assert cell.text_frame.margin_left == px_to_emu(16)
        assert cell.text_frame.margin_top == px_to_emu(14)
        assert "a:gradFill" in cell._tc.xml
        assert 'a:lin ang="18900000"' in cell._tc.xml
        assert "a:lnB" in cell._tc.xml
        assert 'w="12700"' in cell._tc.xml
        assert 'cap="flat"' in cell._tc.xml
        assert 'cmpd="sng"' in cell._tc.xml
        assert 'algn="ctr"' in cell._tc.xml
        assert "<a:round/>" in cell._tc.xml
        assert "<a:headEnd" in cell._tc.xml
        assert "<a:tailEnd" in cell._tc.xml

    def test_table_decoration_adds_shadow_backplate(self, tmp_path: Path) -> None:
        element = _make_table([["H1", "H2"], ["A1", "A2"]])
        element.decoration = BoxDecoration(
            border_radius_px=12,
            box_shadows=[
                BoxShadow(
                    offset_x_px=0,
                    offset_y_px=4,
                    blur_radius_px=24,
                    color=RGBAColor(r=0, g=0, b=0, a=0.08),
                )
            ],
        )
        pres = Presentation(
            slides=[Slide(width_px=1280, height_px=720, elements=[element])]
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        assert any("a:outerShdw" in shape._element.xml for shape in slide.shapes)

    def test_table_style_is_removed_so_custom_borders_win(self, tmp_path: Path) -> None:
        table_data = [["H1", "H2"], ["A1", "A2"]]
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[_make_table(table_data)],
                )
            ],
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        table_shape = next(shape for shape in slide.shapes if shape.has_table)
        assert "tableStyleId" not in table_shape.table._tbl.xml


class TestCodeBlock:
    """Verify code block has background fill."""

    def test_code_block_has_fill(self, tmp_path: Path) -> None:
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[_make_code_block("print('hello')")],
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]

        # Code block should be a textbox with a fill
        for shape in slide.shapes:
            if shape.has_text_frame:
                # Check that fill is set (solid fill)
                fill = shape.fill
                assert fill.type is not None  # Should have some fill type

    def test_code_block_with_decoration_uses_rounded_background_shape(
        self, tmp_path: Path
    ) -> None:
        decoration = BoxDecoration(
            background_color=RGBAColor(r=246, g=248, b=250),
            border_top=BorderSide(
                width_px=1, style="solid", color=RGBAColor(r=209, g=217, b=224)
            ),
            border_right=BorderSide(
                width_px=1, style="solid", color=RGBAColor(r=209, g=217, b=224)
            ),
            border_bottom=BorderSide(
                width_px=1, style="solid", color=RGBAColor(r=209, g=217, b=224)
            ),
            border_left=BorderSide(
                width_px=1, style="solid", color=RGBAColor(r=209, g=217, b=224)
            ),
            border_radius_px=6,
            padding=BoxPadding(top_px=16, right_px=16, bottom_px=16, left_px=16),
        )
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[
                        _make_code_block("print('hello')", decoration=decoration)
                    ],
                )
            ],
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        assert len(slide.shapes) == 2

        background_shape = next(
            shape for shape in slide.shapes if shape.has_text_frame and shape.text == ""
        )
        textbox = next(
            shape for shape in slide.shapes if shape.has_text_frame and shape.text != ""
        )

        assert 'prst="roundRect"' in background_shape._element.xml
        assert textbox.left > px_to_emu(50)
        assert textbox.top > px_to_emu(100)
        assert "<a:noFill/>" in textbox._element.xml


class TestDecorationRendering:
    def test_non_uniform_bottom_border_renders_as_separate_shape(
        self, tmp_path: Path
    ) -> None:
        element = SlideElement(
            element_type=ElementType.DECORATED_BLOCK,
            box=Box(x=50, y=100, width=300, height=40),
            paragraphs=[],
            decoration=BoxDecoration(
                border_bottom=BorderSide(
                    width_px=1,
                    style="solid",
                    color=RGBAColor(r=51, g=65, b=85),
                ),
            ),
        )
        pres = Presentation(
            slides=[Slide(width_px=1280, height_px=720, elements=[element])]
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]

        assert len(slide.shapes) == 2
        border_shape = next(
            shape for shape in slide.shapes if 'val="334155"' in shape._element.xml
        )
        assert border_shape.height == px_to_emu(1)
        assert border_shape.width == px_to_emu(300)


class TestImagePlacement:
    """Verify CSS object-fit placement."""

    def test_contain_fit_preserves_aspect_ratio(self) -> None:
        element = SlideElement(
            element_type=ElementType.IMAGE,
            box=Box(x=10, y=20, width=240, height=120),
            image_src="data:image/png;base64,abc",
            image_natural_width_px=1,
            image_natural_height_px=1,
            object_fit="contain",
            object_position="50% 50%",
        )

        left, top, width, height = _resolve_image_placement(element)
        assert width == px_to_emu(120)
        assert height == px_to_emu(120)
        assert left == px_to_emu(70)
        assert top == px_to_emu(20)


class TestImageRendering:
    """Verify image sources are embedded into PPTX."""

    def test_file_uri_image_is_embedded(self, tmp_path: Path) -> None:
        image_path = tmp_path / "sample image.png"
        image_path.write_bytes(
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
            b"\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx\x9cc\xf8\xcf\xc0"
            b"\x00\x00\x03\x01\x01\x00\xc9\xfe\x92\xef\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[
                        SlideElement(
                            element_type=ElementType.IMAGE,
                            box=Box(x=10, y=20, width=240, height=120),
                            image_src=image_path.as_uri(),
                        )
                    ],
                )
            ],
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        assert len(slide.shapes) == 1
        assert slide.shapes[0].shape_type is not None

    def test_image_with_decoration_adds_background_shape(self, tmp_path: Path) -> None:
        image_path = tmp_path / "sample image.png"
        image_path.write_bytes(
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
            b"\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx\x9cc\xf8\xcf\xc0"
            b"\x00\x00\x03\x01\x01\x00\xc9\xfe\x92\xef\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[
                        SlideElement(
                            element_type=ElementType.IMAGE,
                            box=Box(x=10, y=20, width=240, height=120),
                            image_src=image_path.as_uri(),
                            object_fit="contain",
                            image_natural_width_px=1,
                            image_natural_height_px=1,
                            decoration=BoxDecoration(
                                background_color=RGBAColor(r=255, g=255, b=255),
                                border_top=BorderSide(
                                    width_px=1,
                                    style="solid",
                                    color=RGBAColor(r=203, g=213, b=225),
                                ),
                                border_right=BorderSide(
                                    width_px=1,
                                    style="solid",
                                    color=RGBAColor(r=203, g=213, b=225),
                                ),
                                border_bottom=BorderSide(
                                    width_px=1,
                                    style="solid",
                                    color=RGBAColor(r=203, g=213, b=225),
                                ),
                                border_left=BorderSide(
                                    width_px=1,
                                    style="solid",
                                    color=RGBAColor(r=203, g=213, b=225),
                                ),
                                border_radius_px=12,
                            ),
                        )
                    ],
                )
            ],
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        assert len(slide.shapes) == 2
        background_shape = next(
            shape for shape in slide.shapes if not shape.has_text_frame
        )
        assert "p:style" not in background_shape._element.xml
        assert "effectRef" not in background_shape._element.xml

    def test_standalone_decorated_text_uses_separate_textbox(
        self, tmp_path: Path
    ) -> None:
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[
                        SlideElement(
                            element_type=ElementType.DECORATED_BLOCK,
                            box=Box(x=80, y=120, width=120, height=36),
                            paragraphs=[Paragraph(runs=[TextRun(text="inline code")])],
                            decoration=BoxDecoration(
                                background_color=RGBAColor(r=232, g=237, b=245),
                                border_radius_px=6,
                                padding=BoxPadding(
                                    top_px=2, right_px=6, bottom_px=2, left_px=6
                                ),
                            ),
                        )
                    ],
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        assert len(slide.shapes) == 2
        text_shape = next(
            shape
            for shape in slide.shapes
            if shape.has_text_frame and shape.text == "inline code"
        )
        assert text_shape.has_text_frame
        assert any(shape.has_text_frame and shape.text == "" for shape in slide.shapes)

    def test_image_with_decoration_is_inset_inside_border(self, tmp_path: Path) -> None:
        image_path = tmp_path / "sample image.png"
        image_path.write_bytes(
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
            b"\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx\x9cc\xf8\xcf\xc0"
            b"\x00\x00\x03\x01\x01\x00\xc9\xfe\x92\xef\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        element = SlideElement(
            element_type=ElementType.IMAGE,
            box=Box(x=10, y=20, width=240, height=120),
            image_src=image_path.as_uri(),
            image_natural_width_px=1,
            image_natural_height_px=1,
            object_fit="contain",
            decoration=BoxDecoration(
                background_color=RGBAColor(r=255, g=255, b=255),
                border_top=BorderSide(
                    width_px=2, style="solid", color=RGBAColor(r=0, g=0, b=0)
                ),
                border_right=BorderSide(
                    width_px=2, style="solid", color=RGBAColor(r=0, g=0, b=0)
                ),
                border_bottom=BorderSide(
                    width_px=2, style="solid", color=RGBAColor(r=0, g=0, b=0)
                ),
                border_left=BorderSide(
                    width_px=2, style="solid", color=RGBAColor(r=0, g=0, b=0)
                ),
            ),
        )
        pres = Presentation(
            slides=[Slide(width_px=1280, height_px=720, elements=[element])]
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        picture = next(shape for shape in slide.shapes if shape.shape_type == 13)
        assert picture.left > px_to_emu(10)
        assert picture.top > px_to_emu(20)

    def test_image_with_contain_fit_keeps_picture_geometry_unrounded(
        self, tmp_path: Path
    ) -> None:
        image_path = tmp_path / "sample image.png"
        image_path.write_bytes(
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
            b"\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx\x9cc\xf8\xcf\xc0"
            b"\x00\x00\x03\x01\x01\x00\xc9\xfe\x92\xef\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        element = SlideElement(
            element_type=ElementType.IMAGE,
            box=Box(x=10, y=20, width=240, height=120),
            image_src=image_path.as_uri(),
            image_natural_width_px=1,
            image_natural_height_px=1,
            object_fit="contain",
            decoration=BoxDecoration(
                background_color=RGBAColor(r=255, g=255, b=255),
                border_top=BorderSide(
                    width_px=1, style="solid", color=RGBAColor(r=0, g=0, b=0)
                ),
                border_right=BorderSide(
                    width_px=1, style="solid", color=RGBAColor(r=0, g=0, b=0)
                ),
                border_bottom=BorderSide(
                    width_px=1, style="solid", color=RGBAColor(r=0, g=0, b=0)
                ),
                border_left=BorderSide(
                    width_px=1, style="solid", color=RGBAColor(r=0, g=0, b=0)
                ),
                border_radius_px=12,
            ),
        )
        pres = Presentation(
            slides=[Slide(width_px=1280, height_px=720, elements=[element])]
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        picture = next(shape for shape in slide.shapes if shape.shape_type == 13)
        assert 'prst="roundRect"' not in picture._element.xml

    def test_image_with_cover_fit_uses_rounded_picture_geometry(
        self, tmp_path: Path
    ) -> None:
        image_path = tmp_path / "sample image.png"
        image_path.write_bytes(
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
            b"\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx\x9cc\xf8\xcf\xc0"
            b"\x00\x00\x03\x01\x01\x00\xc9\xfe\x92\xef\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        element = SlideElement(
            element_type=ElementType.IMAGE,
            box=Box(x=10, y=20, width=240, height=120),
            image_src=image_path.as_uri(),
            object_fit="cover",
            decoration=BoxDecoration(
                background_color=RGBAColor(r=255, g=255, b=255),
                border_top=BorderSide(
                    width_px=1, style="solid", color=RGBAColor(r=0, g=0, b=0)
                ),
                border_right=BorderSide(
                    width_px=1, style="solid", color=RGBAColor(r=0, g=0, b=0)
                ),
                border_bottom=BorderSide(
                    width_px=1, style="solid", color=RGBAColor(r=0, g=0, b=0)
                ),
                border_left=BorderSide(
                    width_px=1, style="solid", color=RGBAColor(r=0, g=0, b=0)
                ),
                border_radius_px=12,
            ),
        )
        pres = Presentation(
            slides=[Slide(width_px=1280, height_px=720, elements=[element])]
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        picture = next(shape for shape in slide.shapes if shape.shape_type == 13)
        assert 'prst="roundRect"' in picture._element.xml

    def test_image_opacity_emits_blip_alpha(self, tmp_path: Path) -> None:
        image_path = tmp_path / "sample image.png"
        image_path.write_bytes(
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
            b"\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx\x9cc\xf8\xcf\xc0"
            b"\x00\x00\x03\x01\x01\x00\xc9\xfe\x92\xef\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        element = SlideElement(
            element_type=ElementType.IMAGE,
            box=Box(x=10, y=20, width=240, height=120),
            image_src=image_path.as_uri(),
            image_opacity=0.4,
        )
        pres = Presentation(
            slides=[Slide(width_px=1280, height_px=720, elements=[element])]
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        picture = next(shape for shape in slide.shapes if shape.shape_type == 13)

        assert 'a:alphaModFix amt="40000"' in picture._element.xml

    def test_file_uri_svg_is_rasterized_and_embedded(self, tmp_path: Path) -> None:
        svg_path = tmp_path / "diagram.svg"
        svg_path.write_text(
            """<svg xmlns="http://www.w3.org/2000/svg" width="120" height="60" viewBox="0 0 120 60">
  <rect width="120" height="60" fill="#224488"/>
  <text x="60" y="36" text-anchor="middle" font-size="20" fill="#ffffff">SVG</text>
</svg>
""",
            encoding="utf-8",
        )
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[
                        SlideElement(
                            element_type=ElementType.IMAGE,
                            box=Box(x=10, y=20, width=240, height=120),
                            image_src=svg_path.as_uri(),
                        )
                    ],
                )
            ],
        )

        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        assert len(slide.shapes) == 1
        assert slide.shapes[0].width == px_to_emu(240)

    def test_missing_rsvg_convert_raises_clear_error(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ) -> None:
        svg_path = tmp_path / "diagram.svg"
        svg_path.write_text(
            """<svg xmlns="http://www.w3.org/2000/svg" width="120" height="60" viewBox="0 0 120 60">
  <rect width="120" height="60" fill="#224488"/>
</svg>
""",
            encoding="utf-8",
        )
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[
                        SlideElement(
                            element_type=ElementType.IMAGE,
                            box=Box(x=10, y=20, width=240, height=120),
                            image_src=svg_path.as_uri(),
                        )
                    ],
                )
            ],
        )
        monkeypatch.setattr("marpx.svg_utils.shutil.which", lambda _: None)

        with pytest.raises(
            MissingDependencyError,
            match=r"rsvg-convert.*brew install librsvg.*librsvg2-bin",
        ):
            build_pptx(pres, tmp_path / "out.pptx")


class TestSolidBackground:
    """Verify solid background color is applied."""

    def test_slide_has_background(self, tmp_path: Path) -> None:
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[_make_heading("Dark Slide")],
                    background=Background(color=RGBAColor(r=30, g=30, b=30, a=1.0)),
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]

        # The slide background fill should be set
        bg = slide.background
        fill = bg.fill
        assert fill.type is not None

    def test_slide_linear_gradient_background_renders_picture(
        self, tmp_path: Path
    ) -> None:
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[_make_heading("Gradient Slide")],
                    background=Background(
                        background_gradient="linear-gradient(135deg, #667eea 0%, #764ba2 100%)"
                    ),
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]

        assert any(shape.shape_type == 13 for shape in slide.shapes)

    def test_slide_radial_gradient_background_renders_picture(
        self, tmp_path: Path
    ) -> None:
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[_make_heading("Gradient Slide")],
                    background=Background(
                        background_gradient="radial-gradient(ellipse at 30% 50%, #312e81, #0f172a 70%)"
                    ),
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]

        assert any(shape.shape_type == 13 for shape in slide.shapes)


class TestOutputFile:
    """Verify the PPTX file is written to disk."""

    def test_output_exists(self, tmp_path: Path) -> None:
        pres = Presentation(
            slides=[Slide(width_px=1280, height_px=720, elements=[])],
        )
        out = tmp_path / "result.pptx"
        result = build_pptx(pres, out)
        assert result.exists()
        assert result.stat().st_size > 0


# ---------------------------------------------------------------------------
# Helpers for merged-table tests
# ---------------------------------------------------------------------------


def _make_merged_table(
    rows_data: list[list[tuple[str, int, int]]],
    has_header: bool = True,
) -> SlideElement:
    """Build a table element with colspan/rowspan support.

    Each cell is a tuple of (text, colspan, rowspan).
    """
    table_rows = []
    for r_idx, row in enumerate(rows_data):
        cells = [
            TableCell(
                paragraphs=[Paragraph(runs=[TextRun(text=text)])],
                colspan=cs,
                rowspan=rs,
                is_header=(r_idx == 0 and has_header),
            )
            for text, cs, rs in row
        ]
        table_rows.append(TableRow(cells=cells))
    return SlideElement(
        element_type=ElementType.TABLE,
        box=Box(x=50, y=100, width=600, height=300),
        table_rows=table_rows,
    )


def _find_table(pptx_pres: PptxPresentation):
    """Return the first table object from the first slide."""
    for shape in pptx_pres.slides[0].shapes:
        if shape.has_table:
            return shape.table
    return None


# ---------------------------------------------------------------------------
# Tests for table merge (colspan / rowspan)
# ---------------------------------------------------------------------------


class TestTableMergeColspan:
    """Verify colspan merges produce correct grid dimensions and content."""

    def test_colspan_header(self, tmp_path: Path) -> None:
        """Header spans 2 columns; body has 2 separate cells."""
        rows = [
            [("Header", 2, 1)],  # colspan=2
            [("A", 1, 1), ("B", 1, 1)],  # normal
        ]
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[_make_merged_table(rows)],
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        table = _find_table(pptx)
        assert table is not None

        # Grid must be 2 rows x 2 cols
        assert len(table.rows) == 2
        assert len(table.columns) == 2

        # Top-left merged cell should contain header text
        assert "Header" in table.cell(0, 0).text

        # Body cells
        assert "A" in table.cell(1, 0).text
        assert "B" in table.cell(1, 1).text


class TestTableMergeRowspan:
    """Verify rowspan merges."""

    def test_rowspan_first_column(self, tmp_path: Path) -> None:
        """First column spans 2 rows; second column has separate cells."""
        rows = [
            [("Span", 1, 2), ("R0C1", 1, 1)],
            [("R1C1", 1, 1)],  # only 1 cell because col 0 is occupied
        ]
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[_make_merged_table(rows, has_header=False)],
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        table = _find_table(pptx)
        assert table is not None

        assert len(table.rows) == 2
        assert len(table.columns) == 2

        # Merged cell content
        assert "Span" in table.cell(0, 0).text

        # Non-merged cells
        assert "R0C1" in table.cell(0, 1).text
        assert "R1C1" in table.cell(1, 1).text


class TestTableMergeColspanAndRowspan:
    """Verify combined colspan + rowspan merges."""

    def test_combined_merge(self, tmp_path: Path) -> None:
        """A cell spanning 2 cols x 2 rows in the top-left."""
        rows = [
            [("Big", 2, 2), ("R0C2", 1, 1)],
            [("R1C2", 1, 1)],  # col 0-1 occupied by Big
            [("R2C0", 1, 1), ("R2C1", 1, 1), ("R2C2", 1, 1)],
        ]
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[_make_merged_table(rows, has_header=False)],
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        table = _find_table(pptx)
        assert table is not None

        assert len(table.rows) == 3
        assert len(table.columns) == 3

        # Merged region top-left cell
        assert "Big" in table.cell(0, 0).text

        # Cells outside the merged region
        assert "R0C2" in table.cell(0, 2).text
        assert "R1C2" in table.cell(1, 2).text
        assert "R2C0" in table.cell(2, 0).text
        assert "R2C1" in table.cell(2, 1).text
        assert "R2C2" in table.cell(2, 2).text


class TestTableNoMergeRegression:
    """Ensure simple tables without merge still work correctly."""

    def test_simple_table_unchanged(self, tmp_path: Path) -> None:
        """A plain 2x3 table with no spans should produce correct content."""
        table_data = [
            ["H1", "H2", "H3"],
            ["A1", "A2", "A3"],
        ]
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[_make_table(table_data)],
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        table = _find_table(pptx)
        assert table is not None

        assert len(table.rows) == 2
        assert len(table.columns) == 3

        assert "H1" in table.cell(0, 0).text
        assert "A3" in table.cell(1, 2).text
