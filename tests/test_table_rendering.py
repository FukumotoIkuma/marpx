"""Direct unit tests for table rendering in pptx_builder/table.py.

These tests construct model objects directly -- NO Playwright or marp-cli needed.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation as PptxPresentation
from pptx.shapes.graphfrm import GraphicFrame

from marpx.models import (
    Box,
    ElementType,
    Paragraph,
    Presentation,
    RGBAColor,
    Slide,
    SlideElement,
    TableCell,
    TableElement,
    TableRow,
    TextRun,
)
from marpx.pptx_builder.builder import build_pptx


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _make_slide_with(*elements: SlideElement) -> Presentation:
    return Presentation(
        slides=[Slide(width_px=1280, height_px=720, elements=list(elements))]
    )


def _make_simple_table(
    rows_data: list[list[str]],
    has_header: bool = True,
) -> SlideElement:
    table_rows = []
    for r_idx, row in enumerate(rows_data):
        cells = [
            TableCell(
                paragraphs=[Paragraph(runs=[TextRun(text=cell_text)])],
                is_header=(r_idx == 0 and has_header),
            )
            for cell_text in row
        ]
        table_rows.append(TableRow(cells=cells))
    return TableElement(
        element_type=ElementType.TABLE,
        box=Box(x=50, y=100, width=600, height=300),
        table_rows=table_rows,
    )


def _build_and_read(presentation: Presentation, tmp_path: Path) -> PptxPresentation:
    """Build PPTX and re-open it for assertions."""
    out = tmp_path / "test_output.pptx"
    build_pptx(presentation, out)
    return PptxPresentation(str(out))


def _first_table(pptx: PptxPresentation):
    """Return the first table shape on slide 0."""
    slide = pptx.slides[0]
    for shape in slide.shapes:
        if isinstance(shape, GraphicFrame) and shape.has_table:
            return shape.table
    raise AssertionError("No table found on slide")


# ---------------------------------------------------------------------------
# TestBasicTable
# ---------------------------------------------------------------------------


class TestBasicTable:
    """Verify shape creation, cell text, and header row flags."""

    def test_table_shape_created(self, tmp_path: Path) -> None:
        pres = _make_slide_with(_make_simple_table([["H1", "H2"], ["A", "B"]]))
        pptx = _build_and_read(pres, tmp_path)
        table = _first_table(pptx)
        assert table is not None

    def test_cell_text_content(self, tmp_path: Path) -> None:
        pres = _make_slide_with(_make_simple_table([["Header"], ["Body"]]))
        pptx = _build_and_read(pres, tmp_path)
        table = _first_table(pptx)
        assert table.cell(0, 0).text == "Header"
        assert table.cell(1, 0).text == "Body"

    def test_cell_text_multiple_columns(self, tmp_path: Path) -> None:
        pres = _make_slide_with(_make_simple_table([["A", "B", "C"], ["D", "E", "F"]]))
        pptx = _build_and_read(pres, tmp_path)
        table = _first_table(pptx)
        assert table.cell(0, 0).text == "A"
        assert table.cell(0, 1).text == "B"
        assert table.cell(0, 2).text == "C"
        assert table.cell(1, 0).text == "D"
        assert table.cell(1, 1).text == "E"
        assert table.cell(1, 2).text == "F"

    def test_header_cell_background_can_be_set(self, tmp_path: Path) -> None:
        """Header cells accept a background color without error."""
        header_cell = TableCell(
            paragraphs=[Paragraph(runs=[TextRun(text="Title")])],
            is_header=True,
            background=RGBAColor(r=50, g=50, b=200),
        )
        body_cell = TableCell(
            paragraphs=[Paragraph(runs=[TextRun(text="Data")])],
        )
        el = TableElement(
            element_type=ElementType.TABLE,
            box=Box(x=50, y=100, width=600, height=200),
            table_rows=[
                TableRow(cells=[header_cell]),
                TableRow(cells=[body_cell]),
            ],
        )
        pres = _make_slide_with(el)
        pptx = _build_and_read(pres, tmp_path)
        table = _first_table(pptx)
        assert table.cell(0, 0).text == "Title"


# ---------------------------------------------------------------------------
# TestColspanRowspan
# ---------------------------------------------------------------------------


class TestColspanRowspan:
    """Verify colspan/rowspan merge operations."""

    def test_colspan_merge_spans_cells(self, tmp_path: Path) -> None:
        """A cell with colspan=2 should merge two adjacent columns."""
        merged_cell = TableCell(
            paragraphs=[Paragraph(runs=[TextRun(text="Merged")])],
            colspan=2,
            rowspan=1,
        )
        left = TableCell(paragraphs=[Paragraph(runs=[TextRun(text="Left")])])
        right = TableCell(paragraphs=[Paragraph(runs=[TextRun(text="Right")])])
        el = TableElement(
            element_type=ElementType.TABLE,
            box=Box(x=50, y=100, width=600, height=200),
            table_rows=[
                TableRow(cells=[merged_cell]),
                TableRow(cells=[left, right]),
            ],
        )
        pres = _make_slide_with(el)
        pptx = _build_and_read(pres, tmp_path)
        table = _first_table(pptx)
        # The merged top cell should contain the expected text
        assert table.cell(0, 0).text == "Merged"

    def test_colspan_content_set_after_merge(self, tmp_path: Path) -> None:
        """Content in colspan cells must survive the merge operation."""
        merged_cell = TableCell(
            paragraphs=[Paragraph(runs=[TextRun(text="SpannedContent")])],
            colspan=3,
        )
        data_cells = [
            TableCell(paragraphs=[Paragraph(runs=[TextRun(text=str(i))])])
            for i in range(3)
        ]
        el = TableElement(
            element_type=ElementType.TABLE,
            box=Box(x=50, y=100, width=700, height=200),
            table_rows=[
                TableRow(cells=[merged_cell]),
                TableRow(cells=data_cells),
            ],
        )
        pres = _make_slide_with(el)
        pptx = _build_and_read(pres, tmp_path)
        table = _first_table(pptx)
        assert table.cell(0, 0).text == "SpannedContent"

    def test_rowspan_merge_spans_rows(self, tmp_path: Path) -> None:
        """A cell with rowspan=2 should span two rows in the grid."""
        tall_cell = TableCell(
            paragraphs=[Paragraph(runs=[TextRun(text="Tall")])],
            rowspan=2,
            colspan=1,
        )
        top_right = TableCell(paragraphs=[Paragraph(runs=[TextRun(text="TR")])])
        bottom_right = TableCell(paragraphs=[Paragraph(runs=[TextRun(text="BR")])])
        el = TableElement(
            element_type=ElementType.TABLE,
            box=Box(x=50, y=100, width=400, height=300),
            table_rows=[
                TableRow(cells=[tall_cell, top_right]),
                TableRow(cells=[bottom_right]),
            ],
        )
        pres = _make_slide_with(el)
        pptx = _build_and_read(pres, tmp_path)
        table = _first_table(pptx)
        assert table.cell(0, 0).text == "Tall"
        assert table.cell(0, 1).text == "TR"

    def test_rowspan_content_set_after_merge(self, tmp_path: Path) -> None:
        """Content in rowspan cells must survive the merge operation."""
        tall_cell = TableCell(
            paragraphs=[Paragraph(runs=[TextRun(text="TallContent")])],
            rowspan=2,
        )
        side_top = TableCell(paragraphs=[Paragraph(runs=[TextRun(text="S1")])])
        side_bottom = TableCell(paragraphs=[Paragraph(runs=[TextRun(text="S2")])])
        el = TableElement(
            element_type=ElementType.TABLE,
            box=Box(x=50, y=100, width=400, height=300),
            table_rows=[
                TableRow(cells=[tall_cell, side_top]),
                TableRow(cells=[side_bottom]),
            ],
        )
        pres = _make_slide_with(el)
        pptx = _build_and_read(pres, tmp_path)
        table = _first_table(pptx)
        assert table.cell(0, 0).text == "TallContent"


# ---------------------------------------------------------------------------
# TestTableDimensions
# ---------------------------------------------------------------------------


class TestTableDimensions:
    """Verify row count and column count in the generated table."""

    def test_correct_row_count(self, tmp_path: Path) -> None:
        pres = _make_slide_with(_make_simple_table([["H"], ["R1"], ["R2"], ["R3"]]))
        pptx = _build_and_read(pres, tmp_path)
        table = _first_table(pptx)
        assert len(table.rows) == 4

    def test_correct_column_count(self, tmp_path: Path) -> None:
        pres = _make_slide_with(_make_simple_table([["A", "B", "C"], ["D", "E", "F"]]))
        pptx = _build_and_read(pres, tmp_path)
        table = _first_table(pptx)
        assert len(table.columns) == 3

    def test_single_cell_table(self, tmp_path: Path) -> None:
        pres = _make_slide_with(_make_simple_table([["Only"]]))
        pptx = _build_and_read(pres, tmp_path)
        table = _first_table(pptx)
        assert len(table.rows) == 1
        assert len(table.columns) == 1
        assert table.cell(0, 0).text == "Only"

    def test_colspan_expands_column_count(self, tmp_path: Path) -> None:
        """A colspan=3 header should create a 3-column table."""
        header = TableCell(
            paragraphs=[Paragraph(runs=[TextRun(text="Header")])],
            colspan=3,
        )
        row1 = [
            TableCell(paragraphs=[Paragraph(runs=[TextRun(text=str(i))])])
            for i in range(3)
        ]
        el = TableElement(
            element_type=ElementType.TABLE,
            box=Box(x=50, y=100, width=700, height=200),
            table_rows=[
                TableRow(cells=[header]),
                TableRow(cells=row1),
            ],
        )
        pres = _make_slide_with(el)
        pptx = _build_and_read(pres, tmp_path)
        table = _first_table(pptx)
        assert len(table.columns) == 3
