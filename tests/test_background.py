"""Tests for background image support.

These tests construct model objects directly -- NO Playwright or marp-cli needed.
"""

from __future__ import annotations

import base64
import io
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation as PptxPresentation
from pptx.util import Emu

from marpx.models import (
    Background,
    BackgroundImage,
    Box,
    ElementType,
    Paragraph,
    Presentation,
    RGBAColor,
    Slide,
    SlideElement,
    TextRun,
    TextStyle,
)
from marpx.pptx_builder.background import _add_background_image
from marpx.pptx_builder.builder import build_pptx
from marpx.utils import px_to_emu


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _make_heading(text: str) -> SlideElement:
    return SlideElement(
        element_type=ElementType.HEADING,
        box=Box(x=50, y=20, width=600, height=60),
        heading_level=1,
        paragraphs=[
            Paragraph(
                runs=[TextRun(text=text, style=TextStyle(bold=True, font_size_px=32))],
                alignment="left",
            )
        ],
    )


def _build_and_read(presentation: Presentation, tmp_path: Path) -> PptxPresentation:
    """Build PPTX and re-open it for assertions."""
    out = tmp_path / "test_output.pptx"
    build_pptx(presentation, out)
    return PptxPresentation(str(out))


def _create_tiny_png() -> bytes:
    """Create a minimal valid 1x1 red PNG image for testing."""
    # Minimal 1x1 red PNG (67 bytes)
    return base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR4"
        "nGP4z8BQDwAEgAF/pooBPQAAAABJRU5ErkJggg=="
    )


def _create_png(
    width: int, height: int, color: tuple[int, int, int] = (255, 0, 0)
) -> bytes:
    """Create a simple PNG of the requested size."""
    image = Image.new("RGB", (width, height), color)
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()


# ---------------------------------------------------------------------------
# BackgroundImage model tests
# ---------------------------------------------------------------------------


class TestBackgroundImageModel:
    """Tests for BackgroundImage model creation and defaults."""

    def test_minimal_creation(self) -> None:
        bg = BackgroundImage(url="https://example.com/img.png")
        assert bg.url == "https://example.com/img.png"
        assert bg.size == "cover"
        assert bg.position == "center"
        assert bg.split is None
        assert bg.image_data is None

    def test_contain_size(self) -> None:
        bg = BackgroundImage(url="test.png", size="contain")
        assert bg.size == "contain"

    def test_split_left(self) -> None:
        bg = BackgroundImage(url="test.png", split="left")
        assert bg.split == "left"
        assert bg.split_ratio is None

    def test_split_right(self) -> None:
        bg = BackgroundImage(url="test.png", split="right")
        assert bg.split == "right"

    def test_split_ratio(self) -> None:
        bg = BackgroundImage(url="test.png", split="left", split_ratio=0.4)
        assert bg.split_ratio == pytest.approx(0.4)

    def test_with_image_data(self) -> None:
        data = b"\x89PNG\r\n"
        bg = BackgroundImage(url="test.png", image_data=data)
        assert bg.image_data == data

    def test_custom_position(self) -> None:
        bg = BackgroundImage(url="test.png", position="top left")
        assert bg.position == "top left"


# ---------------------------------------------------------------------------
# Background model tests (with images list)
# ---------------------------------------------------------------------------


class TestBackgroundWithImages:
    """Tests for Background model with the images field."""

    def test_default_empty_images(self) -> None:
        bg = Background()
        assert bg.images == []
        assert bg.color is None
        assert bg.image_path is None

    def test_images_list(self) -> None:
        bg = Background(
            images=[
                BackgroundImage(url="img1.png"),
                BackgroundImage(url="img2.png"),
            ]
        )
        assert len(bg.images) == 2

    def test_color_and_images(self) -> None:
        bg = Background(
            color=RGBAColor(r=30, g=30, b=30),
            images=[BackgroundImage(url="overlay.png")],
        )
        assert bg.color is not None
        assert len(bg.images) == 1

    def test_backward_compatible(self) -> None:
        """Existing code that creates Background(color=...) still works."""
        bg = Background(color=RGBAColor(r=255, g=255, b=255))
        assert bg.images == []


# ---------------------------------------------------------------------------
# _add_background_image tests (PPTX building)
# ---------------------------------------------------------------------------


class TestAddBackgroundImage:
    """Tests for _add_background_image building PPTX shapes."""

    def test_adds_picture_shape_from_image_data(self, tmp_path: Path) -> None:
        """Background image with inline image_data adds a picture shape."""
        png_data = _create_png(1280, 720)
        bg_img = BackgroundImage(url="inline.png", image_data=png_data)

        pptx = PptxPresentation()
        layout = pptx.slide_layouts[6]
        pptx_slide = pptx.slides.add_slide(layout)

        slide_w = px_to_emu(1280)
        slide_h = px_to_emu(720)
        _add_background_image(pptx_slide, bg_img, slide_w, slide_h)

        # Should have exactly 1 shape (the picture)
        assert len(pptx_slide.shapes) == 1
        pic = pptx_slide.shapes[0]
        assert pic.left == 0
        assert pic.top == 0
        assert pic.width == Emu(slide_w)
        assert pic.height == Emu(slide_h)

    def test_split_left_positioning(self, tmp_path: Path) -> None:
        """Split left positions image in left half."""
        png_data = _create_png(640, 720)
        bg_img = BackgroundImage(url="bg.png", image_data=png_data, split="left")

        pptx = PptxPresentation()
        layout = pptx.slide_layouts[6]
        pptx_slide = pptx.slides.add_slide(layout)

        slide_w = px_to_emu(1280)
        slide_h = px_to_emu(720)
        _add_background_image(pptx_slide, bg_img, slide_w, slide_h)

        pic = pptx_slide.shapes[0]
        assert pic.left == 0
        assert pic.top == 0
        assert pic.width == Emu(slide_w // 2)
        assert pic.height == Emu(slide_h)

    def test_split_left_uses_custom_ratio(self, tmp_path: Path) -> None:
        """Split-left respects Marp's explicit split ratio."""
        png_data = _create_png(640, 720)
        bg_img = BackgroundImage(
            url="bg.png",
            image_data=png_data,
            split="left",
            split_ratio=0.4,
        )

        pptx = PptxPresentation()
        layout = pptx.slide_layouts[6]
        pptx_slide = pptx.slides.add_slide(layout)

        slide_w = px_to_emu(1280)
        slide_h = px_to_emu(720)
        _add_background_image(pptx_slide, bg_img, slide_w, slide_h)

        pic = pptx_slide.shapes[0]
        assert pic.left == 0
        assert pic.width == Emu(round(slide_w * 0.4))

    def test_split_right_positioning(self, tmp_path: Path) -> None:
        """Split right positions image in right half."""
        png_data = _create_png(640, 720)
        bg_img = BackgroundImage(url="bg.png", image_data=png_data, split="right")

        pptx = PptxPresentation()
        layout = pptx.slide_layouts[6]
        pptx_slide = pptx.slides.add_slide(layout)

        slide_w = px_to_emu(1280)
        slide_h = px_to_emu(720)
        _add_background_image(pptx_slide, bg_img, slide_w, slide_h)

        pic = pptx_slide.shapes[0]
        assert pic.left == Emu(slide_w // 2)
        assert pic.top == 0
        assert pic.width == Emu(slide_w // 2)
        assert pic.height == Emu(slide_h)

    def test_split_right_uses_custom_ratio(self, tmp_path: Path) -> None:
        """Split-right respects Marp's explicit split ratio."""
        png_data = _create_png(640, 720)
        bg_img = BackgroundImage(
            url="bg.png",
            image_data=png_data,
            split="right",
            split_ratio=0.4,
        )

        pptx = PptxPresentation()
        layout = pptx.slide_layouts[6]
        pptx_slide = pptx.slides.add_slide(layout)

        slide_w = px_to_emu(1280)
        slide_h = px_to_emu(720)
        _add_background_image(pptx_slide, bg_img, slide_w, slide_h)

        pic = pptx_slide.shapes[0]
        expected_width = round(slide_w * 0.4)
        assert pic.left == Emu(slide_w - expected_width)
        assert pic.width == Emu(expected_width)

    def test_split_right_cover_crops_to_right_half(self, tmp_path: Path) -> None:
        """Wide split-right backgrounds should crop inside the right-half box."""
        png_data = _create_png(1040, 380)
        bg_img = BackgroundImage(url="bg.png", image_data=png_data, split="right")

        pptx = PptxPresentation()
        layout = pptx.slide_layouts[6]
        pptx_slide = pptx.slides.add_slide(layout)

        slide_w = px_to_emu(1280)
        slide_h = px_to_emu(720)
        _add_background_image(pptx_slide, bg_img, slide_w, slide_h)

        pic = pptx_slide.shapes[0]
        assert pic.left == Emu(slide_w // 2)
        assert pic.width == Emu(slide_w // 2)
        assert pic.crop_left > 0
        assert pic.crop_right > 0

    def test_data_uri_image(self, tmp_path: Path) -> None:
        """Background image from a data: URI."""
        png_data = _create_tiny_png()
        b64 = base64.b64encode(png_data).decode()
        data_uri = f"data:image/png;base64,{b64}"
        bg_img = BackgroundImage(url=data_uri)

        pptx = PptxPresentation()
        layout = pptx.slide_layouts[6]
        pptx_slide = pptx.slides.add_slide(layout)

        slide_w = px_to_emu(1280)
        slide_h = px_to_emu(720)
        _add_background_image(pptx_slide, bg_img, slide_w, slide_h)

        assert len(pptx_slide.shapes) == 1

    def test_local_file_image(self, tmp_path: Path) -> None:
        """Background image from a local file path."""
        png_data = _create_tiny_png()
        img_file = tmp_path / "bg_image.png"
        img_file.write_bytes(png_data)

        bg_img = BackgroundImage(url=str(img_file))

        pptx = PptxPresentation()
        layout = pptx.slide_layouts[6]
        pptx_slide = pptx.slides.add_slide(layout)

        slide_w = px_to_emu(1280)
        slide_h = px_to_emu(720)
        _add_background_image(pptx_slide, bg_img, slide_w, slide_h)

        assert len(pptx_slide.shapes) == 1

    def test_nonexistent_file_no_shape(self, tmp_path: Path) -> None:
        """A non-existent local file path should not add any shape."""
        bg_img = BackgroundImage(url="/nonexistent/path/to/image.png")

        pptx = PptxPresentation()
        layout = pptx.slide_layouts[6]
        pptx_slide = pptx.slides.add_slide(layout)

        slide_w = px_to_emu(1280)
        slide_h = px_to_emu(720)
        _add_background_image(pptx_slide, bg_img, slide_w, slide_h)

        assert len(pptx_slide.shapes) == 0

    def test_file_uri_image(self, tmp_path: Path) -> None:
        """Background image from a file:// URI."""
        png_data = _create_tiny_png()
        img_file = tmp_path / "file_uri_bg.png"
        img_file.write_bytes(png_data)

        bg_img = BackgroundImage(url=f"file://{img_file}")

        pptx = PptxPresentation()
        layout = pptx.slide_layouts[6]
        pptx_slide = pptx.slides.add_slide(layout)

        slide_w = px_to_emu(1280)
        slide_h = px_to_emu(720)
        _add_background_image(pptx_slide, bg_img, slide_w, slide_h)

        assert len(pptx_slide.shapes) == 1

    def test_file_uri_with_spaces_image(self, tmp_path: Path) -> None:
        """Background image from a URL-escaped file:// URI."""
        png_data = _create_tiny_png()
        img_file = tmp_path / "file uri bg.png"
        img_file.write_bytes(png_data)

        bg_img = BackgroundImage(url=img_file.as_uri())

        pptx = PptxPresentation()
        layout = pptx.slide_layouts[6]
        pptx_slide = pptx.slides.add_slide(layout)

        slide_w = px_to_emu(1280)
        slide_h = px_to_emu(720)
        _add_background_image(pptx_slide, bg_img, slide_w, slide_h)

        assert len(pptx_slide.shapes) == 1

    def test_contain_respects_aspect_ratio(self) -> None:
        """Contain should preserve image aspect ratio within slide bounds."""
        png_data = _create_png(400, 200)
        bg_img = BackgroundImage(url="bg.png", image_data=png_data, size="contain")

        pptx = PptxPresentation()
        layout = pptx.slide_layouts[6]
        pptx_slide = pptx.slides.add_slide(layout)

        slide_w = px_to_emu(1280)
        slide_h = px_to_emu(720)
        _add_background_image(pptx_slide, bg_img, slide_w, slide_h)

        pic = pptx_slide.shapes[0]
        assert pic.width == Emu(slide_w)
        assert pic.height == Emu(slide_w // 2)
        assert pic.top == Emu((slide_h - (slide_w // 2)) // 2)

    def test_cover_honors_background_position(self) -> None:
        """Cover should anchor overflow according to background-position."""
        png_data = _create_png(400, 200)
        bg_img = BackgroundImage(
            url="bg.png",
            image_data=png_data,
            size="cover",
            position="top left",
        )

        pptx = PptxPresentation()
        layout = pptx.slide_layouts[6]
        pptx_slide = pptx.slides.add_slide(layout)

        slide_w = px_to_emu(1280)
        slide_h = px_to_emu(720)
        _add_background_image(pptx_slide, bg_img, slide_w, slide_h)

        pic = pptx_slide.shapes[0]
        assert pic.left == 0
        assert pic.top == 0
        assert pic.width == Emu(slide_w)
        assert pic.height == Emu(slide_h)
        assert pic.crop_left == 0
        assert pic.crop_top == 0
        assert pic.crop_right > 0

    def test_cover_allows_right_aligned_crop(self) -> None:
        """Cover should offset oversized images when anchored right."""
        png_data = _create_png(400, 200)
        bg_img = BackgroundImage(
            url="bg.png",
            image_data=png_data,
            size="cover",
            position="right center",
        )

        pptx = PptxPresentation()
        layout = pptx.slide_layouts[6]
        pptx_slide = pptx.slides.add_slide(layout)

        slide_w = px_to_emu(1280)
        slide_h = px_to_emu(720)
        _add_background_image(pptx_slide, bg_img, slide_w, slide_h)

        pic = pptx_slide.shapes[0]
        assert pic.left == 0
        assert pic.top == 0
        assert pic.width == Emu(slide_w)
        assert pic.crop_left > 0
        assert pic.crop_right == 0


# ---------------------------------------------------------------------------
# Integration: build_pptx with background images
# ---------------------------------------------------------------------------


class TestBuildPptxWithBackgroundImages:
    """End-to-end tests for build_pptx with background images."""

    def test_slide_with_bg_image_and_text(self, tmp_path: Path) -> None:
        """A slide with both a background image and text content."""
        png_data = _create_tiny_png()
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[_make_heading("Over Image")],
                    background=Background(
                        images=[
                            BackgroundImage(url="inline.png", image_data=png_data),
                        ],
                    ),
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]

        # Should have 2 shapes: 1 background image + 1 textbox
        assert len(slide.shapes) == 2

        # Verify text content exists
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    texts.append(para.text)
        assert any("Over Image" in t for t in texts)

    def test_slide_with_color_and_image(self, tmp_path: Path) -> None:
        """Solid color and image background can coexist."""
        png_data = _create_tiny_png()
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[_make_heading("Combined")],
                    background=Background(
                        color=RGBAColor(r=30, g=30, b=30, a=1.0),
                        images=[
                            BackgroundImage(url="bg.png", image_data=png_data),
                        ],
                    ),
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]

        # Should have solid bg fill + 1 picture + 1 textbox = 2 shapes
        assert len(slide.shapes) == 2

        # The slide background fill should also be set
        bg = slide.background
        fill = bg.fill
        assert fill.type is not None

    def test_multiple_bg_images(self, tmp_path: Path) -> None:
        """Multiple background images on a single slide."""
        png_data = _create_tiny_png()
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[],
                    background=Background(
                        images=[
                            BackgroundImage(url="bg1.png", image_data=png_data),
                            BackgroundImage(url="bg2.png", image_data=png_data),
                        ],
                    ),
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]

        # Should have 2 picture shapes
        assert len(slide.shapes) == 2

    def test_bg_image_z_order_behind_text(self, tmp_path: Path) -> None:
        """Background image shape should be behind text shapes in spTree."""
        png_data = _create_tiny_png()
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[_make_heading("Foreground")],
                    background=Background(
                        images=[
                            BackgroundImage(url="bg.png", image_data=png_data),
                        ],
                    ),
                )
            ],
        )
        out = tmp_path / "zorder.pptx"
        build_pptx(pres, out)
        pptx = PptxPresentation(str(out))
        slide = pptx.slides[0]

        # The spTree children: [0] = grpSpPr, [1] = grpSp or similar,
        # then shapes. The background picture should come before the textbox.
        sp_tree = slide.shapes._spTree
        shape_elements = [
            child
            for child in sp_tree
            if child.tag.endswith("}pic") or child.tag.endswith("}sp")
        ]
        # First shape element should be the picture (background)
        assert len(shape_elements) >= 2
        assert shape_elements[0].tag.endswith("}pic")

    def test_split_left_bg_image_build(self, tmp_path: Path) -> None:
        """Split-left background image in full build pipeline."""
        png_data = _create_tiny_png()
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[_make_heading("Split Slide")],
                    background=Background(
                        images=[
                            BackgroundImage(
                                url="bg.png",
                                image_data=png_data,
                                split="left",
                            ),
                        ],
                    ),
                )
            ],
        )
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        assert len(slide.shapes) == 2

    def test_failed_bg_image_logged_not_crash(self, tmp_path: Path) -> None:
        """A background image that can't be loaded should log warning, not crash."""
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    elements=[_make_heading("Still works")],
                    background=Background(
                        images=[
                            BackgroundImage(url="/nonexistent/image.png"),
                        ],
                    ),
                )
            ],
        )
        # Should not raise
        pptx = _build_and_read(pres, tmp_path)
        slide = pptx.slides[0]
        # Only the textbox shape, no picture
        assert len(slide.shapes) == 1
