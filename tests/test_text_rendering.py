"""Direct unit tests for text rendering in pptx_builder/text.py.

These tests construct model objects directly -- NO Playwright or marp-cli needed.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation as PptxPresentation
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn

import math

from marpx.models import (
    Box,
    ElementType,
    ListElement,
    ListItem,
    Paragraph,
    Presentation,
    RGBAColor,
    Slide,
    SlideElement,
    TextElement,
    TextRun,
    TextShadow,
    TextStyle,
)
from marpx.pptx_builder.builder import build_pptx
from marpx.utils.common import px_to_emu


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _make_slide_with(*elements: SlideElement) -> Presentation:
    return Presentation(
        slides=[Slide(width_px=1280, height_px=720, elements=list(elements))]
    )


def _make_text_element(
    text: str,
    alignment: str = "left",
    vertical_align: str = "top",
    element_type: ElementType = ElementType.PARAGRAPH,
    style: TextStyle | None = None,
) -> SlideElement:
    run_style = style if style is not None else TextStyle()
    return TextElement(
        element_type=element_type,
        box=Box(x=50, y=50, width=600, height=200),
        vertical_align=vertical_align,
        paragraphs=[
            Paragraph(
                runs=[TextRun(text=text, style=run_style)],
                alignment=alignment,
            )
        ],
    )


def _make_unordered_list(items: list[tuple[str, int]]) -> SlideElement:
    return ListElement(
        element_type=ElementType.UNORDERED_LIST,
        box=Box(x=50, y=50, width=600, height=300),
        list_items=[
            ListItem(runs=[TextRun(text=text)], level=level) for text, level in items
        ],
    )


def _make_ordered_list(items: list[tuple[str, int]]) -> SlideElement:
    return ListElement(
        element_type=ElementType.ORDERED_LIST,
        box=Box(x=50, y=50, width=600, height=300),
        list_items=[
            ListItem(
                runs=[TextRun(text=text)],
                level=level,
                order_number=idx + 1,
            )
            for idx, (text, level) in enumerate(items)
        ],
    )


def _build_and_read(presentation: Presentation, tmp_path: Path) -> PptxPresentation:
    """Build PPTX and re-open it for assertions."""
    out = tmp_path / "test_output.pptx"
    build_pptx(presentation, out)
    return PptxPresentation(str(out))


def _first_textbox(pptx: PptxPresentation):
    """Return the first textbox shape on slide 0."""
    slide = pptx.slides[0]
    for shape in slide.shapes:
        if shape.has_text_frame:
            return shape
    raise AssertionError("No textbox found on slide")


def _first_paragraph(pptx: PptxPresentation):
    return _first_textbox(pptx).text_frame.paragraphs[0]


def _first_run(pptx: PptxPresentation):
    return _first_paragraph(pptx).runs[0]


# ---------------------------------------------------------------------------
# TestTextAlignment
# ---------------------------------------------------------------------------


class TestTextAlignment:
    """Verify left/center/right/justify paragraph alignment."""

    def test_left_alignment(self, tmp_path: Path) -> None:
        pres = _make_slide_with(_make_text_element("Hello", alignment="left"))
        pptx = _build_and_read(pres, tmp_path)
        para = _first_paragraph(pptx)
        assert para.alignment == PP_ALIGN.LEFT

    def test_center_alignment(self, tmp_path: Path) -> None:
        pres = _make_slide_with(_make_text_element("Hello", alignment="center"))
        pptx = _build_and_read(pres, tmp_path)
        para = _first_paragraph(pptx)
        assert para.alignment == PP_ALIGN.CENTER

    def test_right_alignment(self, tmp_path: Path) -> None:
        pres = _make_slide_with(_make_text_element("Hello", alignment="right"))
        pptx = _build_and_read(pres, tmp_path)
        para = _first_paragraph(pptx)
        assert para.alignment == PP_ALIGN.RIGHT

    def test_justify_alignment(self, tmp_path: Path) -> None:
        pres = _make_slide_with(_make_text_element("Hello", alignment="justify"))
        pptx = _build_and_read(pres, tmp_path)
        para = _first_paragraph(pptx)
        assert para.alignment == PP_ALIGN.JUSTIFY

    def test_start_maps_to_left(self, tmp_path: Path) -> None:
        pres = _make_slide_with(_make_text_element("Hello", alignment="start"))
        pptx = _build_and_read(pres, tmp_path)
        para = _first_paragraph(pptx)
        assert para.alignment == PP_ALIGN.LEFT

    def test_end_maps_to_right(self, tmp_path: Path) -> None:
        pres = _make_slide_with(_make_text_element("Hello", alignment="end"))
        pptx = _build_and_read(pres, tmp_path)
        para = _first_paragraph(pptx)
        assert para.alignment == PP_ALIGN.RIGHT


# ---------------------------------------------------------------------------
# TestTextStyle
# ---------------------------------------------------------------------------


class TestTextStyle:
    """Verify text run styling: bold, italic, underline, strike, font_size, color."""

    def test_bold(self, tmp_path: Path) -> None:
        style = TextStyle(bold=True)
        pres = _make_slide_with(_make_text_element("Bold", style=style))
        pptx = _build_and_read(pres, tmp_path)
        run = _first_run(pptx)
        assert run.font.bold is True

    def test_not_bold_by_default(self, tmp_path: Path) -> None:
        pres = _make_slide_with(_make_text_element("Normal"))
        pptx = _build_and_read(pres, tmp_path)
        run = _first_run(pptx)
        assert not run.font.bold

    def test_italic(self, tmp_path: Path) -> None:
        style = TextStyle(italic=True)
        pres = _make_slide_with(_make_text_element("Italic", style=style))
        pptx = _build_and_read(pres, tmp_path)
        run = _first_run(pptx)
        assert run.font.italic is True

    def test_underline(self, tmp_path: Path) -> None:
        style = TextStyle(underline=True)
        pres = _make_slide_with(_make_text_element("Underline", style=style))
        pptx = _build_and_read(pres, tmp_path)
        run = _first_run(pptx)
        assert run.font.underline is True

    def test_strikethrough(self, tmp_path: Path) -> None:
        style = TextStyle(strike=True)
        pres = _make_slide_with(_make_text_element("Strike", style=style))
        pptx = _build_and_read(pres, tmp_path)
        run = _first_run(pptx)
        r_pr = run._r.get_or_add_rPr()
        assert r_pr.get("strike") == "sngStrike"

    def test_font_size(self, tmp_path: Path) -> None:
        style = TextStyle(font_size_px=32.0)
        pres = _make_slide_with(_make_text_element("Big", style=style))
        pptx = _build_and_read(pres, tmp_path)
        run = _first_run(pptx)
        # font.size is in EMU (hundredths of a point); just verify it is set
        assert run.font.size is not None
        assert run.font.size > 0

    def test_color_opaque(self, tmp_path: Path) -> None:
        style = TextStyle(color=RGBAColor(r=255, g=0, b=0))
        pres = _make_slide_with(_make_text_element("Red", style=style))
        pptx = _build_and_read(pres, tmp_path)
        run = _first_run(pptx)
        rgb = run.font.color.rgb
        assert rgb[0] == 255
        assert rgb[1] == 0
        assert rgb[2] == 0

    def test_color_with_alpha(self, tmp_path: Path) -> None:
        """Semi-transparent color should embed alpha in XML."""
        style = TextStyle(color=RGBAColor(r=0, g=128, b=255, a=0.5))
        pres = _make_slide_with(_make_text_element("Alpha", style=style))
        pptx = _build_and_read(pres, tmp_path)
        run = _first_run(pptx)
        r_pr = run._r.get_or_add_rPr()
        solid_fill = r_pr.find(qn("a:solidFill"))
        assert solid_fill is not None
        srgb = solid_fill.find(qn("a:srgbClr"))
        assert srgb is not None
        alpha_el = srgb.find(qn("a:alpha"))
        assert alpha_el is not None

    def test_text_content_preserved(self, tmp_path: Path) -> None:
        pres = _make_slide_with(_make_text_element("Hello World"))
        pptx = _build_and_read(pres, tmp_path)
        run = _first_run(pptx)
        assert run.text == "Hello World"


# ---------------------------------------------------------------------------
# TestVerticalAlignment
# ---------------------------------------------------------------------------


class TestVerticalAlignment:
    """Verify top/middle/bottom vertical anchor on textboxes."""

    def test_top_vertical_alignment(self, tmp_path: Path) -> None:
        pres = _make_slide_with(_make_text_element("Top", vertical_align="top"))
        pptx = _build_and_read(pres, tmp_path)
        tf = _first_textbox(pptx).text_frame
        assert tf.vertical_anchor == MSO_VERTICAL_ANCHOR.TOP

    def test_middle_vertical_alignment(self, tmp_path: Path) -> None:
        pres = _make_slide_with(_make_text_element("Middle", vertical_align="middle"))
        pptx = _build_and_read(pres, tmp_path)
        tf = _first_textbox(pptx).text_frame
        assert tf.vertical_anchor == MSO_VERTICAL_ANCHOR.MIDDLE

    def test_bottom_vertical_alignment(self, tmp_path: Path) -> None:
        pres = _make_slide_with(_make_text_element("Bottom", vertical_align="bottom"))
        pptx = _build_and_read(pres, tmp_path)
        tf = _first_textbox(pptx).text_frame
        assert tf.vertical_anchor == MSO_VERTICAL_ANCHOR.BOTTOM

    def test_unknown_vertical_alignment_defaults_to_top(self, tmp_path: Path) -> None:
        pres = _make_slide_with(_make_text_element("X", vertical_align="baseline"))
        pptx = _build_and_read(pres, tmp_path)
        tf = _first_textbox(pptx).text_frame
        assert tf.vertical_anchor == MSO_VERTICAL_ANCHOR.TOP


# ---------------------------------------------------------------------------
# TestListRendering
# ---------------------------------------------------------------------------


class TestListRendering:
    """Verify bullet characters, ordered numbers, and nested indent levels."""

    def test_unordered_list_bullet_char(self, tmp_path: Path) -> None:
        el = _make_unordered_list([("Item A", 0)])
        pres = _make_slide_with(el)
        pptx = _build_and_read(pres, tmp_path)
        tf = _first_textbox(pptx).text_frame
        para = tf.paragraphs[0]
        pPr = para._p.pPr
        assert pPr is not None
        bu_char = pPr.find(qn("a:buChar"))
        assert bu_char is not None
        # Default bullet is U+2022 (bullet dot)
        assert bu_char.get("char") == "\u2022"

    def test_ordered_list_auto_numbering(self, tmp_path: Path) -> None:
        el = _make_ordered_list([("First", 0), ("Second", 0)])
        pres = _make_slide_with(el)
        pptx = _build_and_read(pres, tmp_path)
        tf = _first_textbox(pptx).text_frame
        para = tf.paragraphs[0]
        pPr = para._p.pPr
        assert pPr is not None
        bu_auto = pPr.find(qn("a:buAutoNum"))
        assert bu_auto is not None

    def test_ordered_list_numbering_type(self, tmp_path: Path) -> None:
        el = _make_ordered_list([("A", 0)])
        pres = _make_slide_with(el)
        pptx = _build_and_read(pres, tmp_path)
        tf = _first_textbox(pptx).text_frame
        para = tf.paragraphs[0]
        pPr = para._p.pPr
        bu_auto = pPr.find(qn("a:buAutoNum"))
        # Default decimal ordering maps to "arabicPeriod"
        assert bu_auto.get("type") == "arabicPeriod"

    def test_nested_list_level_0(self, tmp_path: Path) -> None:
        el = _make_unordered_list([("Top", 0)])
        pres = _make_slide_with(el)
        pptx = _build_and_read(pres, tmp_path)
        tf = _first_textbox(pptx).text_frame
        assert tf.paragraphs[0].level == 0

    def test_nested_list_level_1(self, tmp_path: Path) -> None:
        el = _make_unordered_list([("Parent", 0), ("Child", 1)])
        pres = _make_slide_with(el)
        pptx = _build_and_read(pres, tmp_path)
        tf = _first_textbox(pptx).text_frame
        assert tf.paragraphs[1].level == 1

    def test_nested_list_multiple_levels(self, tmp_path: Path) -> None:
        el = _make_unordered_list([("L0", 0), ("L1", 1), ("L2", 2)])
        pres = _make_slide_with(el)
        pptx = _build_and_read(pres, tmp_path)
        tf = _first_textbox(pptx).text_frame
        assert tf.paragraphs[0].level == 0
        assert tf.paragraphs[1].level == 1
        assert tf.paragraphs[2].level == 2

    def test_list_item_text_content(self, tmp_path: Path) -> None:
        el = _make_unordered_list([("Alpha", 0), ("Beta", 0)])
        pres = _make_slide_with(el)
        pptx = _build_and_read(pres, tmp_path)
        tf = _first_textbox(pptx).text_frame
        texts = [p.runs[0].text for p in tf.paragraphs if p.runs]
        assert "Alpha" in texts
        assert "Beta" in texts


# ---------------------------------------------------------------------------
# TestTextShadow
# ---------------------------------------------------------------------------


def _get_effect_lst(pptx: PptxPresentation):
    """Return the a:effectLst element from the first run's rPr, or None."""
    run = _first_run(pptx)
    r_pr = run._r.find(qn("a:rPr"))
    if r_pr is None:
        return None
    return r_pr.find(qn("a:effectLst"))


class TestTextShadow:
    """Verify text-shadow CSS renders as PPTX glow / outerShdw effects."""

    def test_text_shadow_glow_effect(self, tmp_path: Path) -> None:
        """Offset=0 shadow produces an a:glow element with correct rad and color."""
        shadow = TextShadow(
            offset_x_px=0.0,
            offset_y_px=0.0,
            blur_radius_px=10.0,
            color=RGBAColor(r=0, g=255, b=255),
        )
        style = TextStyle(text_shadows=[shadow])
        pres = _make_slide_with(_make_text_element("Glow", style=style))
        pptx = _build_and_read(pres, tmp_path)

        effect_lst = _get_effect_lst(pptx)
        assert effect_lst is not None, "effectLst not found in rPr"

        glow = effect_lst.find(qn("a:glow"))
        assert glow is not None, "a:glow element not found"

        expected_rad = str(px_to_emu(10.0))
        assert glow.get("rad") == expected_rad

        srgb = glow.find(qn("a:srgbClr"))
        assert srgb is not None
        assert srgb.get("val") == "00FFFF"

    def test_text_shadow_outer_shadow(self, tmp_path: Path) -> None:
        """Offset > 1px shadow produces an a:outerShdw with correct attributes."""
        shadow = TextShadow(
            offset_x_px=3.0,
            offset_y_px=3.0,
            blur_radius_px=5.0,
            color=RGBAColor(r=0, g=0, b=0),
        )
        style = TextStyle(text_shadows=[shadow])
        pres = _make_slide_with(_make_text_element("Shadow", style=style))
        pptx = _build_and_read(pres, tmp_path)

        effect_lst = _get_effect_lst(pptx)
        assert effect_lst is not None

        outer = effect_lst.find(qn("a:outerShdw"))
        assert outer is not None, "a:outerShdw element not found"

        expected_blur = str(px_to_emu(5.0))
        expected_dist = str(px_to_emu(math.hypot(3.0, 3.0)))
        angle_rad = math.atan2(3.0, 3.0)
        expected_dir = str(int(angle_rad * 180 / math.pi * 60000) % 21600000)

        assert outer.get("blurRad") == expected_blur
        assert outer.get("dist") == expected_dist
        assert outer.get("dir") == expected_dir

    def test_text_shadow_mixed_glow_and_shadow(self, tmp_path: Path) -> None:
        """One glow + one outer shadow both appear in effectLst."""
        glow_shadow = TextShadow(
            offset_x_px=0.0,
            offset_y_px=0.0,
            blur_radius_px=15.0,
            color=RGBAColor(r=255, g=255, b=0),
        )
        outer_shadow = TextShadow(
            offset_x_px=5.0,
            offset_y_px=5.0,
            blur_radius_px=8.0,
            color=RGBAColor(r=0, g=0, b=0),
        )
        style = TextStyle(text_shadows=[glow_shadow, outer_shadow])
        pres = _make_slide_with(_make_text_element("Mixed", style=style))
        pptx = _build_and_read(pres, tmp_path)

        effect_lst = _get_effect_lst(pptx)
        assert effect_lst is not None

        assert effect_lst.find(qn("a:glow")) is not None, "a:glow missing"
        assert effect_lst.find(qn("a:outerShdw")) is not None, "a:outerShdw missing"

    def test_text_shadow_glow_picks_largest_blur(self, tmp_path: Path) -> None:
        """When multiple glow candidates exist, the largest blur radius is used."""
        small_glow = TextShadow(
            offset_x_px=0.0,
            offset_y_px=0.0,
            blur_radius_px=10.0,
            color=RGBAColor(r=255, g=0, b=0),
        )
        large_glow = TextShadow(
            offset_x_px=0.0,
            offset_y_px=0.0,
            blur_radius_px=20.0,
            color=RGBAColor(r=0, g=0, b=255),
        )
        style = TextStyle(text_shadows=[small_glow, large_glow])
        pres = _make_slide_with(_make_text_element("LargestGlow", style=style))
        pptx = _build_and_read(pres, tmp_path)

        effect_lst = _get_effect_lst(pptx)
        assert effect_lst is not None

        glow = effect_lst.find(qn("a:glow"))
        assert glow is not None

        expected_rad = str(px_to_emu(20.0))
        assert glow.get("rad") == expected_rad, "Glow rad should use largest blur"

    def test_text_shadow_alpha_preserved(self, tmp_path: Path) -> None:
        """Shadow with alpha < 1.0 generates an a:alpha element with val='50000'."""
        shadow = TextShadow(
            offset_x_px=0.0,
            offset_y_px=0.0,
            blur_radius_px=8.0,
            color=RGBAColor(r=0, g=0, b=0, a=0.5),
        )
        style = TextStyle(text_shadows=[shadow])
        pres = _make_slide_with(_make_text_element("Alpha", style=style))
        pptx = _build_and_read(pres, tmp_path)

        effect_lst = _get_effect_lst(pptx)
        assert effect_lst is not None

        glow = effect_lst.find(qn("a:glow"))
        assert glow is not None

        srgb = glow.find(qn("a:srgbClr"))
        assert srgb is not None

        alpha_el = srgb.find(qn("a:alpha"))
        assert alpha_el is not None, "a:alpha element not found"
        assert alpha_el.get("val") == "50000"

    def test_no_text_shadow_no_effect_lst(self, tmp_path: Path) -> None:
        """TextStyle without text_shadows produces no a:effectLst in rPr."""
        style = TextStyle()
        pres = _make_slide_with(_make_text_element("NoShadow", style=style))
        pptx = _build_and_read(pres, tmp_path)

        effect_lst = _get_effect_lst(pptx)
        assert effect_lst is None, "effectLst should be absent when no shadows set"
