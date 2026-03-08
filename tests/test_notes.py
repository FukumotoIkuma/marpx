"""Tests for speaker notes support."""

from __future__ import annotations

import tempfile
from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation

from marpx.models import Presentation, Slide
from marpx.pptx_builder.builder import build_pptx


class TestSlideModelNotes:
    """Test that Slide model accepts the notes field."""

    def test_slide_with_notes(self) -> None:
        slide = Slide(width_px=1280, height_px=720, notes="Speaker note text")
        assert slide.notes == "Speaker note text"

    def test_slide_without_notes(self) -> None:
        slide = Slide(width_px=1280, height_px=720)
        assert slide.notes is None

    def test_slide_with_none_notes(self) -> None:
        slide = Slide(width_px=1280, height_px=720, notes=None)
        assert slide.notes is None

    def test_slide_with_empty_notes(self) -> None:
        slide = Slide(width_px=1280, height_px=720, notes="")
        assert slide.notes == ""


class TestBuildPptxNotes:
    """Test that build_pptx correctly writes speaker notes."""

    def test_notes_written_to_pptx(self) -> None:
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    slide_number=0,
                    notes="This is a speaker note",
                ),
            ],
        )
        with tempfile.TemporaryDirectory() as tmp:
            out = Path(tmp) / "output.pptx"
            build_pptx(pres, out)

            # Re-read and verify
            pptx = PptxPresentation(str(out))
            slide = pptx.slides[0]
            assert slide.has_notes_slide
            notes_text = slide.notes_slide.notes_text_frame.text
            assert notes_text == "This is a speaker note"

    def test_no_notes_slide_when_notes_absent(self) -> None:
        pres = Presentation(
            slides=[
                Slide(width_px=1280, height_px=720, slide_number=0),
            ],
        )
        with tempfile.TemporaryDirectory() as tmp:
            out = Path(tmp) / "output.pptx"
            build_pptx(pres, out)

            pptx = PptxPresentation(str(out))
            slide = pptx.slides[0]
            assert not slide.has_notes_slide

    def test_multiline_notes_preserved(self) -> None:
        multi_note = "Line 1\nLine 2\n\nLine 4 after blank"
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    slide_number=0,
                    notes=multi_note,
                ),
            ],
        )
        with tempfile.TemporaryDirectory() as tmp:
            out = Path(tmp) / "output.pptx"
            build_pptx(pres, out)

            pptx = PptxPresentation(str(out))
            slide = pptx.slides[0]
            notes_text = slide.notes_slide.notes_text_frame.text
            assert "Line 1" in notes_text
            assert "Line 2" in notes_text
            assert "Line 4 after blank" in notes_text

    def test_notes_written_for_full_slide_fallback(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            fallback_image = Path(tmp) / "fallback.png"
            fallback_image.write_bytes(
                b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
                b"\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx\x9cc\xf8\xcf\xc0"
                b"\x00\x00\x03\x01\x01\x00\xc9\xfe\x92\xef\x00\x00\x00\x00IEND\xaeB`\x82"
            )
            out = Path(tmp) / "output.pptx"
            pres = Presentation(
                slides=[
                    Slide(
                        width_px=1280,
                        height_px=720,
                        slide_number=0,
                        is_fallback=True,
                        fallback_image_path=str(fallback_image),
                        notes="Fallback note",
                    ),
                ],
            )

            build_pptx(pres, out)

            pptx = PptxPresentation(str(out))
            slide = pptx.slides[0]
            assert slide.has_notes_slide
            assert slide.notes_slide.notes_text_frame.text == "Fallback note"

    def test_mixed_slides_with_and_without_notes(self) -> None:
        pres = Presentation(
            slides=[
                Slide(
                    width_px=1280,
                    height_px=720,
                    slide_number=0,
                    notes="Note for slide 1",
                ),
                Slide(width_px=1280, height_px=720, slide_number=1),
                Slide(
                    width_px=1280,
                    height_px=720,
                    slide_number=2,
                    notes="Note for slide 3",
                ),
            ],
        )
        with tempfile.TemporaryDirectory() as tmp:
            out = Path(tmp) / "output.pptx"
            build_pptx(pres, out)

            pptx = PptxPresentation(str(out))
            slides = list(pptx.slides)

            # Slide 1 has notes
            assert slides[0].has_notes_slide
            assert slides[0].notes_slide.notes_text_frame.text == "Note for slide 1"

            # Slide 2 has no notes
            assert not slides[1].has_notes_slide

            # Slide 3 has notes
            assert slides[2].has_notes_slide
            assert slides[2].notes_slide.notes_text_frame.text == "Note for slide 3"


@pytest.mark.integration
class TestNotesIntegration:
    """Integration test using full pipeline with speaker-notes.md fixture."""

    def test_full_pipeline_extracts_notes(self) -> None:
        """Run full pipeline on speaker-notes.md and verify notes in PPTX.

        Requires marp-cli and Playwright to be installed.
        """
        from marpx.extractor import extract_presentation

        fixture = Path(__file__).parent / "fixtures" / "speaker-notes.md"
        if not fixture.exists():
            pytest.skip("Fixture speaker-notes.md not found")

        # Need to convert MD -> HTML first via marp-cli
        import subprocess

        with tempfile.TemporaryDirectory() as tmp:
            html_path = Path(tmp) / "speaker-notes.html"
            result = subprocess.run(
                ["npx", "@marp-team/marp-cli", str(fixture), "-o", str(html_path)],
                capture_output=True,
                text=True,
                timeout=60,
            )
            if result.returncode != 0:
                pytest.skip(f"marp-cli failed: {result.stderr}")

            # Extract presentation
            from marpx.async_utils import run_coroutine_sync

            pres = run_coroutine_sync(extract_presentation(html_path))

            # Slide 0 should have notes
            assert pres.slides[0].notes is not None
            assert "speaker note for slide 1" in pres.slides[0].notes.lower()

            # Slide 1 should have multi-line notes
            assert pres.slides[1].notes is not None
            assert "multi-line" in pres.slides[1].notes.lower()

            # Slide 2 should have no notes
            assert pres.slides[2].notes is None

            # Build PPTX and verify
            pptx_path = Path(tmp) / "output.pptx"
            build_pptx(pres, pptx_path)

            pptx = PptxPresentation(str(pptx_path))
            slides = list(pptx.slides)
            assert slides[0].has_notes_slide
            assert slides[2].has_notes_slide is False
