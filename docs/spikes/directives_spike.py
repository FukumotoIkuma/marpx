"""
Spike: Investigate python-pptx capabilities for Marp directives
(paginate, header, footer, size)

This spike tests:
1. Slide dimensions for 16:9 vs 4:3
2. Adding slide number placeholders
3. Adding header/footer text as shapes
4. Native header/footer/slide-number support via presentation-level settings
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import lxml.etree as etree

OUTPUT_PATH = "/tmp/directives_spike_output.pptx"


def test_slide_dimensions():
    """Test setting slide dimensions for 16:9 and 4:3."""
    print("=== Test: Slide Dimensions ===")

    # 16:9 dimensions (default Marp)
    prs_16_9 = Presentation()
    prs_16_9.slide_width = Emu(12192000)   # 10 inches * 914400 + adjustments
    prs_16_9.slide_height = Emu(6858000)   # 5.625 inches * 914400 + adjustments
    print(f"16:9 - Width: {prs_16_9.slide_width} EMU = {prs_16_9.slide_width / 914400:.2f} inches")
    print(f"16:9 - Height: {prs_16_9.slide_height} EMU = {prs_16_9.slide_height / 914400:.2f} inches")

    # 4:3 dimensions
    prs_4_3 = Presentation()
    prs_4_3.slide_width = Emu(9144000)    # 10 inches
    prs_4_3.slide_height = Emu(6858000)   # 7.5 inches
    print(f"4:3 - Width: {prs_4_3.slide_width} EMU = {prs_4_3.slide_width / 914400:.2f} inches")
    print(f"4:3 - Height: {prs_4_3.slide_height} EMU = {prs_4_3.slide_height / 914400:.2f} inches")

    return prs_16_9


def test_header_footer_as_textboxes(prs):
    """Test adding header/footer as positioned text boxes."""
    print("\n=== Test: Header/Footer as Text Boxes ===")

    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Header text box (top of slide)
    header_left = Emu(0)
    header_top = Emu(0)
    header_width = slide_width
    header_height = Emu(457200)  # ~0.5 inch
    header_box = slide.shapes.add_textbox(header_left, header_top, header_width, header_height)
    tf = header_box.text_frame
    tf.text = "My Presentation"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
    print(f"Header added: left={header_left}, top={header_top}, w={header_width}, h={header_height}")

    # Footer text box (bottom-left of slide)
    footer_left = Emu(0)
    footer_top = slide_height - Emu(457200)
    footer_width = Emu(slide_width * 2 // 3)
    footer_height = Emu(457200)
    footer_box = slide.shapes.add_textbox(footer_left, footer_top, footer_width, footer_height)
    tf = footer_box.text_frame
    tf.text = "© 2026 Example Corp"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
    print(f"Footer added: left={footer_left}, top={footer_top}, w={footer_width}, h={footer_height}")

    # Page number text box (bottom-right of slide)
    pn_left = Emu(slide_width * 2 // 3)
    pn_top = slide_height - Emu(457200)
    pn_width = Emu(slide_width // 3)
    pn_height = Emu(457200)
    pn_box = slide.shapes.add_textbox(pn_left, pn_top, pn_width, pn_height)
    tf = pn_box.text_frame
    tf.text = "1"  # Static page number
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    print(f"Page number added: left={pn_left}, top={pn_top}, w={pn_width}, h={pn_height}")

    return slide


def test_native_slide_number_placeholder(prs):
    """Test using native PPTX slide number placeholder (fld element)."""
    print("\n=== Test: Native Slide Number Placeholder ===")

    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Add a text box and insert a slide number field via XML
    from pptx.oxml.ns import qn

    pn_left = Emu(slide_width * 2 // 3)
    pn_top = slide_height - Emu(457200)
    pn_width = Emu(slide_width // 3)
    pn_height = Emu(457200)
    pn_box = slide.shapes.add_textbox(pn_left, pn_top, pn_width, pn_height)

    # Access the XML to insert a slide number field
    txBody = pn_box.text_frame._txBody
    p = txBody.findall(qn('a:p'))[0]

    # Clear existing runs
    for r in p.findall(qn('a:r')):
        p.remove(r)

    # Create a field element for slide number
    fld = etree.SubElement(p, qn('a:fld'))
    fld.set('id', '{B6F15528-F159-4107-2052-3ABBE7C1B152}')
    fld.set('type', 'slidenum')

    # Add run properties and text
    rPr = etree.SubElement(fld, qn('a:rPr'))
    rPr.set('lang', 'en-US')
    rPr.set('sz', '1000')  # 10pt in hundredths of a point

    t = etree.SubElement(fld, qn('a:t'))
    t.text = '<#>'

    print("Native slide number field inserted via XML")
    print(f"Field XML: {etree.tostring(fld, pretty_print=True).decode()}")

    return slide


def test_native_hf_settings(prs):
    """Test native presentation-level header/footer settings via XML."""
    print("\n=== Test: Native Header/Footer Settings (XML exploration) ===")

    # Explore what XML elements exist for header/footer at presentation level
    prs_xml = prs.presentation.xml
    print("Presentation XML namespace elements (first 500 chars):")
    print(prs_xml[:500])

    # Check slide layout for header/footer placeholders
    for i, layout in enumerate(prs.slide_layouts):
        placeholders = list(layout.placeholders)
        if placeholders:
            ph_types = [(ph.placeholder_format.idx, ph.placeholder_format.type) for ph in placeholders]
            print(f"Layout {i} ({layout.name}): placeholders = {ph_types}")

    # Check if any layout has header/footer/slide-number placeholders
    # Placeholder types: 10=date, 11=slide number, 12=footer, 13=header
    print("\nLooking for special placeholder types (10=date, 11=slidenum, 12=footer, 13=header):")
    for i, layout in enumerate(prs.slide_layouts):
        for ph in layout.placeholders:
            idx = ph.placeholder_format.idx
            if idx in (10, 11, 12, 13):
                print(f"  Layout {i} ({layout.name}): idx={idx}, type={ph.placeholder_format.type}")


def main():
    print("Directives Spike: python-pptx header/footer/paginate/size capabilities\n")

    prs = test_slide_dimensions()
    test_header_footer_as_textboxes(prs)
    test_native_slide_number_placeholder(prs)
    test_native_hf_settings(prs)

    prs.save(OUTPUT_PATH)
    print(f"\nOutput saved to: {OUTPUT_PATH}")
    print("\n=== Summary ===")
    print("1. Slide dimensions: FULLY SUPPORTED via prs.slide_width/slide_height (EMU)")
    print("2. Header/Footer as text boxes: FULLY SUPPORTED (manual positioning)")
    print("3. Native slide number field: SUPPORTED via XML manipulation (a:fld type=slidenum)")
    print("4. Native header/footer placeholders: DEPENDS on slide layout (not all layouts have them)")


if __name__ == "__main__":
    main()
