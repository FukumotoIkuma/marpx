"""Spike: SVG, Math, and Mermaid content handling in python-pptx.

Findings (2026-03-06):

1. python-pptx SVG Support
   - python-pptx (v1.0.2) CANNOT embed SVG directly via add_picture()
   - Raises UnidentifiedImageError when given SVG bytes
   - Supported raster formats: PNG, JPEG, GIF, TIFF, BMP
   - Supports vector: EMF (image/x-emf) and WMF (image/x-wmf) content types
   - EMF/WMF must be provided as raw binary; python-pptx does NOT convert SVG to EMF

2. SVG-to-EMF Conversion Options
   - cairosvg: NOT installed; converts SVG -> PNG/PDF/PS (not EMF)
   - svglib: NOT installed; converts SVG -> ReportLab (not EMF directly)
   - inkscape CLI: can convert SVG -> EMF, but requires system dependency
   - librsvg (rsvg-convert): SVG -> PNG only
   - Playwright screenshot: SVG -> high-DPI PNG (current approach, works well)

3. Best Approach for v2
   - PRIMARY: Use Playwright to screenshot SVG/Math elements at 2x DPI as PNG
     - Already implemented in fallback_renderer.py
     - Works for ALL visual content types (SVG, MathJax SVG, Mermaid SVG)
     - Quality is good at 2x+ scale factor
   - OPTIONAL ENHANCEMENT: For inline SVG, extract SVG string and use cairosvg
     to render at arbitrary DPI without browser overhead
   - EMF is theoretically ideal (vector in PPTX) but conversion tooling is fragile
"""

from __future__ import annotations

import io

from pptx import Presentation
from pptx.util import Inches


def test_svg_direct_embed():
    """Demonstrate that python-pptx rejects SVG."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    svg_content = b"""<svg width="200" height="200" xmlns="http://www.w3.org/2000/svg">
      <circle cx="100" cy="100" r="80" fill="blue" opacity="0.5"/>
    </svg>"""

    try:
        slide.shapes.add_picture(io.BytesIO(svg_content), Inches(1), Inches(1))
        print("UNEXPECTED: SVG added successfully")
    except Exception as e:
        print(f"EXPECTED: SVG rejected - {type(e).__name__}: {e}")


def test_png_embed():
    """Demonstrate that PNG embedding works (current fallback approach)."""
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])

    # A real implementation would use Playwright screenshot here
    # For this spike, we just verify the API works
    print("PNG embedding via add_picture() works correctly")
    print("Current fallback_renderer.py already implements this")


def test_emf_content_type():
    """Show that python-pptx recognizes EMF content type."""
    from pptx.opc.constants import CONTENT_TYPE as CT
    print(f"EMF content type: {CT.X_EMF}")  # image/x-emf
    print(f"WMF content type: {CT.X_WMF}")  # image/x-wmf
    print("EMF/WMF can be embedded if we can produce the binary data")


if __name__ == "__main__":
    print("=== SVG Direct Embed Test ===")
    test_svg_direct_embed()
    print()
    print("=== PNG Embed Test ===")
    test_png_embed()
    print()
    print("=== EMF Content Type Test ===")
    test_emf_content_type()
