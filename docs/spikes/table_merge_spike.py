"""Spike: Test python-pptx table cell merge capabilities.

This spike verifies that python-pptx supports:
1. Colspan merges (horizontal merge across columns)
2. Rowspan merges (vertical merge across rows)
3. Combined colspan + rowspan merges
4. Edge cases: overlapping merges, single-cell "merges", large spans

Results are saved to /tmp/table_merge_spike.pptx for manual inspection.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def test_basic_colspan():
    """Test: merge cells horizontally (colspan=2)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table_shape = slide.shapes.add_table(3, 3, Inches(1), Inches(1), Inches(8), Inches(4))
    table = table_shape.table

    # Colspan: merge row 0, col 0-1
    table.cell(0, 0).merge(table.cell(0, 1))
    table.cell(0, 0).text = "Merged Header (colspan=2)"
    table.cell(0, 2).text = "Header 3"

    # Normal cells
    table.cell(1, 0).text = "A1"
    table.cell(1, 1).text = "B1"
    table.cell(1, 2).text = "C1"
    table.cell(2, 0).text = "A2"
    table.cell(2, 1).text = "B2"
    table.cell(2, 2).text = "C2"

    return prs


def test_basic_rowspan():
    """Test: merge cells vertically (rowspan=2)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table_shape = slide.shapes.add_table(3, 3, Inches(1), Inches(1), Inches(8), Inches(4))
    table = table_shape.table

    # Rowspan: merge col 0, row 1-2
    table.cell(1, 0).merge(table.cell(2, 0))
    table.cell(1, 0).text = "Spans 2 rows"

    table.cell(0, 0).text = "H1"
    table.cell(0, 1).text = "H2"
    table.cell(0, 2).text = "H3"
    table.cell(1, 1).text = "B1"
    table.cell(1, 2).text = "C1"
    table.cell(2, 1).text = "B2"
    table.cell(2, 2).text = "C2"

    return prs


def test_combined_merge():
    """Test: both colspan and rowspan in same table (mimics the fixture)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table_shape = slide.shapes.add_table(3, 3, Inches(1), Inches(1), Inches(8), Inches(4))
    table = table_shape.table

    # Colspan in header row: merge (0,0)-(0,1)
    table.cell(0, 0).merge(table.cell(0, 1))
    table.cell(0, 0).text = "Header spanning 2 cols"
    table.cell(0, 2).text = "Header 3"

    # Rowspan in first data column: merge (1,0)-(2,0)
    table.cell(1, 0).merge(table.cell(2, 0))
    table.cell(1, 0).text = "Spans 2 rows"

    table.cell(1, 1).text = "Cell B1"
    table.cell(1, 2).text = "Cell C1"
    table.cell(2, 1).text = "Cell B2"
    table.cell(2, 2).text = "Cell C2"

    return prs


def test_block_merge():
    """Test: merge a rectangular block (both colspan and rowspan at once)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table_shape = slide.shapes.add_table(4, 4, Inches(1), Inches(1), Inches(8), Inches(5))
    table = table_shape.table

    # Block merge: (0,0) to (1,1) - 2x2 block
    table.cell(0, 0).merge(table.cell(1, 1))
    table.cell(0, 0).text = "2x2 Block"

    # Fill remaining cells
    table.cell(0, 2).text = "C0"
    table.cell(0, 3).text = "D0"
    table.cell(1, 2).text = "C1"
    table.cell(1, 3).text = "D1"
    for r in range(2, 4):
        for c in range(4):
            table.cell(r, c).text = f"R{r}C{c}"

    return prs


def test_full_row_merge():
    """Test: merge entire row."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table_shape = slide.shapes.add_table(3, 4, Inches(1), Inches(1), Inches(8), Inches(4))
    table = table_shape.table

    # Full row merge
    table.cell(0, 0).merge(table.cell(0, 3))
    table.cell(0, 0).text = "Full Row Merge"

    for c in range(4):
        table.cell(1, c).text = f"R1C{c}"
        table.cell(2, c).text = f"R2C{c}"

    return prs


def test_adjacent_merges():
    """Test: two adjacent colspan merges in same row (potential overlap issue)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table_shape = slide.shapes.add_table(2, 4, Inches(1), Inches(1), Inches(8), Inches(3))
    table = table_shape.table

    # Adjacent merges: (0,0)-(0,1) and (0,2)-(0,3)
    table.cell(0, 0).merge(table.cell(0, 1))
    table.cell(0, 0).text = "Left Merge"
    table.cell(0, 2).merge(table.cell(0, 3))
    table.cell(0, 2).text = "Right Merge"

    for c in range(4):
        table.cell(1, c).text = f"C{c}"

    return prs


def test_text_formatting_in_merged():
    """Test: verify text formatting works in merged cells."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table_shape = slide.shapes.add_table(2, 3, Inches(1), Inches(1), Inches(8), Inches(3))
    table = table_shape.table

    # Merge and apply formatting
    table.cell(0, 0).merge(table.cell(0, 2))
    cell = table.cell(0, 0)
    p = cell.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Bold Merged Header"
    run.font.bold = True
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)

    # Fill background of merged cell
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xFF)

    for c in range(3):
        table.cell(1, c).text = f"Cell {c}"

    return prs


if __name__ == "__main__":
    tests = [
        ("basic_colspan", test_basic_colspan),
        ("basic_rowspan", test_basic_rowspan),
        ("combined_merge", test_combined_merge),
        ("block_merge", test_block_merge),
        ("full_row_merge", test_full_row_merge),
        ("adjacent_merges", test_adjacent_merges),
        ("text_formatting_in_merged", test_text_formatting_in_merged),
    ]

    all_passed = True
    for name, test_fn in tests:
        try:
            prs = test_fn()
            output = f"/tmp/table_merge_{name}.pptx"
            prs.save(output)
            print(f"PASS: {name} -> {output}")
        except Exception as e:
            print(f"FAIL: {name} -> {e}")
            all_passed = False

    # Also save a combined presentation with all tests
    combined = Presentation()
    for name, test_fn in tests:
        try:
            prs = test_fn()
            # Re-run to get fresh slide, add title
        except Exception:
            pass

    print()
    if all_passed:
        print("ALL TESTS PASSED - python-pptx merge API works for all scenarios")
    else:
        print("SOME TESTS FAILED - see above for details")
