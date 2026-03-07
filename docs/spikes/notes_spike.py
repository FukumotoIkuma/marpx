"""Spike: Verify python-pptx speaker notes API.

Creates a PPTX with 2 slides, adds speaker notes to each,
saves it, then re-opens to verify the notes are readable.
"""

from pathlib import Path
from pptx import Presentation

OUTPUT_PATH = Path("/tmp/notes_spike_output.pptx")


def create_pptx_with_notes() -> None:
    prs = Presentation()

    # Slide 1 - simple note
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    notes_slide1 = slide1.notes_slide
    notes_tf1 = notes_slide1.notes_text_frame
    notes_tf1.text = "This is a simple speaker note for slide 1."

    # Slide 2 - multi-line note
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    notes_slide2 = slide2.notes_slide
    notes_tf2 = notes_slide2.notes_text_frame
    notes_tf2.text = (
        "Multi-line speaker note.\n"
        "It has multiple paragraphs.\n"
        "\n"
        "And a list:\n"
        "- Item 1\n"
        "- Item 2"
    )

    prs.save(str(OUTPUT_PATH))
    print(f"Saved PPTX to {OUTPUT_PATH}")


def verify_notes() -> None:
    prs = Presentation(str(OUTPUT_PATH))
    for i, slide in enumerate(prs.slides, start=1):
        try:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text
            print(f"Slide {i} notes: {notes_text!r}")
        except Exception as e:
            print(f"Slide {i} notes error: {e}")


if __name__ == "__main__":
    create_pptx_with_notes()
    verify_notes()
