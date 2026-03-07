"""
Spike: Background Image Capabilities in python-pptx

This spike tests two approaches for setting slide backgrounds:
1. Slide background fill with image (XML-level manipulation)
2. Full-slide picture shape at z-index 0 (shape-level approach)

Also tests solid color background fill.

Results (all tested successfully):
- Approach 1 (XML bg fill): True slide background, no z-order issues,
  but requires XML manipulation for images. Image relationship must be
  created via `slide.part.get_or_add_image_part()` which returns (part, rId).
- Approach 2 (picture shape): Simpler API via add_picture(), needs z-order
  management via spTree insertion order (insert at index 2).
- Solid color: Natively supported via slide.background.fill.solid()
- Split backgrounds: Achievable via half-width picture shapes
- Contain: Requires aspect ratio calculation (Pillow or manual)
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import urllib.request

OUTPUT_DIR = Path("/tmp/bg_spike_output")
OUTPUT_DIR.mkdir(exist_ok=True)

# Download a test image
TEST_IMAGE = OUTPUT_DIR / "test_bg.jpg"
if not TEST_IMAGE.exists():
    print("Downloading test image...")
    urllib.request.urlretrieve("https://picsum.photos/1280/720", str(TEST_IMAGE))
    print(f"Downloaded to {TEST_IMAGE}")


def approach1_background_fill_solid():
    """Test: Set slide background to a solid color using python-pptx.

    Uses the native python-pptx API: slide.background.fill.solid()
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x33, 0x66, 0x99)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Solid Background Color (#336699)"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    out = OUTPUT_DIR / "approach1_solid_bg.pptx"
    prs.save(str(out))
    print(f"[OK] Approach 1 (solid bg): saved to {out}")
    return True


def approach1_background_fill_image():
    """Test: Set slide background to an image using XML manipulation.

    IMPORTANT: `slide.part.get_or_add_image_part()` returns a tuple (image_part, rId).
    The rId is already a valid relationship ID - do NOT call `relate_to()` again,
    as that creates a broken relationship that causes AssertionError on save.
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # get_or_add_image_part returns (ImagePart, rId) - rId is ready to use
    _image_part, rId = slide.part.get_or_add_image_part(str(TEST_IMAGE))

    # Build XML for background image fill
    bg = slide.background._element
    # Clear existing children (bgRef, etc.)
    for child in list(bg):
        bg.remove(child)

    bgPr = bg.makeelement(qn("p:bgPr"), {})

    blipFill = bgPr.makeelement(
        qn("a:blipFill"), {"dpi": "0", "rotWithShape": "1"}
    )
    blip = blipFill.makeelement(qn("a:blip"), {qn("r:embed"): rId})
    blipFill.append(blip)

    # Stretch fill = cover behavior
    stretch = blipFill.makeelement(qn("a:stretch"), {})
    fillRect = stretch.makeelement(qn("a:fillRect"), {})
    stretch.append(fillRect)
    blipFill.append(stretch)

    bgPr.append(blipFill)

    # effectLst is required by the schema
    effectLst = bgPr.makeelement(qn("a:effectLst"), {})
    bgPr.append(effectLst)

    bg.append(bgPr)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Background Image via XML Fill"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    out = OUTPUT_DIR / "approach1_image_bg.pptx"
    prs.save(str(out))
    print(f"[OK] Approach 1 (image bg via XML): saved to {out}")
    return True


def approach2_fullslide_picture():
    """Test: Add a full-slide picture shape behind other content.

    Uses the standard add_picture API and repositions via spTree insertion.
    Shapes in OOXML are rendered in document order; index 2 = behind all others
    (after nvGrpSpPr and grpSpPr which occupy indices 0 and 1).
    """
    prs = Presentation()
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    pic = slide.shapes.add_picture(
        str(TEST_IMAGE), left=0, top=0, width=slide_width, height=slide_height
    )

    # Move picture to back (z-order = 0)
    sp_tree = slide.shapes._spTree
    pic_element = pic._element
    sp_tree.remove(pic_element)
    sp_tree.insert(2, pic_element)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Background via Full-Slide Picture Shape"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    out = OUTPUT_DIR / "approach2_picture_bg.pptx"
    prs.save(str(out))
    print(f"[OK] Approach 2 (full-slide picture): saved to {out}")
    return True


def approach2_split_background():
    """Test: Simulate Marp's 'bg left' split background.

    Place picture on left half, content on right half.
    For 'bg right', mirror the positions.
    """
    prs = Presentation()
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    pic = slide.shapes.add_picture(
        str(TEST_IMAGE),
        left=0,
        top=0,
        width=slide_width // 2,
        height=slide_height,
    )

    sp_tree = slide.shapes._spTree
    sp_tree.remove(pic._element)
    sp_tree.insert(2, pic._element)

    txBox = slide.shapes.add_textbox(
        slide_width // 2 + Inches(0.5),
        Inches(1),
        slide_width // 2 - Inches(1),
        Inches(5),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Split Background Left"
    p.font.size = Pt(28)

    p2 = tf.add_paragraph()
    p2.text = "Text on the right side."
    p2.font.size = Pt(18)

    out = OUTPUT_DIR / "approach2_split_bg.pptx"
    prs.save(str(out))
    print(f"[OK] Approach 2 (split bg left): saved to {out}")
    return True


def approach_contain_background():
    """Test: Simulate 'bg contain' - image centered, aspect-ratio preserved.

    Requires calculating the scale factor to fit within slide dimensions
    without cropping. Uses Pillow to read image dimensions.
    """
    from PIL import Image

    prs = Presentation()
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    img = Image.open(str(TEST_IMAGE))
    img_w, img_h = img.size

    # Scale to fit (contain = no cropping)
    scale = min(slide_width / img_w, slide_height / img_h)
    new_w = int(img_w * scale)
    new_h = int(img_h * scale)

    # Center on slide
    left = (slide_width - new_w) // 2
    top = (slide_height - new_h) // 2

    pic = slide.shapes.add_picture(
        str(TEST_IMAGE), left=left, top=top, width=new_w, height=new_h
    )

    sp_tree = slide.shapes._spTree
    sp_tree.remove(pic._element)
    sp_tree.insert(2, pic._element)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Background Contain (aspect-ratio preserved)"
    p.font.size = Pt(24)

    out = OUTPUT_DIR / "approach_contain_bg.pptx"
    prs.save(str(out))
    print(f"[OK] Contain bg: saved to {out}")
    return True


if __name__ == "__main__":
    print("=" * 60)
    print("Background Image Spike - python-pptx Capabilities")
    print("=" * 60)

    results = {}

    print("\n--- Approach 1: Slide Background Fill ---")
    results["solid_bg"] = approach1_background_fill_solid()
    results["image_bg_xml"] = approach1_background_fill_image()

    print("\n--- Approach 2: Full-Slide Picture Shape ---")
    results["fullslide_picture"] = approach2_fullslide_picture()
    results["split_bg"] = approach2_split_background()

    try:
        print("\n--- Contain Background (requires Pillow) ---")
        results["contain_bg"] = approach_contain_background()
    except ImportError:
        print("[SKIP] Pillow not installed - contain calculation skipped")
        results["contain_bg"] = "skipped"

    print("\n" + "=" * 60)
    print("RESULTS SUMMARY")
    print("=" * 60)
    for k, v in results.items():
        status = "PASS" if v is True else ("SKIP" if v == "skipped" else "FAIL")
        print(f"  {k}: {status}")

    print(f"\nOutput files in: {OUTPUT_DIR}")
